"""
Title slide handler.

Handles title slides with main title, subtitle, and author information,
along with decorative shapes.
"""

from typing import Any
from lxml import etree
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

from .base import SlideHandler
from ..renderers import BaseRenderer, TextRenderer, ShapeRenderer
from ..config import LayoutConfig, FontConfig, ColorConfig


class TitleSlideHandler(SlideHandler):
    """
    Handler for title slides.

    Title slides contain:
    - Main title (large, centered-ish)
    - Subtitle (medium size)
    - Author/description (smaller)
    - Decorative shapes at bottom-left
    """

    @property
    def priority(self) -> int:
        """Title slides have high priority (checked early)."""
        return 10

    def can_handle(self, html_slide: etree.Element) -> bool:
        """
        Check if this is a title slide.

        Args:
            html_slide: HTML slide element

        Returns:
            True if slide has 'title-slide' class
        """
        return self._has_class(html_slide, 'title-slide')

    def handle(self, slide: Any, html_slide: etree.Element) -> None:
        """
        Process title slide using renderers.

        Args:
            slide: PowerPoint slide to populate
            html_slide: HTML slide element with content
        """
        # Initialize renderers
        base_renderer = BaseRenderer(slide)
        text_renderer = TextRenderer(slide)
        shape_renderer = ShapeRenderer(slide)

        # Apply cream background
        base_renderer.apply_background(ColorConfig.CREAM)

        # Extract elements
        title_elem = html_slide.find('.//*[@class="title-content"]/h1')
        subtitle_elem = html_slide.find('.//*[@class="subtitle"]')
        author_elem = html_slide.find('.//*[@class="author"]')

        # Add title (large, vertically centered)
        if title_elem is not None:
            title_text = self.converter.extract_text_content(title_elem)
            text_renderer.render_cover_title(title_text, y_center=2.8)

        # Add subtitle (medium size, orange)
        if subtitle_elem is not None:
            subtitle_box = base_renderer.add_textbox(
                LayoutConfig.PADDING,
                4.6,  # Below title (with proper gap)
                LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING,
                0.5
            )
            subtitle_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            p = subtitle_box.text_frame.paragraphs[0]
            p.text = self.converter.extract_text_content(subtitle_elem)
            p.alignment = PP_ALIGN.LEFT
            p.font.name = FontConfig.BODY_FONT
            p.font.size = FontConfig.HEADING_SMALL
            p.font.color.rgb = ColorConfig.ORANGE

        # Add author (smaller, muted) - handle line breaks
        if author_elem is not None:
            author_box = base_renderer.add_textbox(
                LayoutConfig.PADDING,
                5.4,  # Below subtitle (with proper gap)
                LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING,
                1.5  # Taller to accommodate multiple lines
            )
            author_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            # Parse text with <br> tags - create separate paragraphs
            import re
            from lxml import etree as ET

            # Get the HTML string
            author_html = ET.tostring(author_elem, encoding='unicode', method='html')

            # Split the HTML by <br> tags to get line segments
            line_segments = re.split(r'<br\s*/?>', author_html)

            for idx, line_html in enumerate(line_segments):
                # Clean up the line HTML
                line_html = re.sub(r'^<[^>]+>', '', line_html)  # Remove opening tag at start
                line_html = re.sub(r'<[^>]+>$', '', line_html)  # Remove closing tag at end

                if not line_html.strip():
                    continue

                # Use existing paragraph for first line, add new for others
                if idx == 0:
                    p = author_box.text_frame.paragraphs[0]
                else:
                    p = author_box.text_frame.add_paragraph()

                p.alignment = PP_ALIGN.LEFT

                # Check if this line has <strong> or <b> tags
                if '<strong>' in line_html or '<b>' in line_html:
                    # Parse with formatting - split by strong/b tags
                    parts = re.split(r'(<strong>.*?</strong>|<b>.*?</b>)', line_html)

                    for part in parts:
                        if not part.strip():
                            continue

                        # Check if this part is bold
                        is_bold = '<strong>' in part or '<b>' in part

                        # Strip tags and get text
                        text = re.sub(r'<[^>]+>', '', part).strip()

                        if text:
                            run = p.add_run()
                            run.text = text
                            run.font.name = FontConfig.BODY_FONT
                            run.font.size = FontConfig.BODY_LARGE
                            run.font.color.rgb = ColorConfig.MUTED_GRAY
                            if is_bold:
                                run.font.bold = True
                else:
                    # No formatting tags, just add plain text
                    text = re.sub(r'<[^>]+>', '', line_html).strip()
                    if text:
                        p.text = text
                        p.font.name = FontConfig.BODY_FONT
                        p.font.size = FontConfig.BODY_LARGE
                        p.font.color.rgb = ColorConfig.MUTED_GRAY

        # Add decorative shapes at bottom-left
        shape_renderer.render_decorative_shapes(
            x=LayoutConfig.DECORATIVE_SHAPE_X_START,
            y=LayoutConfig.DECORATIVE_SHAPE_Y,
            count=3,
            color=ColorConfig.TAN
        )
