"""Reflection slide handler.

Handles contemplative reflection/thinking prompt slides.
"""

from typing import Any
from lxml import etree
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

from .base import SlideHandler
from ..renderers import TextRenderer, BaseRenderer
from ..config import LayoutConfig, ColorConfig, FontConfig, SpacingConfig


class ReflectionSlideHandler(SlideHandler):
    """Handler for reflection/thinking prompt slides."""

    def __init__(self, converter: Any):
        """Initialize with converter reference.

        Args:
            converter: HTMLToPPTXConverter instance
        """
        super().__init__(converter)

    @property
    def priority(self) -> int:
        """Reflection slides have medium priority."""
        return 35

    def can_handle(self, html_slide: etree.Element) -> bool:
        """Check if this is a reflection slide.

        Args:
            html_slide: HTML slide element

        Returns:
            True if slide has 'reflection-slide' class or reflection elements
        """
        if self._has_class(html_slide, 'reflection-slide'):
            return True

        # Check for reflection-specific elements
        reflection_elem = html_slide.find('.//*[@class="reflection-question"]')
        if reflection_elem is not None:
            return True

        thinking_elem = html_slide.find('.//*[@class="thinking-prompt"]')
        if thinking_elem is not None:
            return True

        return False

    def handle(self, slide: Any, html_slide: etree.Element) -> None:
        """Process reflection slide.

        Args:
            slide: PowerPoint slide to populate
            html_slide: HTML slide element with content
        """
        # Initialize renderers
        base_renderer = BaseRenderer(slide)

        # Apply cream background
        base_renderer.apply_background(ColorConfig.CREAM)

        # Find reflection question
        question_elem = self._find_reflection_question(html_slide)
        question_text = self._extract_question_text(html_slide, question_elem)

        if not question_text:
            # Fall back to content slide if no question found
            self.converter.handle_content_slide(slide, html_slide)
            return

        # Add decorative icon/emoji at top
        icon_box = base_renderer.add_textbox(
            LayoutConfig.SLIDE_WIDTH / 2 - 0.5,
            1.8,
            1.0,
            0.8
        )
        icon_frame = icon_box.text_frame
        icon_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = icon_frame.paragraphs[0]
        p.text = "💭"  # Thought bubble emoji
        p.font.size = Pt(72)
        p.alignment = PP_ALIGN.CENTER

        # Add reflection question
        question_box = base_renderer.add_textbox(
            LayoutConfig.PADDING + 0.5,
            2.8,
            LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING - 1.0,
            2.0
        )
        question_frame = question_box.text_frame
        question_frame.word_wrap = True
        question_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        p = question_frame.paragraphs[0]
        p.text = question_text
        p.font.name = FontConfig.BODY_FONT
        p.font.size = Pt(32)
        p.font.color.rgb = ColorConfig.DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = SpacingConfig.LINE_SPACING

        # Add optional instruction text
        instruction_elem = html_slide.find('.//*[@class="reflection-instruction"]')
        if instruction_elem is not None:
            instruction_text = self.converter.extract_text_content(instruction_elem)
            instruction_box = base_renderer.add_textbox(
                LayoutConfig.PADDING + 1.0,
                5.0,
                LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING - 2.0,
                1.0
            )
            p = instruction_box.text_frame.paragraphs[0]
            p.text = instruction_text
            p.font.name = FontConfig.BODY_FONT
            p.font.size = Pt(18)
            p.font.color.rgb = ColorConfig.MUTED_GRAY
            p.alignment = PP_ALIGN.CENTER
            p.font.italic = True

        # Add footer if present
        footer_elem = html_slide.find('.//*[@class="slide-footer"]')
        if footer_elem is not None:
            text_renderer = TextRenderer(slide)
            footer_text = self.converter.extract_text_content(footer_elem)
            text_renderer.render_footer(footer_text)

    def _find_reflection_question(self, html_slide: etree.Element):
        """Find reflection question element.

        Args:
            html_slide: HTML slide element

        Returns:
            Element containing reflection question or None
        """
        # Try standard class names
        question_elem = html_slide.find('.//*[@class="reflection-question"]')
        if question_elem is not None:
            return question_elem

        question_elem = html_slide.find('.//*[@class="thinking-prompt"]')
        if question_elem is not None:
            return question_elem

        # Try finding by class containing 'reflection' or 'question'
        for elem in html_slide.iter():
            elem_class = elem.get('class', '')
            if 'reflection' in elem_class or 'question' in elem_class:
                return elem

        return None

    def _extract_question_text(self, html_slide: etree.Element, question_elem):
        """Extract question text from element.

        Args:
            html_slide: HTML slide element
            question_elem: Question element (may be None)

        Returns:
            Question text string or None
        """
        if question_elem is not None:
            return ''.join(question_elem.itertext()).strip()

        # Look for question indicators in any content
        for elem in html_slide.iter():
            if elem.tag in ['h2', 'h3', 'p', 'div']:
                text = ''.join(elem.itertext()).strip()
                if '?' in text and any(word in text.lower() for word in [
                    'reflect', 'think', 'consider', 'what', 'how', 'why'
                ]):
                    return text

        return None
