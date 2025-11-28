"""Checklist slide handler.

Handles assessment checklist slides with categorized items.
"""

from typing import Any
from lxml import etree
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .base import SlideHandler
from ..renderers import TextRenderer, BaseRenderer
from ..config import LayoutConfig, ColorConfig, FontConfig, SpacingConfig
from ..utils import XMLHelper


class ChecklistSlideHandler(SlideHandler):
    """Handler for assessment checklist slides."""

    def __init__(self, converter: Any):
        """Initialize with converter reference.

        Args:
            converter: HTMLToPPTXConverter instance
        """
        super().__init__(converter)

    @property
    def priority(self) -> int:
        """Checklist slides have medium priority."""
        return 35

    def can_handle(self, html_slide: etree.Element) -> bool:
        """Check if this is a checklist slide.

        Args:
            html_slide: HTML slide element

        Returns:
            True if slide has 'checklist-slide' class or checklist categories
        """
        if self._has_class(html_slide, 'checklist-slide'):
            return True

        # Check for checklist categories
        categories = html_slide.findall('.//*[@class="checklist-category"]')
        return len(categories) > 0

    def handle(self, slide: Any, html_slide: etree.Element) -> None:
        """Process checklist slide.

        Args:
            slide: PowerPoint slide to populate
            html_slide: HTML slide element with content
        """
        # Initialize renderers
        base_renderer = BaseRenderer(slide)
        text_renderer = TextRenderer(slide)

        # Apply background
        base_renderer.apply_background(ColorConfig.CREAM)

        # Add title
        title_elem = html_slide.find('.//h2')
        if title_elem is not None:
            title_text = self.converter.extract_text_content(title_elem)
            title_box = base_renderer.add_textbox(
                LayoutConfig.PADDING,
                LayoutConfig.PADDING,
                LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING,
                0.6
            )
            p = title_box.text_frame.paragraphs[0]
            p.text = title_text
            p.font.name = FontConfig.HEADER_FONT
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = ColorConfig.DARK_GRAY
            p.alignment = PP_ALIGN.LEFT

        # Find checklist categories
        categories = html_slide.findall('.//*[@class="checklist-category"]')
        y_pos = 1.3

        for category in categories:
            # Category header (if present)
            header = category.find('.//*[@class="category-header"]')
            if header is not None:
                h3 = header.find('.//h3')
                if h3 is not None:
                    # Add colored header box
                    header_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(LayoutConfig.PADDING),
                        Inches(y_pos),
                        Inches(LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING),
                        Inches(0.4)
                    )
                    header_shape.fill.solid()
                    header_shape.fill.fore_color.rgb = ColorConfig.ORANGE
                    header_shape.line.width = 0
                    header_shape.shadow.inherit = False

                    # Add header text
                    header_text = self.converter.extract_text_content(h3)
                    header_box = base_renderer.add_textbox(
                        LayoutConfig.PADDING + 0.2,
                        y_pos + 0.05,
                        LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING - 0.4,
                        0.3
                    )
                    p = header_box.text_frame.paragraphs[0]
                    p.text = header_text
                    p.font.name = FontConfig.HEADER_FONT
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = ColorConfig.WHITE

                    y_pos += 0.45

            # Checklist items
            items = category.find('.//*[@class="checklist-items"]')
            if items is not None:
                li_items = items.findall('.//li')
                if li_items:
                    # Calculate height
                    item_height = len(li_items) * 0.4 + 0.2
                    item_box = base_renderer.add_textbox(
                        LayoutConfig.PADDING,
                        y_pos,
                        LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING,
                        item_height
                    )
                    tf = item_box.text_frame
                    tf.word_wrap = True
                    tf.clear()

                    # Add all items as checkbox bulleted paragraphs
                    first_item = True
                    for li in li_items:
                        if first_item:
                            p = tf.paragraphs[0] if len(tf.paragraphs) > 0 else tf.add_paragraph()
                            first_item = False
                        else:
                            p = tf.add_paragraph()

                        p.level = 0
                        p.alignment = PP_ALIGN.LEFT
                        p.line_spacing = SpacingConfig.LINE_SPACING
                        p.space_before = SpacingConfig.CHECKLIST_SPACE_BEFORE
                        p.space_after = SpacingConfig.CHECKLIST_SPACE_AFTER

                        # Add checkbox bullet
                        XMLHelper.add_bullet_with_indent(p, marker='☐', font_name=FontConfig.BODY_FONT)

                        # Add formatted text
                        self.converter.add_formatted_text(p, li, skip_leading_marker=True)

                        # Apply styling to runs
                        for run in p.runs:
                            if not run.font.name:
                                run.font.name = FontConfig.BODY_FONT
                            if not run.font.size:
                                run.font.size = Pt(16)
                            # Check for existing color (from formatting)
                            try:
                                _ = run.font.color.rgb
                            except AttributeError:
                                # No color set, apply default
                                run.font.color.rgb = ColorConfig.DARK_GRAY

                    y_pos += item_height + 0.2

        # Add footer if present
        footer_elem = html_slide.find('.//*[@class="slide-footer"]')
        if footer_elem is not None:
            footer_text = self.converter.extract_text_content(footer_elem)
            text_renderer.render_footer(footer_text)
