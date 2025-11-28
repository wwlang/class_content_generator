"""Learning objectives slide handler.

Handles slides displaying learning objectives with formatted boxes.
"""

from typing import Any
from lxml import etree
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .base import SlideHandler
from ..renderers import TextRenderer, BaseRenderer
from ..config import LayoutConfig, ColorConfig, FontConfig


class ObjectivesSlideHandler(SlideHandler):
    """Handler for learning objectives slides."""

    def __init__(self, converter: Any):
        """Initialize with converter reference.

        Args:
            converter: HTMLToPPTXConverter instance
        """
        super().__init__(converter)

    @property
    def priority(self) -> int:
        """Objectives slides have medium priority."""
        return 35

    def can_handle(self, html_slide: etree.Element) -> bool:
        """Check if this is an objectives slide.

        Args:
            html_slide: HTML slide element

        Returns:
            True if slide has 'objectives-slide' class or objective items
        """
        if self._has_class(html_slide, 'objectives-slide'):
            return True

        # Check for objective items
        objective_items = html_slide.findall('.//*[@class="objective-item"]')
        return len(objective_items) > 0

    def handle(self, slide: Any, html_slide: etree.Element) -> None:
        """Process objectives slide.

        Args:
            slide: PowerPoint slide to populate
            html_slide: HTML slide element with content
        """
        # Initialize renderers
        base_renderer = BaseRenderer(slide)
        text_renderer = TextRenderer(slide)

        # Apply background (usually cream or white)
        is_dark = self._has_class(html_slide, 'dark-slide')
        if is_dark:
            base_renderer.apply_background(ColorConfig.DARK_GRAY)
        else:
            base_renderer.apply_background(ColorConfig.CREAM)

        # Add title
        title_elem = html_slide.find('.//h2')
        if title_elem is not None:
            title_text = self.converter.extract_text_content(title_elem)
            title_box = base_renderer.add_textbox(
                LayoutConfig.PADDING,
                LayoutConfig.PADDING,
                LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING,
                0.6
            )
            p = title_box.text_frame.paragraphs[0]
            p.text = title_text
            p.font.name = FontConfig.HEADER_FONT
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = ColorConfig.DARK_GRAY if not is_dark else ColorConfig.CREAM
            p.alignment = PP_ALIGN.LEFT

        # Add intro text if present
        intro = html_slide.find('.//*[@class="objectives-intro"]')
        if intro is not None:
            intro_text = self.converter.extract_text_content(intro)
            intro_box = base_renderer.add_textbox(
                LayoutConfig.PADDING,
                1.2,
                LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING,
                0.5
            )
            p = intro_box.text_frame.paragraphs[0]
            p.text = intro_text
            p.font.name = FontConfig.BODY_FONT
            p.font.size = Pt(20)
            p.font.color.rgb = ColorConfig.DARK_GRAY if not is_dark else ColorConfig.CREAM

        # Add objective items
        objective_items = html_slide.findall('.//*[@class="objective-item"]')
        y_pos = 1.8

        for item in objective_items:
            # Add background box
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(LayoutConfig.PADDING),
                Inches(y_pos),
                Inches(LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING),
                Inches(0.7)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = ColorConfig.WHITE
            bg_shape.line.color.rgb = ColorConfig.LIGHT_GRAY
            bg_shape.shadow.inherit = False

            # Extract objective text
            obj_text = item.find('.//p')
            if obj_text is not None:
                obj_text_content = self.converter.extract_text_content(obj_text)
                text_box = base_renderer.add_textbox(
                    LayoutConfig.PADDING + 0.3,
                    y_pos + 0.15,
                    LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING - 0.6,
                    0.5
                )
                p = text_box.text_frame.paragraphs[0]
                p.text = obj_text_content
                p.font.name = FontConfig.BODY_FONT
                p.font.size = Pt(18)
                p.font.color.rgb = ColorConfig.DARK_GRAY

            y_pos += 0.85

        # Add footer if present
        footer_elem = html_slide.find('.//*[@class="slide-footer"]')
        if footer_elem is not None:
            footer_text = self.converter.extract_text_content(footer_elem)
            text_renderer.render_footer(footer_text)
