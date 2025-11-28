"""Framework slide handler.

Handles framework slides with component grids showing structured information.
"""

from typing import Any, List, Dict
from lxml import etree

from .base import SlideHandler
from ..renderers import TextRenderer, ShapeRenderer, BaseRenderer
from ..config import LayoutConfig, ColorConfig, FontConfig, SpacingConfig


class FrameworkSlideHandler(SlideHandler):
    """Handler for framework slides with component grids."""

    def __init__(self, converter: Any):
        """Initialize with converter reference.

        Args:
            converter: HTMLToPPTXConverter instance
        """
        super().__init__(converter)
        self.text_renderer = None
        self.shape_renderer = None
        self.base_renderer = None

    @property
    def priority(self) -> int:
        """Framework slides have high priority (specific type)."""
        return 30

    def can_handle(self, html_slide: etree.Element) -> bool:
        """Check if this is a framework slide.

        Args:
            html_slide: HTML slide element

        Returns:
            True if slide has 'framework-slide' class
        """
        return self._has_class(html_slide, 'framework-slide')

    def handle(self, slide: Any, html_slide: etree.Element) -> None:
        """Process framework slide using renderers.

        Args:
            slide: PowerPoint slide to populate
            html_slide: HTML slide element with content
        """
        # Initialize renderers
        self.text_renderer = TextRenderer(slide)
        self.shape_renderer = ShapeRenderer(slide)
        self.base_renderer = BaseRenderer(slide)

        # Apply cream background
        self.base_renderer.apply_background(ColorConfig.CREAM)

        # Extract and render title
        title_elem = html_slide.find('.//*[@class="framework-title"]')
        if title_elem is None:
            title_elem = html_slide.find('.//h2')

        if title_elem is not None:
            title_text = self.converter.extract_text_content(title_elem)
            self.text_renderer.render_slide_title(title_text, is_dark_bg=False)

        # Extract and render subtitle if present
        subtitle_elem = html_slide.find('.//*[@class="framework-subtitle"]')
        y_start = LayoutConfig.CONTENT_START_Y

        if subtitle_elem is not None:
            subtitle_text = self.converter.extract_text_content(subtitle_elem)
            subtitle_box = self.base_renderer.add_textbox(
                LayoutConfig.PADDING,
                y_start,
                LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING,
                0.5
            )
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Pt
            p = subtitle_box.text_frame.paragraphs[0]
            p.text = subtitle_text
            p.alignment = PP_ALIGN.LEFT  # Left-aligned instead of centered
            p.font.name = FontConfig.HEADER_FONT
            p.font.size = Pt(16)  # Reduced by ~33% from 24pt
            p.font.bold = True
            p.font.color.rgb = ColorConfig.DARK_GRAY
            p.line_spacing = 1.3

            # Adjust y_start to account for subtitle (smaller now)
            y_start += 0.6

        # Extract framework components
        components_container = html_slide.find('.//*[@class="framework-components"]')
        if components_container is None:
            return

        components = components_container.findall('./div[@class="component"]')
        if not components:
            return

        # Prepare component data with structured content
        cards_data = []
        for component in components:
            title_elem = component.find('.//h3')
            title_text = self.converter.extract_text_content(title_elem) if title_elem is not None else ""

            # Extract component parts separately
            component_data = {
                'header': title_text,
                'question': None,
                'bullets': [],
                'description': []
            }

            # Extract question if present
            question_elem = component.find('.//*[@class="component-question"]')
            if question_elem is not None:
                component_data['question'] = self.converter.extract_text_content(question_elem)

            # Extract description paragraphs (excluding question)
            for p in component.findall('.//p'):
                if 'component-question' not in p.get('class', ''):
                    desc_text = self.converter.extract_text_content(p)
                    if desc_text:
                        component_data['description'].append(desc_text)

            # Extract bullets if present
            bullets_elem = component.find('.//*[@class="component-bullets"]')
            if bullets_elem is not None:
                bullet_items = bullets_elem.findall('.//li')
                for li in bullet_items:
                    bullet_text = self.converter.extract_text_content(li)
                    if bullet_text:
                        component_data['bullets'].append(bullet_text)

            cards_data.append(component_data)

        # Render component grid using framework-specific renderer
        # Calculate columns based on number of components
        num_components = len(cards_data)
        if num_components <= 3:
            cols = num_components
        elif num_components == 4:
            cols = 2
        else:
            cols = 3

        # Grid dimensions
        grid_width = LayoutConfig.SLIDE_WIDTH - 2 * LayoutConfig.PADDING
        # Further reduced height to keep cards compact with subtitle and auxiliary content
        grid_height = 3.15  # Reduced by 25% from 4.2" for more compact layout

        self.shape_renderer.render_framework_grid(
            cards_data=cards_data,
            x=LayoutConfig.PADDING,
            y=y_start,
            width=grid_width,
            height=grid_height,
            columns=cols,
            gap=0.3  # Gap to prevent card content overflow overlaps
        )

        # Handle framework auxiliary content if present
        commentary_elem = html_slide.find('.//*[@class="framework-auxiliary"]')
        if commentary_elem is None:
            # Fallback to old class name for backward compatibility
            commentary_elem = html_slide.find('.//*[@class="framework-commentary"]')
        if commentary_elem is not None:
            # Tighter spacing to keep content compact
            commentary_y = y_start + grid_height + 0.25
            commentary_box = self.base_renderer.add_textbox(
                LayoutConfig.PADDING,
                commentary_y,
                grid_width,
                1.5  # Taller box for multi-paragraph content
            )

            from pptx.enum.text import PP_ALIGN
            from pptx.util import Pt

            # Process each paragraph and list in auxiliary content
            text_frame = commentary_box.text_frame
            text_frame.clear()  # Remove default paragraph
            text_frame.word_wrap = True

            # Render paragraphs and lists
            for elem in commentary_elem:
                if elem.tag == 'p':
                    # Check if paragraph has bold heading (like "Key insight:", "With:")
                    p = text_frame.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.line_spacing = 1.3
                    p.space_after = Pt(6)

                    # Process text and strong tags
                    text_parts = []
                    if elem.text:
                        text_parts.append(('normal', elem.text))

                    for child in elem:
                        if child.tag == 'strong':
                            text_parts.append(('bold', child.text or ''))
                            if child.tail:
                                text_parts.append(('normal', child.tail))
                        elif child.tail:
                            text_parts.append(('normal', child.tail))

                    # Add text runs with appropriate formatting
                    for style, text in text_parts:
                        if not text.strip():
                            continue
                        run = p.add_run()
                        run.text = text
                        run.font.name = FontConfig.BODY_FONT
                        run.font.size = FontConfig.BODY_NORMAL  # Normal body size (not BODY_LARGE)
                        if style == 'bold':
                            run.font.bold = True
                            run.font.color.rgb = ColorConfig.ORANGE  # Orange for bold headings
                        else:
                            run.font.color.rgb = ColorConfig.DARK_GRAY

                elif elem.tag == 'ul':
                    # Render bullet list
                    for li in elem.findall('.//li'):
                        p = text_frame.add_paragraph()
                        p.alignment = PP_ALIGN.LEFT
                        p.line_spacing = 1.3
                        p.space_after = Pt(4)
                        p.level = 0  # First level bullet

                        # Process li content (may have bold text)
                        text_parts = []
                        if li.text:
                            text_parts.append(('normal', li.text))

                        for child in li:
                            if child.tag == 'strong':
                                text_parts.append(('bold', child.text or ''))
                                if child.tail:
                                    text_parts.append(('normal', child.tail))
                            elif child.tail:
                                text_parts.append(('normal', child.tail))

                        # Add text runs
                        for idx, (style, text) in enumerate(text_parts):
                            if not text.strip():
                                continue
                            if idx == 0:
                                # First run - add bullet
                                run = p.add_run()
                                run.text = '• ' + text
                            else:
                                run = p.add_run()
                                run.text = text

                            run.font.name = FontConfig.BODY_FONT
                            run.font.size = FontConfig.BODY_NORMAL
                            if style == 'bold':
                                run.font.bold = True
                                run.font.color.rgb = ColorConfig.ORANGE
                            else:
                                run.font.color.rgb = ColorConfig.DARK_GRAY

        # Add footer
        footer_elem = html_slide.find('.//*[@class="slide-footer"]')
        if footer_elem is not None:
            footer_text = self.converter.extract_text_content(footer_elem)
            self.text_renderer.render_footer(footer_text)
