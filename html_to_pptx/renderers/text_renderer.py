"""Text rendering operations for PowerPoint slides.

This module handles all text-based rendering including titles, paragraphs,
lists, and formatted text with HTML elements.
"""

from typing import List, Optional, Any
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from .base_renderer import BaseRenderer
from ..config import LayoutConfig, FontConfig, ColorConfig, SpacingConfig
from ..utils import FontStyler, XMLHelper


class TextRenderer(BaseRenderer):
    """Renderer for text-based content (titles, paragraphs, lists)."""

    def render_slide_title(
        self,
        title_text: str,
        is_dark_bg: bool = False,
        x: Optional[float] = None,
        y: Optional[float] = None,
        width: Optional[float] = None,
        height: Optional[float] = None
    ) -> Any:
        """Render consistent slide title across all slide types.

        Args:
            title_text: Title text content
            is_dark_bg: Whether slide has dark background
            x: Optional X position (defaults to padding)
            y: Optional Y position (defaults to LayoutConfig.TITLE_Y)
            width: Optional width (defaults to full width minus padding)
            height: Optional height (defaults to LayoutConfig.TITLE_HEIGHT)

        Returns:
            Title text box shape
        """
        # Use defaults from config if not specified
        if x is None:
            x = LayoutConfig.PADDING
        if y is None:
            y = LayoutConfig.TITLE_Y
        if width is None:
            width = self.config.SLIDE_WIDTH - 2 * LayoutConfig.PADDING
        if height is None:
            height = LayoutConfig.TITLE_HEIGHT

        # Create title text box
        title_box = self.add_textbox(x, y, width, height, word_wrap=True)

        # Set text frame properties
        self.set_text_frame_properties(
            title_box.text_frame,
            word_wrap=True,
            margin_top=0,
            margin_bottom=0
        )

        # Set title text and styling
        p = title_box.text_frame.paragraphs[0]
        p.text = title_text
        p.line_spacing = SpacingConfig.TITLE_LINE_SPACING
        p.alignment = PP_ALIGN.LEFT

        # Apply font styling
        p.font.name = FontConfig.HEADER_FONT
        p.font.size = FontConfig.HEADING_MEDIUM
        p.font.bold = True

        # Apply color based on background
        if is_dark_bg:
            p.font.color.rgb = ColorConfig.CREAM
        else:
            p.font.color.rgb = ColorConfig.DARK_GRAY

        # Apply condensed letter spacing
        spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
        p.font._element.set('spc', str(spacing_hundredths))

        return title_box

    def render_cover_title(
        self,
        title_text: str,
        y_center: float = 2.8,
        x: Optional[float] = None,
        width: Optional[float] = None,
        height: float = 1.5
    ) -> Any:
        """Render large centered title for cover/title slides.

        Args:
            title_text: Title text content
            y_center: Vertical center position (default 2.8 for centered)
            x: Optional X position (defaults to padding)
            width: Optional width (defaults to full width minus padding)
            height: Text box height (default 1.5)

        Returns:
            Title text box shape
        """
        from pptx.enum.text import MSO_VERTICAL_ANCHOR

        # Use defaults from config if not specified
        if x is None:
            x = LayoutConfig.PADDING
        if width is None:
            width = self.config.SLIDE_WIDTH - 2 * LayoutConfig.PADDING

        # Create title text box
        title_box = self.add_textbox(x, y_center, width, height, word_wrap=True)
        title_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Set title text and styling
        p = title_box.text_frame.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.LEFT

        # Apply font styling - larger size for cover
        p.font.name = FontConfig.HEADER_FONT
        p.font.size = FontConfig.TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = ColorConfig.DARK_GRAY

        # Apply condensed letter spacing
        spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
        p.font._element.set('spc', str(spacing_hundredths))

        return title_box

    def render_section_title(
        self,
        title_text: str,
        y: Optional[float] = None,
        x: Optional[float] = None,
        width: Optional[float] = None,
        height: Optional[float] = None
    ) -> Any:
        """Render section break title (large white text on colored background).

        Handles multi-line text (with \n) by creating multiple paragraphs.

        Args:
            title_text: Title text content (may contain \n for line breaks)
            y: Optional Y position (defaults to LayoutConfig.TITLE_Y for consistency)
            x: Optional X position (defaults to padding)
            width: Optional width (defaults to full width minus padding)
            height: Optional height (defaults to larger section title height)

        Returns:
            Title text box shape
        """
        from pptx.dml.color import RGBColor

        # Use defaults from config if not specified
        if x is None:
            x = LayoutConfig.PADDING
        if y is None:
            y = LayoutConfig.TITLE_Y  # Use standard title position for consistency
        if width is None:
            width = self.config.SLIDE_WIDTH - 2 * LayoutConfig.PADDING
        if height is None:
            height = 1.5  # Taller for section titles

        # Create title text box
        title_box = self.add_textbox(x, y, width, height, word_wrap=True)
        tf = title_box.text_frame

        # Handle multi-line text by splitting on newlines
        lines = title_text.split('\n')

        # First paragraph
        p = tf.paragraphs[0]
        p.text = lines[0] if lines else title_text
        p.alignment = PP_ALIGN.LEFT

        # Apply font styling to first paragraph
        p.font.name = FontConfig.HEADER_FONT
        p.font.size = FontConfig.SECTION_TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # White on colored background

        # Apply condensed letter spacing
        spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
        p.font._element.set('spc', str(spacing_hundredths))

        # Add additional lines as new paragraphs
        for line in lines[1:]:
            if line.strip():  # Skip empty lines
                p = tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.LEFT

                # Apply same styling to additional paragraphs
                p.font.name = FontConfig.HEADER_FONT
                p.font.size = FontConfig.SECTION_TITLE_SIZE
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font._element.set('spc', str(spacing_hundredths))

        return title_box

    def render_paragraph(
        self,
        text_frame: Any,
        text: str,
        is_first: bool = True,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        font_size: Pt = Pt(16),
        is_dark_bg: bool = False,
        space_before: Pt = None,
        space_after: Pt = None
    ) -> Any:
        """Render a paragraph in text frame.

        Args:
            text_frame: PowerPoint text frame to add paragraph to
            text: Paragraph text
            is_first: Whether this is first paragraph (uses existing)
            alignment: Text alignment
            font_size: Font size
            is_dark_bg: Whether background is dark
            space_before: Space before paragraph
            space_after: Space after paragraph

        Returns:
            Paragraph object
        """
        # Use first paragraph if available, otherwise add new
        if is_first and len(text_frame.paragraphs) > 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = text
        p.alignment = alignment
        p.line_spacing = SpacingConfig.LINE_SPACING
        p.space_before = space_before if space_before is not None else SpacingConfig.PARAGRAPH_SPACE_BEFORE
        p.space_after = space_after if space_after is not None else SpacingConfig.PARAGRAPH_SPACE_AFTER

        # Apply font styling
        for run in p.runs:
            run.font.name = FontConfig.BODY_FONT
            run.font.size = font_size

            # Apply color based on background
            if is_dark_bg:
                run.font.color.rgb = ColorConfig.CREAM
            else:
                run.font.color.rgb = ColorConfig.DARK_GRAY

        return p

    def render_bullet_list(
        self,
        text_frame: Any,
        items: List[str],
        is_dark_bg: bool = False,
        marker: str = '•',
        is_first: bool = True
    ) -> List[Any]:
        """Render a bulleted list in text frame.

        Args:
            text_frame: PowerPoint text frame
            items: List of item texts
            is_dark_bg: Whether background is dark
            marker: Bullet character
            is_first: Whether first item uses existing paragraph

        Returns:
            List of paragraph objects
        """
        paragraphs = []
        first_item = is_first

        for item_text in items:
            # Use first paragraph or add new
            if first_item and len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
                first_item = False
            else:
                p = text_frame.add_paragraph()

            p.text = item_text
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = SpacingConfig.LINE_SPACING
            p.space_before = SpacingConfig.LIST_ITEM_SPACE_BEFORE
            p.space_after = SpacingConfig.LIST_ITEM_SPACE_AFTER

            # Add bullet with indentation using helper
            XMLHelper.add_bullet_with_indent(p, marker=marker, font_name=FontConfig.BODY_FONT)

            # Apply font styling using helper
            FontStyler.apply_body_font_styling(p.runs, is_dark_bg=is_dark_bg)

            paragraphs.append(p)

        return paragraphs

    def render_numbered_list(
        self,
        text_frame: Any,
        items: List[str],
        is_dark_bg: bool = False,
        style: str = 'arabicPeriod',
        is_first: bool = True
    ) -> List[Any]:
        """Render a numbered list in text frame.

        Args:
            text_frame: PowerPoint text frame
            items: List of item texts
            is_dark_bg: Whether background is dark
            style: Numbering style
            is_first: Whether first item uses existing paragraph

        Returns:
            List of paragraph objects
        """
        paragraphs = []
        first_item = is_first

        for item_text in items:
            # Use first paragraph or add new
            if first_item and len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
                first_item = False
            else:
                p = text_frame.add_paragraph()

            p.text = item_text
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = SpacingConfig.LINE_SPACING
            p.space_before = SpacingConfig.LIST_ITEM_SPACE_BEFORE
            p.space_after = SpacingConfig.LIST_ITEM_SPACE_AFTER

            # Add numbering with indentation using helper
            XMLHelper.add_numbering_with_indent(p, style=style)

            # Apply font styling using helper
            FontStyler.apply_body_font_styling(p.runs, is_dark_bg=is_dark_bg)

            paragraphs.append(p)

        return paragraphs

    def render_footer(
        self,
        footer_text: str,
        slide_number: Optional[int] = None,
        y_position: Optional[float] = None
    ) -> Any:
        """Render slide footer with optional page number.

        Args:
            footer_text: Footer text content
            slide_number: Optional slide number to display
            y_position: Optional Y position (defaults to LayoutConfig.FOOTER_Y)

        Returns:
            Footer text box shape
        """
        if y_position is None:
            y_position = LayoutConfig.FOOTER_Y

        footer_box = self.add_textbox(
            LayoutConfig.PADDING,
            y_position,
            self.config.SLIDE_WIDTH - 2 * LayoutConfig.PADDING,
            LayoutConfig.FOOTER_HEIGHT
        )

        p = footer_box.text_frame.paragraphs[0]
        p.text = footer_text
        p.alignment = PP_ALIGN.LEFT

        p.font.name = FontConfig.BODY_FONT
        p.font.size = FontConfig.FOOTER_SIZE
        p.font.color.rgb = ColorConfig.MUTED_GRAY

        # Add slide number if provided
        if slide_number:
            number_box = self.add_textbox(
                self.config.SLIDE_WIDTH - LayoutConfig.FOOTER_NUMBER_X_OFFSET,
                y_position,
                LayoutConfig.FOOTER_NUMBER_WIDTH,
                LayoutConfig.FOOTER_HEIGHT
            )

            p_num = number_box.text_frame.paragraphs[0]
            p_num.text = str(slide_number)
            p_num.alignment = PP_ALIGN.RIGHT

            p_num.font.name = FontConfig.BODY_FONT
            p_num.font.size = FontConfig.FOOTER_SIZE
            p_num.font.color.rgb = ColorConfig.MUTED_GRAY

        return footer_box

    def render_quote(
        self,
        quote_text: str,
        author: Optional[str] = None,
        x: Optional[float] = None,
        y: Optional[float] = None,
        width: Optional[float] = None,
        height: Optional[float] = None
    ) -> Any:
        """Render a quote with optional attribution.

        Args:
            quote_text: Quote text
            author: Optional author/attribution
            x: Optional X position (defaults to centered)
            y: Optional Y position (defaults to centered)
            width: Optional width (defaults to 70% of slide width)
            height: Optional height

        Returns:
            Quote text box shape
        """
        # Use defaults for centered quote
        if width is None:
            width = self.config.SLIDE_WIDTH * 0.7
        if x is None:
            x = (self.config.SLIDE_WIDTH - width) / 2
        if y is None:
            y = LayoutConfig.CONTENT_START_Y + 0.5
        if height is None:
            height = 3.0

        quote_box = self.add_textbox(x, y, width, height)

        # Set margins for quote text
        quote_box.text_frame.margin_left = self.inches(0.1)  # 0.25cm
        quote_box.text_frame.margin_right = self.inches(0.1)
        quote_box.text_frame.margin_top = 0
        quote_box.text_frame.margin_bottom = 0

        p = quote_box.text_frame.paragraphs[0]
        p.text = quote_text
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2  # Tighter line spacing for better readability

        # Set font properties on the run (not paragraph)
        if p.runs:
            run = p.runs[0]
            run.font.name = FontConfig.BODY_FONT
            run.font.size = Pt(36)  # Match HTML (36px) - large, prominent quote text
            run.font.italic = True
            run.font.color.rgb = ColorConfig.DARK_GRAY

        # Add author if provided
        if author:
            p_author = quote_box.text_frame.add_paragraph()
            p_author.text = f"— {author}"
            p_author.alignment = PP_ALIGN.RIGHT
            p_author.space_before = Pt(12)  # Extra space before author attribution

            # Set font properties on the run
            if p_author.runs:
                run_author = p_author.runs[0]
                run_author.font.name = FontConfig.BODY_FONT
                run_author.font.size = Pt(18)
                run_author.font.color.rgb = ColorConfig.MUTED_GRAY

        return quote_box
