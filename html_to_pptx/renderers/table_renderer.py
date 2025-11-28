"""Table rendering operations for PowerPoint slides.

This module handles rendering of tables with proper styling, column widths,
and cell formatting.
"""

from typing import List, Dict, Any, Optional
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

from .base_renderer import BaseRenderer
from ..config import LayoutConfig, FontConfig, ColorConfig, SpacingConfig


class TableRenderer(BaseRenderer):
    """Renderer for table-based layouts."""

    def render_table(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        rows: int,
        cols: int,
        data: Optional[List[List[str]]] = None,
        has_header: bool = True
    ) -> Any:
        """Render a table with specified dimensions.

        Args:
            x: X position in inches
            y: Y position in inches
            width: Table width in inches
            height: Table height in inches
            rows: Number of rows
            cols: Number of columns
            data: Optional 2D list of cell data
            has_header: Whether first row is a header

        Returns:
            PowerPoint table shape
        """
        table_shape = self.slide.shapes.add_table(
            rows, cols,
            self.inches(x),
            self.inches(y),
            self.inches(width),
            self.inches(height)
        )

        table = table_shape.table

        # Distribute width evenly across columns for autofit
        col_width = width / cols
        for col in table.columns:
            col.width = self.inches(col_width)

        # Set tighter row heights for better auto-sizing
        # Use smaller heights to reduce excessive vertical spacing
        for row_idx, row in enumerate(table.rows):
            # Header rows slightly taller, data rows compact
            if has_header and row_idx == 0:
                row.height = self.inches(0.35)  # Compact header
            else:
                row.height = self.inches(0.45)  # Compact data rows

        # Add data if provided
        if data:
            for row_idx, row_data in enumerate(data):
                for col_idx, cell_data in enumerate(row_data):
                    if row_idx < rows and col_idx < cols:
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_data)

                        # Style based on whether it's header or body
                        is_header_cell = has_header and row_idx == 0
                        self.style_table_cell(cell, is_header=is_header_cell)

        return table_shape

    def style_table_cell(
        self,
        cell: Any,
        is_header: bool = False,
        bg_color: Optional[Any] = None,
        text_color: Optional[Any] = None,
        font_size: Optional[Pt] = None
    ) -> None:
        """Apply styling to a table cell.

        Args:
            cell: PowerPoint table cell
            is_header: Whether this is a header cell
            bg_color: Optional background color
            text_color: Optional text color
            font_size: Optional font size
        """
        # Set background color
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
        elif is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ColorConfig.ORANGE
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ColorConfig.WHITE

        # Style text
        if cell.text_frame and cell.text_frame.paragraphs:
            # Enable text wrapping
            cell.text_frame.word_wrap = True

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                # Use unified line spacing for tables
                paragraph.line_spacing = SpacingConfig.LINE_SPACING
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(0)

                for run in paragraph.runs:
                    run.font.name = FontConfig.BODY_FONT

                    if font_size:
                        run.font.size = font_size
                    elif is_header:
                        run.font.size = Pt(14)
                        run.font.bold = True
                    else:
                        run.font.size = Pt(13)

                    if text_color:
                        run.font.color.rgb = text_color
                    elif is_header:
                        run.font.color.rgb = ColorConfig.WHITE
                    else:
                        run.font.color.rgb = ColorConfig.DARK_GRAY

        # Set cell margins - reduced for tighter spacing
        cell.margin_left = self.inches(0.08)
        cell.margin_right = self.inches(0.08)
        cell.margin_top = self.inches(0.03)
        cell.margin_bottom = self.inches(0.03)

        # Set vertical alignment
        if hasattr(cell, 'vertical_anchor'):
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    def render_vocab_table(
        self,
        x: float,
        y: float,
        width: float,
        vocab_data: List[Dict[str, str]],
        headers: Optional[List[str]] = None
    ) -> Any:
        """Render a vocabulary table with 2 or 3 columns.

        Args:
            x: X position
            y: Y position
            width: Table width
            vocab_data: List of dicts with 'term', optional 'translation', and 'definition' keys
            headers: Optional list of header labels

        Returns:
            Table shape
        """
        rows = len(vocab_data) + 1  # +1 for header

        # Detect number of columns from data structure
        has_translation = any('translation' in item for item in vocab_data)
        cols = 3 if has_translation else 2

        # Use headers from HTML if provided, otherwise use defaults
        if headers and len(headers) >= cols:
            # Use headers from HTML exactly as provided
            table_data = [headers[:cols]]
        else:
            # Fallback only if no headers were provided in HTML
            if has_translation:
                table_data = [['Term', 'Translation', 'Definition']]
            else:
                table_data = [['Term', 'Definition']]

        # Add data rows
        for item in vocab_data:
            if has_translation:
                table_data.append([
                    item.get('term', ''),
                    item.get('translation', ''),
                    item.get('definition', '')
                ])
            else:
                table_data.append([
                    item.get('term', ''),
                    item.get('definition', '')
                ])

        # Calculate height based on number of rows for better auto-sizing
        height = min(5.5, 0.4 + (rows * 0.6))  # Base + per-row estimate

        return self.render_table(
            x, y, width, height,
            rows, cols,
            data=table_data,
            has_header=True
        )

    def render_comparison_table(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        comparison_data: List[Dict[str, str]],
        headers: Optional[List[str]] = None
    ) -> Any:
        """Render a comparison table.

        Args:
            x: X position
            y: Y position
            width: Table width
            height: Table height
            comparison_data: List of row dicts
            headers: Optional list of header labels

        Returns:
            Table shape
        """
        if not comparison_data:
            return None

        # Determine columns from first row
        cols = len(comparison_data[0])
        rows = len(comparison_data) + (1 if headers else 0)

        # Prepare table data
        table_data = []
        if headers:
            table_data.append(headers)

        for row in comparison_data:
            table_data.append(list(row.values()))

        return self.render_table(
            x, y, width, height,
            rows, cols,
            data=table_data,
            has_header=bool(headers)
        )

    def set_column_widths(
        self,
        table: Any,
        widths: List[float]
    ) -> None:
        """Set specific column widths for table.

        Args:
            table: PowerPoint table
            widths: List of column widths in inches
        """
        for col_idx, width in enumerate(widths):
            if col_idx < len(table.columns):
                table.columns[col_idx].width = self.inches(width)

    def set_row_heights(
        self,
        table: Any,
        heights: List[float]
    ) -> None:
        """Set specific row heights for table.

        Args:
            table: PowerPoint table
            heights: List of row heights in inches
        """
        for row_idx, height in enumerate(heights):
            if row_idx < len(table.rows):
                table.rows[row_idx].height = self.inches(height)

    def merge_cells(
        self,
        table: Any,
        start_row: int,
        start_col: int,
        end_row: int,
        end_col: int
    ) -> Any:
        """Merge a range of cells in table.

        Args:
            table: PowerPoint table
            start_row: Starting row index
            start_col: Starting column index
            end_row: Ending row index
            end_col: Ending column index

        Returns:
            Merged cell
        """
        start_cell = table.cell(start_row, start_col)
        end_cell = table.cell(end_row, end_col)
        return start_cell.merge(end_cell)
