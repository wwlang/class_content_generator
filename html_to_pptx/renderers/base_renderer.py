"""Base renderer class for PowerPoint rendering operations.

This module provides the abstract base class that all specific renderers
inherit from, defining common operations and the rendering interface.
"""

from abc import ABC
from typing import Optional, Any
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..config import LayoutConfig, FontConfig, ColorConfig
from ..utils import FontStyler, XMLHelper


class BaseRenderer(ABC):
    """Abstract base class for all renderers.

    Provides common operations for creating shapes, text boxes, and
    applying basic styling. Specific renderers inherit and extend this.
    """

    def __init__(self, slide: Any, config: LayoutConfig = None):
        """Initialize renderer with slide and configuration.

        Args:
            slide: PowerPoint slide object to render to
            config: Layout configuration (uses defaults if not provided)
        """
        self.slide = slide
        self.config = config or LayoutConfig()

    def add_textbox(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        text: str = "",
        word_wrap: bool = True
    ) -> Any:
        """Create a text box at specified position.

        Args:
            x: X position in inches
            y: Y position in inches
            width: Width in inches
            height: Height in inches
            text: Optional initial text content
            word_wrap: Enable word wrapping

        Returns:
            PowerPoint text box shape
        """
        text_box = self.slide.shapes.add_textbox(
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height)
        )

        text_box.text_frame.word_wrap = word_wrap

        if text:
            text_box.text = text

        return text_box

    def add_rectangle(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        fill_color: Optional[RGBColor] = None,
        line_color: Optional[RGBColor] = None
    ) -> Any:
        """Create a rectangle shape.

        Args:
            x: X position in inches
            y: Y position in inches
            width: Width in inches
            height: Height in inches
            fill_color: Optional fill color
            line_color: Optional line color

        Returns:
            PowerPoint rectangle shape
        """
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height)
        )

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color

        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()  # No line

        return shape

    def apply_text_frame_margins(
        self,
        text_frame: Any,
        top: float = 0,
        bottom: float = 0,
        left: float = 0,
        right: float = 0
    ) -> None:
        """Apply margins to text frame.

        Args:
            text_frame: PowerPoint text frame
            top: Top margin in inches
            bottom: Bottom margin in inches
            left: Left margin in inches
            right: Right margin in inches
        """
        text_frame.margin_top = Inches(top)
        text_frame.margin_bottom = Inches(bottom)
        text_frame.margin_left = Inches(left)
        text_frame.margin_right = Inches(right)

    def set_text_frame_properties(
        self,
        text_frame: Any,
        word_wrap: bool = True,
        vertical_anchor: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
        margin_top: float = 0,
        margin_bottom: float = 0
    ) -> None:
        """Set common text frame properties.

        Args:
            text_frame: PowerPoint text frame
            word_wrap: Enable word wrapping
            vertical_anchor: Vertical alignment
            margin_top: Top margin in inches
            margin_bottom: Bottom margin in inches
        """
        text_frame.word_wrap = word_wrap
        text_frame.vertical_anchor = vertical_anchor
        text_frame.margin_top = Inches(margin_top)
        text_frame.margin_bottom = Inches(margin_bottom)

    def apply_background(
        self,
        color: RGBColor,
        width: Optional[float] = None,
        height: Optional[float] = None
    ) -> Any:
        """Apply solid background color to slide.

        Args:
            color: Background color
            width: Optional width (defaults to slide width)
            height: Optional height (defaults to slide height)

        Returns:
            Background rectangle shape
        """
        if width is None:
            width = self.config.SLIDE_WIDTH
        if height is None:
            height = self.config.SLIDE_HEIGHT

        bg_shape = self.add_rectangle(0, 0, width, height, fill_color=color)

        # Move to back so it doesn't cover content
        self.slide.shapes._spTree.remove(bg_shape._element)
        self.slide.shapes._spTree.insert(2, bg_shape._element)

        return bg_shape

    def inches(self, value: float) -> Inches:
        """Convert float to Inches object (convenience method).

        Args:
            value: Value in inches

        Returns:
            Inches object
        """
        return Inches(value)

    def points(self, value: float) -> Pt:
        """Convert float to Pt object (convenience method).

        Args:
            value: Value in points

        Returns:
            Pt object
        """
        return Pt(value)
