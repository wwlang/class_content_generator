"""Shape rendering operations for PowerPoint slides.

This module handles shape-based layouts including cards, grids, stats banners,
and decorative elements.
"""

from typing import List, Dict, Any, Optional
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE

from .base_renderer import BaseRenderer
from ..config import LayoutConfig, FontConfig, ColorConfig, SpacingConfig
from ..utils import LayoutCalculator, LayoutBox


class ShapeRenderer(BaseRenderer):
    """Renderer for shape-based layouts (cards, grids, stats)."""

    def render_card(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        header_text: str,
        content_text: str,
        header_color: Optional[Any] = None,
        autofit: bool = False
    ) -> Dict[str, Any]:
        """Render a card with colored header and content.

        Args:
            x: X position in inches
            y: Y position in inches
            width: Card width in inches
            height: Card height in inches
            header_text: Header text
            content_text: Content text
            header_color: Optional header background color
            autofit: Enable shape autofit for content

        Returns:
            Dict with 'card', 'header', and 'content' shape references
        """
        if header_color is None:
            header_color = ColorConfig.ORANGE

        # Card background (white)
        card_bg = self.add_rectangle(x, y, width, height, fill_color=ColorConfig.WHITE)

        # Header bar
        header_height = LayoutConfig.CARD_HEADER_HEIGHT
        header_shape = self.add_rectangle(
            x, y, width, header_height,
            fill_color=header_color
        )

        # Header text
        header_text_box = self.add_textbox(
            x + 0.15, y + 0.1,
            width - 0.3, header_height - 0.2,
            text=header_text
        )

        p = header_text_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.name = FontConfig.HEADER_FONT
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ColorConfig.WHITE

        # Content text
        content_y = y + header_height + 0.1
        content_height = height - header_height - 0.2

        content_box = self.add_textbox(
            x + 0.15, content_y,
            width - 0.3, content_height,
            text=content_text
        )

        content_box.text_frame.margin_left = self.inches(LayoutConfig.CARD_PADDING)
        content_box.text_frame.margin_right = self.inches(LayoutConfig.CARD_PADDING)
        content_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        # Enable autofit if requested
        if autofit:
            content_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        p_content = content_box.text_frame.paragraphs[0]
        p_content.alignment = PP_ALIGN.LEFT
        p_content.font.name = FontConfig.BODY_FONT
        p_content.font.size = Pt(14)
        p_content.font.color.rgb = ColorConfig.DARK_GRAY

        return {
            'card': card_bg,
            'header': header_shape,
            'content': content_box
        }

    def render_card_grid(
        self,
        cards_data: List[Dict[str, str]],
        x: float,
        y: float,
        width: float,
        height: float,
        columns: int = 2,
        gap: float = 0.4
    ) -> List[Dict[str, Any]]:
        """Render a grid of cards.

        Args:
            cards_data: List of dicts with 'header' and 'content' keys
            x: Container X position
            y: Container Y position
            width: Container width
            height: Container height
            columns: Number of columns
            gap: Gap between cards

        Returns:
            List of card shape references
        """
        # Calculate card positions using layout helper
        num_cards = len(cards_data)
        card_height = LayoutConfig.CARD_MIN_HEIGHT

        layout_boxes = LayoutCalculator.calculate_card_layout(
            num_cards, x, y, width, card_height, columns, gap
        )

        rendered_cards = []
        for i, (box, card_data) in enumerate(zip(layout_boxes, cards_data)):
            card = self.render_card(
                box.x, box.y, box.width, box.height,
                header_text=card_data.get('header', ''),
                content_text=card_data.get('content', ''),
                autofit=True
            )
            rendered_cards.append(card)

        return rendered_cards

    def render_framework_card(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        title_text: str,
        content_text: str = None,
        component_data: Dict = None
    ) -> Dict[str, Any]:
        """Render a framework component card (white with border, structured content).

        Framework cards match HTML design:
        - White background with gray border
        - Orange title (h3 style, 27px → 20pt)
        - Italic question (12px → 9pt)
        - Left-aligned bullets (17px → 13pt)
        - Centered description (17px → 13pt)

        Args:
            x: X position
            y: Y position
            width: Card width
            height: Card height
            title_text: Card title (orange, centered)
            content_text: Legacy content text (for backward compatibility)
            component_data: Structured data with 'question', 'bullets', 'description' lists

        Returns:
            Dict with 'card' shape reference
        """
        # White card with border - use rounded rectangle
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        card_bg = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height)
        )

        # Set fill color
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = ColorConfig.WHITE

        # Add gray border (1px in HTML = thin line in PPTX)
        card_bg.line.color.rgb = ColorConfig.LIGHT_GRAY
        card_bg.line.width = Pt(1)

        # Remove shadow
        card_bg.shadow.inherit = False

        # Reduce roundedness by half (default is ~0.16667, set to 0.08)
        card_bg.adjustments[0] = 0.08

        # Use the shape's text frame directly (not separate text boxes)
        text_frame = card_bg.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        # Set internal margins for proper spacing
        text_frame.margin_left = self.inches(0.15)
        text_frame.margin_right = self.inches(0.15)
        text_frame.margin_top = self.inches(0.12)
        text_frame.margin_bottom = self.inches(0.12)

        # Add title paragraph (orange, centered, 27px → 20pt)
        p_title = text_frame.paragraphs[0]
        p_title.text = title_text
        p_title.alignment = PP_ALIGN.CENTER
        p_title.line_spacing = 1.2
        p_title.space_after = Pt(6)
        p_title.font.name = FontConfig.HEADER_FONT
        p_title.font.size = Pt(20)  # 27px * 0.75 ≈ 20pt
        p_title.font.bold = True
        p_title.font.color.rgb = ColorConfig.ORANGE

        # If structured data provided, use it
        if component_data:
            # Add question if present (italic, centered, 12px → 9pt)
            if component_data.get('question'):
                p_q = text_frame.add_paragraph()
                p_q.text = component_data['question']
                p_q.alignment = PP_ALIGN.CENTER
                p_q.line_spacing = 1.2
                p_q.space_after = Pt(4)
                p_q.font.name = FontConfig.BODY_FONT
                p_q.font.size = Pt(11)  # 12px * 0.75 ≈ 9pt, increased slightly for readability
                p_q.font.italic = True
                p_q.font.color.rgb = ColorConfig.DARK_GRAY

            # Add bullets if present (left-aligned, 17px → 13pt)
            if component_data.get('bullets'):
                for bullet in component_data['bullets']:
                    p_b = text_frame.add_paragraph()
                    p_b.text = '• ' + bullet
                    p_b.alignment = PP_ALIGN.LEFT
                    p_b.line_spacing = 1.3
                    p_b.space_after = Pt(2)
                    p_b.font.name = FontConfig.BODY_FONT
                    p_b.font.size = Pt(11)  # 17px * 0.75 ≈ 13pt, reduced slightly to fit
                    p_b.font.color.rgb = ColorConfig.MUTED_GRAY

            # Add description paragraphs if present (centered, 17px → 13pt)
            if component_data.get('description'):
                for desc in component_data['description']:
                    p_d = text_frame.add_paragraph()
                    p_d.text = desc
                    p_d.alignment = PP_ALIGN.CENTER
                    p_d.line_spacing = 1.3
                    p_d.space_after = Pt(4)
                    p_d.font.name = FontConfig.BODY_FONT
                    p_d.font.size = Pt(13)  # 17px * 0.75 ≈ 13pt
                    p_d.font.color.rgb = ColorConfig.MUTED_GRAY

        # Fallback to legacy content_text for backward compatibility
        elif content_text:
            p_content = text_frame.add_paragraph()
            p_content.text = content_text
            p_content.alignment = PP_ALIGN.CENTER
            p_content.line_spacing = 1.3
            p_content.space_after = Pt(0)
            p_content.font.name = FontConfig.BODY_FONT
            p_content.font.size = Pt(13)  # Updated from 16 to match HTML
            p_content.font.color.rgb = ColorConfig.DARK_GRAY

        return {
            'card': card_bg
        }

    def render_framework_grid(
        self,
        cards_data: List[Dict[str, str]],
        x: float,
        y: float,
        width: float,
        height: float,
        columns: int = 3,
        gap: float = 0.2
    ) -> List[Dict[str, Any]]:
        """Render a grid of framework component cards.

        Args:
            cards_data: List of dicts with 'header' and 'content' keys
            x: Container X position
            y: Container Y position
            width: Container width
            height: Container height
            columns: Number of columns (default 3 for framework)
            gap: Gap between cards (smaller for framework, default 0.2)

        Returns:
            List of card shape references
        """
        from ..utils import LayoutCalculator

        # Calculate card positions
        num_cards = len(cards_data)

        layout_boxes = LayoutCalculator.calculate_grid_layout(
            num_cards, x, y, width, height, columns, gap
        )

        rendered_cards = []
        for box, card_data in zip(layout_boxes, cards_data):
            card = self.render_framework_card(
                box.x, box.y, box.width, box.height,
                title_text=card_data.get('header', ''),
                component_data=card_data  # Pass structured data instead of concatenated content
            )
            rendered_cards.append(card)

        return rendered_cards

    def render_grid_item(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        title: str,
        description: str,
        accent_color: Optional[Any] = None
    ) -> Dict[str, Any]:
        """Render a grid item with accent bar, title, and description.

        Args:
            x: X position
            y: Y position
            width: Item width
            height: Item height
            title: Item title
            description: Item description
            accent_color: Optional accent bar color

        Returns:
            Dict with shape references
        """
        if accent_color is None:
            accent_color = ColorConfig.ORANGE

        # Background
        bg = self.add_rectangle(x, y, width, height, fill_color=ColorConfig.WHITE)

        # Accent bar at top
        accent_bar = self.add_rectangle(
            x, y, width, LayoutConfig.GRID_ITEM_ACCENT_HEIGHT,
            fill_color=accent_color
        )

        # Title
        title_y = y + 0.15
        title_box = self.add_textbox(
            x + 0.15, title_y,
            width - 0.3, 0.4,
            text=title
        )

        p_title = title_box.text_frame.paragraphs[0]
        p_title.font.name = FontConfig.HEADER_FONT
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = ColorConfig.DARK_GRAY

        # Description
        desc_y = title_y + 0.5
        desc_box = self.add_textbox(
            x + 0.15, desc_y,
            width - 0.3, height - 0.8,
            text=description
        )

        p_desc = desc_box.text_frame.paragraphs[0]
        p_desc.font.name = FontConfig.BODY_FONT
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = ColorConfig.DARK_GRAY

        return {
            'background': bg,
            'accent': accent_bar,
            'title': title_box,
            'description': desc_box
        }

    def render_stat_item(
        self,
        x: float,
        y: float,
        width: float,
        number: str,
        label: str,
        description: Optional[str] = None
    ) -> Dict[str, Any]:
        """Render a statistics item with large number and description.

        Args:
            x: X position
            y: Y position
            width: Item width
            number: Statistic number (e.g., "85%")
            label: Short label for stat
            description: Optional longer description

        Returns:
            Dict with shape references
        """
        # Number
        number_box = self.add_textbox(
            x, y, width, LayoutConfig.STATS_NUMBER_HEIGHT,
            text=number
        )

        p_num = number_box.text_frame.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        p_num.font.name = FontConfig.HEADER_FONT
        p_num.font.size = FontConfig.STAT_NUMBER_SIZE
        p_num.font.bold = True
        p_num.font.color.rgb = ColorConfig.ORANGE

        # Label
        label_y = y + LayoutConfig.STATS_NUMBER_HEIGHT
        label_box = self.add_textbox(
            x, label_y, width, LayoutConfig.STATS_LABEL_HEIGHT,
            text=label
        )

        p_label = label_box.text_frame.paragraphs[0]
        p_label.alignment = PP_ALIGN.CENTER
        p_label.font.name = FontConfig.BODY_FONT
        p_label.font.size = Pt(14)
        p_label.font.bold = True
        p_label.font.color.rgb = ColorConfig.DARK_GRAY

        shapes = {
            'number': number_box,
            'label': label_box
        }

        # Description (optional)
        if description:
            desc_y = label_y + LayoutConfig.STATS_LABEL_HEIGHT + 0.1
            desc_box = self.add_textbox(
                x, desc_y, width, LayoutConfig.STATS_DESCRIPTION_HEIGHT,
                text=description
            )

            p_desc = desc_box.text_frame.paragraphs[0]
            p_desc.alignment = PP_ALIGN.CENTER
            p_desc.font.name = FontConfig.BODY_FONT
            p_desc.font.size = Pt(12)
            p_desc.font.color.rgb = ColorConfig.MUTED_GRAY

            shapes['description'] = desc_box

        return shapes

    def render_stats_banner(
        self,
        stats_data: List[Dict[str, str]],
        x: float,
        y: float,
        width: float,
        gap: float = 0.4
    ) -> List[Dict[str, Any]]:
        """Render a horizontal banner of statistics.

        Args:
            stats_data: List of dicts with 'number', 'label', 'description'
            x: Container X position
            y: Container Y position
            width: Container width
            gap: Gap between stat items

        Returns:
            List of stat shape references
        """
        num_stats = len(stats_data)
        stat_height = LayoutConfig.STATS_NUMBER_HEIGHT + LayoutConfig.STATS_LABEL_HEIGHT

        # Calculate stat positions
        layout_boxes = LayoutCalculator.calculate_stats_banner_layout(
            num_stats, x, y, width, stat_height, gap
        )

        rendered_stats = []
        for box, stat_data in zip(layout_boxes, stats_data):
            stat = self.render_stat_item(
                box.x, box.y, box.width,
                number=stat_data.get('number', ''),
                label=stat_data.get('label', ''),
                description=stat_data.get('description')
            )
            rendered_stats.append(stat)

        return rendered_stats

    def render_decorative_shapes(
        self,
        x: float,
        y: float,
        count: int = 3,
        color: Optional[Any] = None
    ) -> List[Any]:
        """Render decorative shapes (for title slide).

        Args:
            x: Starting X position
            y: Y position
            count: Number of shapes to render
            color: Shape color

        Returns:
            List of shape references
        """
        if color is None:
            color = ColorConfig.TAN

        shapes = []
        current_x = x

        for i in range(count):
            shape = self.add_rectangle(
                current_x, y,
                LayoutConfig.DECORATIVE_SHAPE_WIDTH,
                LayoutConfig.DECORATIVE_SHAPE_HEIGHT,
                fill_color=color
            )
            shapes.append(shape)
            current_x += LayoutConfig.DECORATIVE_SHAPE_WIDTH + LayoutConfig.DECORATIVE_SHAPE_SPACING

        return shapes
