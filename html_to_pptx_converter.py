#!/usr/bin/env python3
"""
HTML to PPTX Converter
Converts HTML presentation slides to PowerPoint format

Usage:
    python3 html_to_pptx_converter.py <input.html> <output.pptx>
"""

import sys
import re
import os
import io
import time
import hashlib
from pathlib import Path
from lxml import html as lhtml, etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Import configuration constants
from html_to_pptx.config import (
    LayoutConfig,
    FontConfig,
    ColorConfig,
    SpacingConfig,
)
from html_to_pptx.css_parser import CSSStyleParser
from html_to_pptx.handlers import create_handler_registry

# Import utility helpers
from html_to_pptx.utils import FontStyler, XMLHelper, IndentationPresets

# Image processing imports
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("Warning: PIL/Pillow not available. Image features will be limited.")

try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False
    print("Warning: requests library not available. API image fetching disabled.")


class HTMLToPPTXConverter:
    """Converts HTML slides to PowerPoint presentation"""

    def __init__(self):
        self.prs = Presentation()
        # Set 4:3 aspect ratio (1024x768px at 96 DPI)
        self.prs.slide_width = 9753600  # 1024px in EMU
        self.prs.slide_height = 7315200  # 768px in EMU

        # Import layout constants from config
        self.SLIDE_WIDTH = LayoutConfig.SLIDE_WIDTH
        self.SLIDE_HEIGHT = LayoutConfig.SLIDE_HEIGHT
        self.PADDING = LayoutConfig.PADDING

        # Import font constants from config
        self.HEADER_FONT = FontConfig.HEADER_FONT
        self.BODY_FONT = FontConfig.BODY_FONT

        # Import color mappings from config
        self.COLORS = ColorConfig.COLORS

        # Initialize CSS parser with CSS variables
        self.css_parser = CSSStyleParser(css_variables=ColorConfig.CSS_VARIABLES)

        # Initialize slide handler registry
        self.handler_registry = create_handler_registry(self)

    def parse_color(self, color_str):
        """Parse CSS color to RGBColor"""
        if not color_str:
            return None

        color_str = color_str.strip()

        # Handle hex colors
        if color_str.startswith('#'):
            if color_str in self.COLORS:
                return self.COLORS[color_str]
            # Parse hex manually
            hex_color = color_str.lstrip('#')
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)

        # Handle rgb() colors
        if color_str.startswith('rgb'):
            match = re.search(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
            if match:
                r, g, b = map(int, match.groups())
                return RGBColor(r, g, b)

        return None

    def parse_font_size(self, size_str):
        """Parse CSS font-size to Pt"""
        if not size_str:
            return None

        # Extract number from strings like "40px", "2em", etc.
        match = re.search(r'(\d+)', size_str)
        if match:
            size = int(match.group(1))
            # Assume px for now
            return Pt(size * 0.75)  # Rough conversion from px to pt

        return None

    def add_textbox(self, slide, left, top, width, height):
        """Add a text box to the slide"""
        textbox = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )
        # Enable text wrapping
        textbox.text_frame.word_wrap = True
        # Disable shadow
        textbox.shadow.inherit = False
        return textbox

    def add_text(self, text_frame, text, font_name=None, font_size=None,
                 bold=False, italic=False, color=None, alignment=None, spacing=None):
        """Add text to a text frame with formatting

        Args:
            text_frame: PowerPoint text frame
            text: Text content to add
            font_name: Font family name
            font_size: Font size (Pt object)
            bold: Bold text
            italic: Italic text
            color: Text color (RGBColor)
            alignment: Text alignment (PP_ALIGN constant)
            spacing: Character spacing in points (Pt object). Negative = condensed.
        """
        p = text_frame.paragraphs[0] if len(text_frame.paragraphs) > 0 else text_frame.add_paragraph()
        run = p.add_run()
        run.text = text

        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = font_size
        if bold:
            run.font.bold = True
        if italic:
            run.font.italic = True
        if color:
            run.font.color.rgb = color
        if alignment:
            p.alignment = alignment
        if spacing is not None:
            # Set character spacing via XML (python-pptx doesn't support spacing property)
            rPr = run.font._element
            # Convert Pt to 1/100th of a point (OOXML character spacing unit)
            spacing_hundredths = int(spacing.pt * 100) if hasattr(spacing, 'pt') else int(spacing * 100)
            rPr.set('spc', str(spacing_hundredths))

        return p, run

    def add_footer_system(self, slide, slide_number, total_slides, course_name="", is_dark_bg=False):
        """Add consistent footer with course name and slide number"""
        footer_y = LayoutConfig.FOOTER_Y

        # Course name (left)
        if course_name:
            course_box = self.add_textbox(slide, 0.49, footer_y, 6.5, LayoutConfig.FOOTER_HEIGHT)
            tf = course_box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            run = p.add_run()
            run.text = course_name
            FontStyler.apply_footer_styling(run, is_dark_bg=is_dark_bg, is_slide_number=False)

        # Slide number (right)
        number_box = self.add_textbox(
            slide,
            self.SLIDE_WIDTH - LayoutConfig.FOOTER_NUMBER_X_OFFSET,
            footer_y,
            LayoutConfig.FOOTER_NUMBER_WIDTH,
            LayoutConfig.FOOTER_HEIGHT
        )
        tf = number_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT

        run = p.add_run()
        run.text = str(slide_number)
        FontStyler.apply_footer_styling(run, is_slide_number=True)

    def add_decorative_shapes(self, slide, position='bottom-left'):
        """Add small decorative geometric shapes"""
        if position == 'bottom-left':
            # Calculate positions for three shapes
            x_positions = [
                LayoutConfig.DECORATIVE_SHAPE_X_START,
                LayoutConfig.DECORATIVE_SHAPE_X_START - LayoutConfig.DECORATIVE_SHAPE_SPACING,
                LayoutConfig.DECORATIVE_SHAPE_X_START + LayoutConfig.DECORATIVE_SHAPE_SPACING,
            ]

            # Create three tan rectangles
            for x_pos in x_positions:
                rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x_pos),
                    Inches(LayoutConfig.DECORATIVE_SHAPE_Y),
                    Inches(LayoutConfig.DECORATIVE_SHAPE_WIDTH),
                    Inches(LayoutConfig.DECORATIVE_SHAPE_HEIGHT)
                )
                rect.fill.solid()
                rect.fill.fore_color.rgb = ColorConfig.TAN
                rect.line.fill.background()
                rect.shadow.inherit = False

    def add_separator_line(self, slide, x, y, width=None, height=None, vertical=False, color='#e2e8f0'):
        """Add horizontal or vertical separator line"""
        line_color = self.parse_color(color) or self.COLORS['#e2e8f0']

        if vertical:
            # Vertical line (thin width, tall height)
            width = width or 0.01
            height = height or 3.0
        else:
            # Horizontal line (wide width, thin height)
            width = width or (self.SLIDE_WIDTH - 2*self.PADDING)
            height = height or 0.005

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = line_color
        line.line.fill.background()
        line.shadow.inherit = False

        return line

    def add_icon_above_stat(self, slide, x, y, icon_type='circle', size=0.6):
        """Add simple geometric icon above stat (circle, square, or from resources/icons/)"""
        # Try to load PNG icon first (if available)
        icon_path = f"resources/icons/{icon_type}.png"
        if os.path.exists(icon_path):
            try:
                picture = slide.shapes.add_picture(
                    icon_path,
                    Inches(x), Inches(y),
                    width=Inches(size)
                )
                return picture
            except Exception as e:
                print(f"Warning: Could not load icon {icon_path}: {e}")

        # Fallback to simple geometric shapes
        if icon_type in ['circle', 'dot']:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y),
                Inches(size), Inches(size)
            )
        elif icon_type in ['square', 'box']:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(size), Inches(size)
            )
        else:
            # Default to circle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y),
                Inches(size), Inches(size)
            )

        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['#131313']  # Dark gray for icons
        shape.line.fill.background()
        shape.shadow.inherit = False

        return shape

    def add_image_from_path(self, slide, image_path, x, y, width=None, height=None, send_to_back=True):
        """Add image to slide with optional resizing"""
        if not os.path.exists(image_path):
            print(f"Warning: Image not found: {image_path}")
            return None

        try:
            picture = slide.shapes.add_picture(
                image_path,
                Inches(x), Inches(y),
                width=Inches(width) if width else None,
                height=Inches(height) if height else None
            )

            # Send to back so it doesn't cover text
            if send_to_back:
                slide.shapes._spTree.remove(picture._element)
                slide.shapes._spTree.insert(2, picture._element)

            return picture
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
            return None

    def prepare_image(self, source_path, target_width, target_height, output_path):
        """Resize and crop image to exact dimensions using Pillow"""
        if not PIL_AVAILABLE:
            print("Warning: Pillow not available, cannot prepare image")
            return source_path

        try:
            img = Image.open(source_path)

            # Convert to RGB if necessary (handles RGBA, P, etc.)
            if img.mode not in ['RGB', 'L']:
                img = img.convert('RGB')

            # Calculate aspect ratios
            source_aspect = img.width / img.height
            target_aspect = target_width / target_height

            # Crop to match target aspect ratio (center crop)
            if source_aspect > target_aspect:
                # Image is wider, crop horizontally
                new_width = int(img.height * target_aspect)
                left = (img.width - new_width) // 2
                img = img.crop((left, 0, left + new_width, img.height))
            elif source_aspect < target_aspect:
                # Image is taller, crop vertically
                new_height = int(img.width / target_aspect)
                top = (img.height - new_height) // 2
                img = img.crop((0, top, img.width, top + new_height))

            # Resize to target dimensions
            target_px_width = int(target_width * 96)  # Convert inches to pixels at 96 DPI
            target_px_height = int(target_height * 96)
            img = img.resize((target_px_width, target_px_height), Image.Resampling.LANCZOS)

            # Save
            img.save(output_path, 'JPEG', quality=90, optimize=True)
            return output_path
        except Exception as e:
            print(f"Error preparing image: {e}")
            return source_path

    def fetch_topic_image(self, topic_keyword, output_dir='temp_images', max_width=1280):
        """Fetch image from Pexels API with caching (requires PEXELS_API_KEY env var)"""
        if not REQUESTS_AVAILABLE or not PIL_AVAILABLE:
            return None

        api_key = os.getenv('PEXELS_API_KEY')
        if not api_key:
            # Silently skip if no API key (not an error, just not configured)
            return None

        # Create cache directory
        os.makedirs(output_dir, exist_ok=True)

        # Create cache key from topic keyword
        cache_key = hashlib.md5(topic_keyword.lower().encode()).hexdigest()
        cache_path = os.path.join(output_dir, f"{cache_key}.jpg")

        # Check cache first
        if os.path.exists(cache_path):
            return cache_path

        try:
            # Search Pexels
            url = "https://api.pexels.com/v1/search"
            headers = {'Authorization': api_key}
            params = {
                'query': topic_keyword,
                'per_page': 1,
                'orientation': 'landscape',
                'size': 'large'
            }

            # Rate limiting: simple delay
            time.sleep(0.5)  # Be nice to the API

            response = requests.get(url, headers=headers, params=params, timeout=10)

            if response.status_code == 200:
                data = response.json()
                if data.get('photos') and len(data['photos']) > 0:
                    # Get the large size URL (1280px width)
                    image_url = data['photos'][0]['src']['large']

                    # Download image
                    img_response = requests.get(image_url, timeout=10)
                    img_response.raise_for_status()

                    # Open with Pillow
                    img = Image.open(io.BytesIO(img_response.content))

                    # Convert to RGB if needed
                    if img.mode not in ['RGB', 'L']:
                        img = img.convert('RGB')

                    # Save to cache
                    img.save(cache_path, 'JPEG', quality=90, optimize=True)
                    print(f"✓ Fetched image for '{topic_keyword}'")
                    return cache_path
                else:
                    print(f"No images found for '{topic_keyword}'")
            else:
                print(f"Pexels API error: {response.status_code}")
        except Exception as e:
            print(f"Error fetching image for '{topic_keyword}': {e}")

        return None

    def layout_grid_container(self, container_elem, x_start, y_start, width, height, row_height=None):
        """
        Calculate positions for children in CSS grid layout

        Args:
            container_elem: HTML element with display: grid
            x_start, y_start: Top-left position in inches
            width, height: Available space in inches
            row_height: Height for each row (optional, uses height if not specified)

        Returns:
            List of dicts with {'element': elem, 'x': float, 'y': float, 'width': float, 'height': float}
        """
        # Get computed styles (CSS rules + inline styles)
        styles = self.css_parser.get_computed_style(container_elem)

        # Get grid configuration
        grid_cols = styles.get('grid-template-columns', '1fr')
        gap_value = styles.get('gap', '0px')

        # Parse gap (can be single value or "row-gap column-gap")
        gap_parts = gap_value.split()
        if len(gap_parts) == 2:
            row_gap = self.css_parser.parse_size(gap_parts[0], context='spacing') or 0
            col_gap = self.css_parser.parse_size(gap_parts[1], context='spacing') or 0
        else:
            gap_inches = self.css_parser.parse_size(gap_value, context='spacing') or 0
            row_gap = col_gap = gap_inches

        # Calculate column widths
        col_fractions = self.css_parser.parse_grid(grid_cols)
        num_cols = len(col_fractions)

        total_col_gap = col_gap * (num_cols - 1)
        available_width = width - total_col_gap

        col_widths = [available_width * frac for frac in col_fractions]

        # Get children
        children = list(container_elem)
        positions = []

        # Calculate row height
        if row_height is None:
            num_rows = (len(children) + num_cols - 1) // num_cols  # Ceiling division
            if num_rows > 0:
                total_row_gap = row_gap * (num_rows - 1)
                row_height = (height - total_row_gap) / num_rows
            else:
                row_height = height

        # Position each child
        current_x = x_start
        current_y = y_start
        col_index = 0

        for child in children:
            positions.append({
                'element': child,
                'x': current_x,
                'y': current_y,
                'width': col_widths[col_index],
                'height': row_height
            })

            # Move to next column
            current_x += col_widths[col_index] + col_gap
            col_index += 1

            # Wrap to next row
            if col_index >= num_cols:
                col_index = 0
                current_x = x_start
                current_y += row_height + row_gap

        return positions

    def layout_flex_container(self, container_elem, x_start, y_start, width, height):
        """
        Calculate positions for children in flexbox layout (horizontal only for now)

        Args:
            container_elem: HTML element with display: flex
            x_start, y_start: Top-left position in inches
            width, height: Available space in inches

        Returns:
            List of dicts with {'element': elem, 'x': float, 'y': float, 'width': float, 'height': float}
        """
        # Get computed styles (CSS rules + inline styles)
        styles = self.css_parser.get_computed_style(container_elem)

        # Get flex configuration
        gap_value = styles.get('gap', '0px')
        gap_inches = self.css_parser.parse_size(gap_value, context='spacing') or 0
        justify_content = styles.get('justify-content', 'flex-start')

        # Get children
        children = list(container_elem)
        num_children = len(children)

        if num_children == 0:
            return []

        # Calculate item widths (equal distribution for now)
        total_gap = gap_inches * (num_children - 1)
        item_width = (width - total_gap) / num_children

        positions = []
        current_x = x_start

        # Adjust starting position based on justify-content
        if justify_content == 'center':
            total_content_width = (item_width * num_children) + total_gap
            current_x = x_start + (width - total_content_width) / 2
        elif justify_content == 'flex-end':
            total_content_width = (item_width * num_children) + total_gap
            current_x = x_start + (width - total_content_width)
        elif justify_content == 'space-between' and num_children > 1:
            item_width = width / num_children
            gap_inches = 0  # Space is distributed between items

        # Position each child
        for child in children:
            positions.append({
                'element': child,
                'x': current_x,
                'y': y_start,
                'width': item_width,
                'height': height
            })

            current_x += item_width + gap_inches

        return positions

    def apply_shadow(self, shape, shadow_params):
        """
        Apply shadow effect to PPTX shape

        Args:
            shape: python-pptx shape object
            shadow_params: Dict from parse_box_shadow() with offset_x, offset_y, blur, color, transparency
        """
        if not shadow_params:
            return

        try:
            shadow = shape.shadow
            shadow.inherit = False
            shadow.visible = True

            # Shadow position (use offset_y as distance, calculate angle)
            offset_x = shadow_params.get('offset_x', 0)
            offset_y = shadow_params.get('offset_y', 0)
            blur = shadow_params.get('blur', 0)

            # Calculate distance and angle from offsets
            import math
            distance = math.sqrt(offset_x**2 + offset_y**2)
            if distance > 0:
                # Convert to EMU (914400 EMU = 1 inch)
                shadow.distance = int(distance * 914400)

                # Calculate angle (0 = right, 90 = down, 180 = left, 270 = up)
                angle = math.degrees(math.atan2(offset_y, offset_x))
                # PPTX uses 0 = right, 90 = down
                shadow.angle = angle

            # Shadow blur
            if blur > 0:
                shadow.blur_radius = int(blur * 914400)

            # Shadow color and transparency (requires oxml manipulation)
            # Note: python-pptx doesn't expose shadow color directly, would need oxml
            # For now, skip color setting - shadows will use default

        except Exception as e:
            print(f"Warning: Could not apply shadow: {e}")

    def apply_border(self, shape, border_params):
        """
        Apply border to PPTX shape

        Args:
            shape: python-pptx shape object
            border_params: Dict from parse_border() with width, color, style
        """
        if not border_params:
            # No border
            shape.line.fill.background()
            return

        try:
            # Set border color
            color = border_params.get('color')
            if color:
                shape.line.color.rgb = color

            # Set border width
            width = border_params.get('width')
            if width:
                shape.line.width = width

            # Border style (solid, dashed, etc.)
            style = border_params.get('style', 'solid')
            # Note: python-pptx has limited line style support
            # MSO_LINE_DASH_STYLE could be used for dashed lines

        except Exception as e:
            print(f"Warning: Could not apply border: {e}")

    def apply_slide_background(self, slide, background_color):
        """
        Apply background color to slide

        Args:
            slide: python-pptx slide object
            background_color: RGBColor object
        """
        if not background_color:
            return

        try:
            bg_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                self.prs.slide_width,
                self.prs.slide_height
            )
            bg_rect.fill.solid()
            bg_rect.fill.fore_color.rgb = background_color
            bg_rect.line.fill.background()
            bg_rect.shadow.inherit = False

            # Move to back
            slide.shapes._spTree.remove(bg_rect._element)
            slide.shapes._spTree.insert(2, bg_rect._element)

        except Exception as e:
            print(f"Warning: Could not apply background: {e}")

    def extract_text_content(self, element):
        """Extract text content from HTML element, normalizing whitespace like HTML does"""
        if element is None:
            return ""
        text = etree.tostring(element, method='text', encoding='unicode')
        # Normalize whitespace: collapse newlines and multiple spaces to single space
        import re
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def detect_list_marker(self, li_element):
        """Detect what kind of marker this list item uses"""
        # Get the full text content
        full_text = self.extract_text_content(li_element).strip()

        # Check for markers at the start
        if full_text.startswith('✓'):
            return '✓'
        elif full_text.startswith('✗'):
            return '✗'
        elif full_text.startswith('☐'):
            return '☐'
        elif full_text.startswith('☑'):
            return '☑'
        elif full_text.startswith('☒'):
            return '☒'
        else:
            return '•'  # Default bullet

    def is_numbered_list(self, ul_element):
        """Check if this unordered list should be treated as numbered"""
        li_items = ul_element.findall('./li')
        if len(li_items) < 2:
            return False

        # Check if first few items start with sequential numbers
        for idx, li in enumerate(li_items[:min(3, len(li_items))], 1):
            text = self.extract_text_content(li).strip()

            # Check for patterns like "1.", "2.", "3." at the start
            if text.startswith(f'{idx}.'):
                continue
            # Check for patterns like "1)", "2)", "3)"
            elif text.startswith(f'{idx})'):
                continue
            # Check for word patterns: "First", "Second", "Third", etc.
            elif idx == 1 and (text.lower().startswith('first') or text.lower().startswith('1st')):
                continue
            elif idx == 2 and (text.lower().startswith('second') or text.lower().startswith('2nd')):
                continue
            elif idx == 3 and (text.lower().startswith('third') or text.lower().startswith('3rd')):
                continue
            else:
                return False

        return True

    def add_formatted_text(self, paragraph, element, skip_leading_marker=False, skip_leading_number=False):
        """Add formatted text from HTML element to paragraph, preserving bold/italic/color"""
        if element is None:
            return

        import re  # Import once at the top

        # Handle direct text
        if element.text:
            # Normalize whitespace: collapse newlines and multiple spaces to single space
            text = re.sub(r'\s+', ' ', element.text)
            text = text.lstrip()  # Only strip leading space
            if text:
                # Skip leading marker if requested (✓, ✗, ☐, etc.)
                if skip_leading_marker and text and text[0] in ['✓', '✗', '☐', '☑', '☒']:
                    text = text[1:].lstrip()

                # Skip leading number if requested (1., 2., First, etc.)
                if skip_leading_number:
                    import re
                    # Remove patterns like "1.", "2)", "First,", "First:", etc.
                    text = re.sub(r'^\d+[\.\)]\s*', '', text)  # "1. " or "2) "
                    text = re.sub(r'^(First|Second|Third|Fourth|Fifth|Sixth|Seventh|Eighth|Ninth|Tenth)[:\,\s]\s*', '', text, flags=re.IGNORECASE)
                    text = re.sub(r'^(1st|2nd|3rd|4th|5th|6th|7th|8th|9th|10th)[:\,\s]\s*', '', text, flags=re.IGNORECASE)

                if text:  # Only add if there's text left
                    run = paragraph.add_run()
                    run.text = text

                    # Check for styling
                    if element.tag in ['strong', 'b']:
                        run.font.bold = True
                        run.font.color.rgb = ColorConfig.ORANGE  # Bold text is orange
                    if element.tag in ['em', 'i']:
                        run.font.italic = True

        # Handle child elements recursively
        for child in element:
            # Skip <span> elements that only contain bullet characters
            if child.tag == 'span':
                span_text = (child.text or '').strip()
                if span_text in ['•', '◦', '▪', '■', '□', '○', '●']:
                    # Skip this span entirely - it's just a custom bullet
                    continue

            if child.tag in ['strong', 'b']:
                if child.text:
                    child_text = child.text
                    # Skip leading marker in strong tags too
                    if skip_leading_marker and child_text and child_text[0] in ['✓', '✗', '☐', '☑', '☒']:
                        child_text = child_text[1:].strip()

                    # Skip leading number if requested
                    if skip_leading_number:
                        import re
                        child_text = re.sub(r'^\d+[\.\)]\s*', '', child_text)
                        child_text = re.sub(r'^(First|Second|Third|Fourth|Fifth|Sixth|Seventh|Eighth|Ninth|Tenth)[:\,\s]\s*', '', child_text, flags=re.IGNORECASE)
                        child_text = re.sub(r'^(1st|2nd|3rd|4th|5th|6th|7th|8th|9th|10th)[:\,\s]\s*', '', child_text, flags=re.IGNORECASE)

                    if child_text:
                        run = paragraph.add_run()
                        run.text = child_text
                        run.font.bold = True

                        # Check for inline color styles (takes precedence)
                        style_attr = child.get('style', '')
                        color_applied = False
                        if 'color:' in style_attr:
                            # Extract color from inline style
                            color = self.css_parser.parse_inline_style(style_attr).get('color')
                            if color:
                                rgb_color = self.css_parser.parse_color(color)
                                if rgb_color:
                                    run.font.color.rgb = rgb_color
                                    color_applied = True

                        # If no inline color, apply default orange for bold
                        if not color_applied:
                            run.font.color.rgb = ColorConfig.ORANGE
            elif child.tag in ['em', 'i']:
                if child.text:
                    run = paragraph.add_run()
                    run.text = child.text
                    run.font.italic = True
            else:
                # Regular text in child
                if child.text:
                    run = paragraph.add_run()
                    run.text = child.text

            # Handle tail text (text after closing tag)
            if child.tail:
                # Normalize whitespace: collapse newlines and multiple spaces to single space
                tail_text = re.sub(r'\s+', ' ', child.tail)
                tail_text = tail_text.rstrip()  # Only strip trailing space
                if tail_text:
                    run = paragraph.add_run()
                    run.text = tail_text

    def handle_title_slide(self, slide, html_slide):
        """Handle title slide layout"""
        # Extract title
        title_elem = html_slide.find('.//*[@class="title-content"]/h1')
        subtitle_elem = html_slide.find('.//*[@class="subtitle"]')
        author_elem = html_slide.find('.//*[@class="author"]')

        # Add title - top left aligned
        if title_elem is not None:
            title_box = self.add_textbox(slide, self.PADDING, self.PADDING,
                                        self.SLIDE_WIDTH - 2*self.PADDING, 1.2)
            title_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            self.add_text(title_box.text_frame, self.extract_text_content(title_elem),
                         font_name=self.HEADER_FONT, font_size=Pt(52), bold=True,
                         color=ColorConfig.ORANGE, alignment=PP_ALIGN.LEFT,
                         spacing=FontConfig.TITLE_LETTER_SPACING)

        # Add subtitle
        if subtitle_elem is not None:
            subtitle_box = self.add_textbox(slide, self.PADDING, 1.5,
                                           self.SLIDE_WIDTH - 2*self.PADDING, 0.8)
            subtitle_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            self.add_text(subtitle_box.text_frame, self.extract_text_content(subtitle_elem),
                         font_name=self.BODY_FONT, font_size=Pt(27),
                         color=ColorConfig.ORANGE, alignment=PP_ALIGN.LEFT)

        # Add author with line breaks preserved
        if author_elem is not None:
            author_box = self.add_textbox(slide, self.PADDING, 2.6,
                                         self.SLIDE_WIDTH - 2*self.PADDING, 1.5)
            tf = author_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            tf.clear()

            # Parse HTML to preserve line breaks
            import html as html_module
            # Get the inner HTML with br tags
            author_html = etree.tostring(author_elem, encoding='unicode', method='html')
            # Split by <br> or <br/>
            parts = re.split(r'<br\s*/?>', author_html)

            for i, part in enumerate(parts):
                # Extract text from HTML fragment
                # Remove opening and closing tags
                text = re.sub(r'<[^>]+>', '', part).strip()
                text = html_module.unescape(text)

                if text:
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(4)

                    # Check if text has <strong> tags for bold
                    if '<strong>' in part or '<b>' in part:
                        # Parse bold sections
                        bold_pattern = r'<(?:strong|b)>([^<]+)</(?:strong|b)>'
                        remaining = part
                        last_end = 0

                        for match in re.finditer(bold_pattern, remaining):
                            # Add text before bold
                            before_text = re.sub(r'<[^>]+>', '', remaining[last_end:match.start()]).strip()
                            before_text = html_module.unescape(before_text)
                            if before_text:
                                run = p.add_run()
                                run.text = before_text
                                FontStyler.apply_muted_text_styling(run)

                            # Add bold text
                            bold_text = html_module.unescape(match.group(1))
                            run = p.add_run()
                            run.text = bold_text
                            FontStyler.apply_muted_text_styling(run, bold=True)

                            last_end = match.end()

                        # Add remaining text after last bold
                        after_text = re.sub(r'<[^>]+>', '', remaining[last_end:]).strip()
                        after_text = html_module.unescape(after_text)
                        if after_text:
                            run = p.add_run()
                            run.text = after_text
                            FontStyler.apply_muted_text_styling(run)
                    else:
                        # Simple text, no bold
                        run = p.add_run()
                        run.text = text
                        FontStyler.apply_muted_text_styling(run)

        # Add decorative shapes
        self.add_decorative_shapes(slide, 'bottom-left')

    def handle_section_break_slide(self, slide, html_slide):
        """Handle section break slide layout"""
        # Add colored background
        background = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,
            self.prs.slide_width,
            self.prs.slide_height
        )
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['#7373b0']
        background.shadow.inherit = False

        # Move to back
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)

        # Extract title
        title_elem = html_slide.find('.//*[@class="section-title"]')
        if title_elem is not None:
            title_box = self.add_textbox(slide, self.PADDING, 3.0,
                                        self.SLIDE_WIDTH - 2*self.PADDING, 2.0)
            tf = title_box.text_frame
            tf.word_wrap = True

            # Handle <br> tags by splitting into paragraphs
            # Get text before first <br>
            first_p = tf.paragraphs[0]
            first_p.alignment = PP_ALIGN.LEFT

            if title_elem.text:
                run = first_p.add_run()
                run.text = title_elem.text
                run.font.name = self.HEADER_FONT
                run.font.size = Pt(38)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                # Apply condensed letter spacing via XML (1/100th of a point)
                spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
                run.font._element.set('spc', str(spacing_hundredths))

            # Process <br> elements and text after them
            for br in title_elem:
                if br.tag == 'br':
                    # Add a new paragraph for text after <br>
                    if br.tail:
                        p = tf.add_paragraph()
                        p.alignment = PP_ALIGN.LEFT
                        run = p.add_run()
                        run.text = br.tail
                        run.font.name = self.HEADER_FONT
                        run.font.size = Pt(38)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        # Apply condensed letter spacing via XML (1/100th of a point)
                        spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
                        run.font._element.set('spc', str(spacing_hundredths))

        # Extract and add subtitle if present
        subtitle_elem = html_slide.find('.//*[@class="section-subtitle"]')
        if subtitle_elem is not None:
            subtitle_box = self.add_textbox(slide, self.PADDING, 4.3,
                                          self.SLIDE_WIDTH - 2*self.PADDING, 1.0)
            sub_tf = subtitle_box.text_frame
            sub_tf.word_wrap = True
            sub_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            sub_p = sub_tf.paragraphs[0]
            sub_p.alignment = PP_ALIGN.LEFT

            sub_run = sub_p.add_run()
            sub_run.text = self.extract_text_content(subtitle_elem)
            sub_run.font.name = self.BODY_FONT
            sub_run.font.size = Pt(24)
            sub_run.font.color.rgb = RGBColor(255, 255, 255)
            sub_run.font.italic = True

    def handle_big_number_slide(self, slide, html_slide):
        """Handle big number slide with centered statistic"""
        # Apply cream background
        self.apply_slide_background(slide, ColorConfig.CREAM)

        # Find number content container
        number_content = html_slide.find('.//*[@class="number-content"]')
        if number_content is None:
            return

        # Extract big number
        big_number_elem = number_content.find('.//*[@class="big-number"]')
        if big_number_elem is not None:
            number_text = self.extract_text_content(big_number_elem).strip()

            # Create centered text box for big number
            number_box = self.add_textbox(
                slide,
                1.0,  # Centered horizontally with padding
                LayoutConfig.BIG_NUMBER_Y,
                self.SLIDE_WIDTH - 2.0,  # Full width with padding
                LayoutConfig.BIG_NUMBER_HEIGHT
            )
            tf = number_box.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            run = p.add_run()
            run.text = number_text
            run.font.name = FontConfig.HEADER_FONT
            run.font.size = FontConfig.BIG_NUMBER_SIZE
            run.font.bold = True
            run.font.color.rgb = ColorConfig.ORANGE

        # Extract explanation text
        explanation_elem = number_content.find('.//*[@class="number-explanation"]')
        if explanation_elem is not None:
            explanation_text = self.extract_text_content(explanation_elem).strip()

            # Create text box for explanation below number
            explanation_box = self.add_textbox(
                slide,
                1.0,
                LayoutConfig.BIG_NUMBER_Y + LayoutConfig.BIG_NUMBER_EXPLANATION_Y_OFFSET,
                self.SLIDE_WIDTH - 2.0,
                LayoutConfig.BIG_NUMBER_EXPLANATION_HEIGHT
            )
            tf = explanation_box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            run = p.add_run()
            run.text = explanation_text
            FontStyler.apply_muted_text_styling(run, font_size=FontConfig.BODY_MEDIUM)

    def handle_content_slide(self, slide, html_slide):
        """Handle standard content slide with title and bullets"""
        # Get slide classes to check for dark background
        slide_classes = html_slide.get('class', '')
        is_dark = 'dark-bg' in slide_classes

        # Apply slide background (cream for light slides, dark for dark slides)
        if is_dark:
            self.apply_slide_background(slide, ColorConfig.DARK_GRAY)
        else:
            self.apply_slide_background(slide, ColorConfig.CREAM)

        # Extract title
        title_elem = html_slide.find('.//h2[@class="slide-title"]')
        if title_elem is None:
            title_elem = html_slide.find('.//h2')

        # Delete placeholder shapes if using Title and Content layout
        # We'll add our own text boxes for better control
        shapes_to_delete = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                shapes_to_delete.append(shape)

        for shape in shapes_to_delete:
            sp = shape.element
            sp.getparent().remove(sp)

        # Add title using TextRenderer for consistency
        if title_elem is not None:
            from html_to_pptx.renderers import TextRenderer
            text_renderer = TextRenderer(slide)
            title_text = self.extract_text_content(title_elem)
            text_renderer.render_slide_title(title_text, is_dark_bg=is_dark)

        # Extract content body (try multiple class names)
        content_body = html_slide.find('.//*[@class="content-body"]')
        if content_body is None:
            content_body = html_slide.find('.//*[@class="slide-content"]')
        if content_body is None:
            content_body = html_slide

        # Check for card layout (support both class names)
        cards_container = html_slide.find('.//*[@class="content-cards"]')
        if cards_container is None:
            cards_container = html_slide.find('.//*[@class="card-container"]')
        if cards_container is not None:
            # Start cards after title with proper gap
            self.handle_card_layout(slide, cards_container, y_start=LayoutConfig.CONTENT_START_Y)
            return

        # Check for stats banner
        stats_banner = html_slide.find('.//*[@class="stats-banner"]')
        if stats_banner is not None:
            # Start stats after title with proper gap
            self.handle_stats_banner(slide, stats_banner, y_start=LayoutConfig.CONTENT_START_Y, is_dark_bg=is_dark)
            return

        # Check for grid layout (using computed styles to detect CSS classes)
        grid_div = None
        for child in content_body:
            if child.tag == 'div':
                computed = self.css_parser.get_computed_style(child)
                if computed.get('display') == 'grid':
                    grid_div = child
                    break

        # Handle grid layout separately - process in document order
        if grid_div is not None:
            y_pos = LayoutConfig.CONTENT_START_Y

            # Process all children in document order
            for elem in content_body:
                # Handle paragraphs
                if elem.tag == 'p':
                    p_text = self.extract_text_content(elem).strip()
                    if p_text:
                        style = elem.get('style', '')

                        # Create text box for the paragraph
                        p_box = self.add_textbox(slide, self.PADDING, y_pos,
                                               self.SLIDE_WIDTH - 2*self.PADDING, 0.8)
                        p_tf = p_box.text_frame
                        p_tf.word_wrap = True
                        p = p_tf.paragraphs[0]

                        if 'text-align: center' in style:
                            p.alignment = PP_ALIGN.CENTER
                        else:
                            p.alignment = PP_ALIGN.LEFT

                        self.add_formatted_text(p, elem)

                        # Apply styling
                        for run in p.runs:
                            # Check for font family in style first
                            if 'font-family:' in style:
                                if 'var(--font-header)' in style or self.HEADER_FONT.lower() in style.lower():
                                    run.font.name = self.HEADER_FONT
                                elif not run.font.name:
                                    run.font.name = self.BODY_FONT
                            elif not run.font.name:
                                run.font.name = self.BODY_FONT

                            if 'font-size:' in style:
                                import re
                                size_match = re.search(r'font-size:\s*(\d+)px', style)
                                if size_match:
                                    run.font.size = Pt(int(size_match.group(1)) * 0.75)
                            elif not run.font.size:
                                run.font.size = Pt(16)

                            if 'font-weight: 600' in style or 'font-weight: bold' in style:
                                run.font.bold = True

                            if 'color:' in style:
                                color_match = re.search(r'color:\s*([^;]+)', style)
                                if color_match:
                                    color = self.parse_color(color_match.group(1).strip())
                                    if color:
                                        run.font.color.rgb = color
                                    else:
                                        # Set colors: orange for bold, appropriate color for regular based on background
                                        if run.font.bold:
                                            run.font.color.rgb = ColorConfig.ORANGE
                                        else:
                                            run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']
                            else:
                                # Set colors: orange for bold, appropriate color for regular based on background
                                if run.font.bold:
                                    run.font.color.rgb = ColorConfig.ORANGE
                                else:
                                    run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

                        y_pos += 0.8

                # Handle grid div
                elif elem == grid_div:
                    grid_columns = grid_div.findall('./div')
                    if not grid_columns:
                        continue

                    # Check if grid items have the .grid-item class (styled containers)
                    first_column = grid_columns[0] if grid_columns else None
                    is_styled_grid = first_column is not None and 'grid-item' in first_column.get('class', '')

                    if is_styled_grid:
                        # Use CSS grid layout engine for styled grid items
                        num_cols = len(grid_columns)
                        col_width = (self.SLIDE_WIDTH - 2*self.PADDING - (num_cols-1)*0.3) / num_cols
                        grid_item_height = 2.5  # Height for each styled grid item

                        for idx, column in enumerate(grid_columns):
                            x_pos = self.PADDING + idx * (col_width + 0.3)
                            self._create_grid_item(slide, column, x_pos, y_pos, col_width, grid_item_height)

                        y_pos += grid_item_height + 0.3
                    else:
                        # Legacy grid rendering (text only, no styled containers)
                        num_cols = len(grid_columns)
                        col_width = (self.SLIDE_WIDTH - 2*self.PADDING - (num_cols-1)*0.3) / num_cols

                        for idx, column in enumerate(grid_columns):
                            x_pos = self.PADDING + idx * (col_width + 0.3)
                            col_y_pos = y_pos

                            # Extract heading from column (h3, h4, or p)
                            heading_elem = column.find('./h3')
                            if heading_elem is None:
                                heading_elem = column.find('./h4')
                            if heading_elem is None:
                                heading_elem = column.find('./p')

                            # Extract list or paragraph content
                            list_elem = column.find('./ul')
                            para_elem = None
                            if list_elem is None:
                                # Look for paragraph after heading
                                for child in column:
                                    if child.tag == 'p' and child != heading_elem:
                                        para_elem = child
                                        break

                            if heading_elem is not None:
                                # Add heading
                                heading_text = self.extract_text_content(heading_elem).strip()
                                if heading_text:
                                    heading_box = self.add_textbox(slide, x_pos, col_y_pos,
                                                                  col_width, 0.4)
                                    h_tf = heading_box.text_frame
                                    h_tf.word_wrap = True
                                    h_p = h_tf.paragraphs[0]
                                    h_p.alignment = PP_ALIGN.LEFT

                                    self.add_formatted_text(h_p, heading_elem)

                                    # Apply heading styles
                                    style = heading_elem.get('style', '')
                                    for run in h_p.runs:
                                        # Use header font for grid headings
                                        run.font.name = FontConfig.HEADER_FONT

                                        if 'font-size:' in style:
                                            import re
                                            size_match = re.search(r'font-size:\s*(\d+)px', style)
                                            if size_match:
                                                run.font.size = Pt(int(size_match.group(1)) * 0.75)
                                        else:
                                            run.font.size = Pt(16)

                                        run.font.bold = True

                                        # Apply color: check inline style first, then default to orange
                                        if 'color:' in style:
                                            color_match = re.search(r'color:\s*([^;]+)', style)
                                            if color_match:
                                                color = self.parse_color(color_match.group(1).strip())
                                                if color:
                                                    run.font.color.rgb = color
                                                else:
                                                    # Fallback to orange if parse fails
                                                    run.font.color.rgb = ColorConfig.ORANGE
                                        else:
                                            # Default: grid headings are orange (accent color)
                                            run.font.color.rgb = ColorConfig.ORANGE

                                    col_y_pos += 0.35  # Reduced gap between heading and list

                            if list_elem is not None:
                                # Count items to determine height (reduced for tighter spacing)
                                li_items = list_elem.findall('./li')
                                list_height = len(li_items) * 0.3 + 0.2

                                # Add list
                                list_box = self.add_textbox(slide, x_pos, col_y_pos,
                                                           col_width, list_height)
                                l_tf = list_box.text_frame
                                l_tf.word_wrap = True
                                l_tf.clear()

                                # Add list items
                                first_item = True
                                for li in li_items:
                                    if first_item:
                                        l_p = l_tf.paragraphs[0] if len(l_tf.paragraphs) > 0 else l_tf.add_paragraph()
                                        first_item = False
                                    else:
                                        l_p = l_tf.add_paragraph()

                                    l_p.level = 0
                                    l_p.alignment = PP_ALIGN.LEFT
                                    l_p.line_spacing = SpacingConfig.LINE_SPACING
                                    l_p.space_before = SpacingConfig.LIST_ITEM_SPACE_BEFORE
                                    l_p.space_after = SpacingConfig.LIST_ITEM_SPACE_AFTER

                                    # Enable bullet formatting
                                    marker = self.detect_list_marker(li)
                                    XMLHelper.add_bullet_with_indent(l_p, marker=marker, font_name=self.BODY_FONT)

                                    # Add bullet color (orange)
                                    pPr = l_p._element.get_or_add_pPr()
                                    buClr_xml = f'''<a:buClr {nsdecls("a")}>
                                        <a:srgbClr val="ed5e29"/>
                                    </a:buClr>'''
                                    buClr = parse_xml(buClr_xml)
                                    pPr.append(buClr)

                                    # Add formatted text
                                    self.add_formatted_text(l_p, li, skip_leading_marker=True)

                                    # Apply styling
                                    list_style = list_elem.get('style', '')
                                    for run in l_p.runs:
                                        if not run.font.name:
                                            run.font.name = self.BODY_FONT

                                        if 'font-size:' in list_style:
                                            import re
                                            size_match = re.search(r'font-size:\s*(\d+)px', list_style)
                                            if size_match:
                                                run.font.size = Pt(int(size_match.group(1)) * 0.75)
                                        elif not run.font.size:
                                            run.font.size = Pt(16)

                                        # Set colors: orange for bold, appropriate color for regular based on background
                                        if run.font.bold:
                                            run.font.color.rgb = ColorConfig.ORANGE
                                        else:
                                            run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

                            # Render paragraph if present (and no list)
                            if para_elem is not None and list_elem is None:
                                para_text = self.extract_text_content(para_elem).strip()
                                if para_text:
                                    # Add paragraph
                                    para_box = self.add_textbox(slide, x_pos, col_y_pos,
                                                               col_width, 0.6)
                                    p_tf = para_box.text_frame
                                    p_tf.word_wrap = True
                                    p = p_tf.paragraphs[0]
                                    p.alignment = PP_ALIGN.LEFT

                                    self.add_formatted_text(p, para_elem)

                                    # Apply paragraph styles
                                    para_style = para_elem.get('style', '')
                                    for run in p.runs:
                                        if not run.font.name:
                                            run.font.name = self.BODY_FONT

                                        if 'font-size:' in para_style:
                                            import re
                                            size_match = re.search(r'font-size:\s*(\d+)px', para_style)
                                            if size_match:
                                                run.font.size = Pt(int(size_match.group(1)) * 0.75)
                                        elif not run.font.size:
                                            run.font.size = Pt(14)

                                        # Set colors: orange for bold, appropriate color for regular based on background
                                        if run.font.bold:
                                            run.font.color.rgb = ColorConfig.ORANGE
                                        else:
                                            run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

                        y_pos += 2.5  # Move down after grid

                # Handle other div elements (like Vietnamese context boxes)
                elif elem.tag == 'div' and elem != grid_div:
                    # This is a styled div (not the grid itself)
                    div_paras = elem.findall('./p')
                    if div_paras:
                        # Calculate height based on content
                        div_height = len(div_paras) * 0.4 + 0.3

                        # Create text box for div content
                        div_box = self.add_textbox(slide, self.PADDING, y_pos,
                                                   self.SLIDE_WIDTH - 2*self.PADDING, div_height)
                        div_tf = div_box.text_frame
                        div_tf.word_wrap = True
                        div_tf.clear()

                        # Add all paragraphs from the div
                        first_para = True
                        for div_p in div_paras:
                            if first_para:
                                p = div_tf.paragraphs[0] if len(div_tf.paragraphs) > 0 else div_tf.add_paragraph()
                                first_para = False
                            else:
                                p = div_tf.add_paragraph()

                            p.alignment = PP_ALIGN.LEFT
                            p.line_spacing = SpacingConfig.LINE_SPACING
                            p.space_after = SpacingConfig.PARAGRAPH_SPACE_AFTER

                            self.add_formatted_text(p, div_p)

                            # Apply styling
                            style = div_p.get('style', '')
                            for run in p.runs:
                                if not run.font.name:
                                    run.font.name = self.BODY_FONT

                                if 'font-size:' in style:
                                    import re
                                    size_match = re.search(r'font-size:\s*(\d+)px', style)
                                    if size_match:
                                        run.font.size = Pt(int(size_match.group(1)) * 0.75)
                                elif not run.font.size:
                                    run.font.size = Pt(16)

                                if 'font-weight: 600' in style or 'font-weight: bold' in style:
                                    run.font.bold = True

                                if 'color:' in style:
                                    color_match = re.search(r'color:\s*([^;]+)', style)
                                    if color_match:
                                        color = self.parse_color(color_match.group(1).strip())
                                        if color:
                                            run.font.color.rgb = color
                                        else:
                                            # Set colors: orange for bold, appropriate color for regular based on background
                                            if run.font.bold:
                                                run.font.color.rgb = ColorConfig.ORANGE
                                            else:
                                                run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']
                                else:
                                    # Set colors: orange for bold, appropriate color for regular based on background
                                    if run.font.bold:
                                        run.font.color.rgb = ColorConfig.ORANGE
                                    else:
                                        run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

                        y_pos += div_height + 0.2

            return  # Grid layout handled, exit method

        # Check for mixed content (paragraphs + styled divs) before standard handling
        direct_children = list(content_body)
        has_styled_divs = False
        for child in direct_children:
            if child.tag == 'div':
                style = child.get('style', '')
                # Check direct child - only trigger for cards (background or border)
                if ('background:' in style or 'border-left:' in style):
                    has_styled_divs = True
                    break
                # Check nested children (one level deep)
                for nested in child:
                    if nested.tag == 'div':
                        nested_style = nested.get('style', '')
                        if ('background:' in nested_style or 'border-left:' in nested_style):
                            has_styled_divs = True
                            break
            if has_styled_divs:
                break

        if has_styled_divs:
            # Process content in document order
            import re
            y_pos = LayoutConfig.CONTENT_START_Y

            for elem in direct_children:
                # Handle direct paragraphs
                if elem.tag == 'p':
                    p_text = self.extract_text_content(elem).strip()
                    if p_text:
                        style = elem.get('style', '')
                        # Parse font size from style
                        import re
                        font_size = Pt(16)
                        if 'font-size:' in style:
                            size_match = re.search(r'font-size:\s*(\d+)px', style)
                            if size_match:
                                font_size = Pt(int(size_match.group(1)) * 0.75)

                        # Create text box for paragraph
                        p_box = self.add_textbox(slide, LayoutConfig.PADDING, y_pos,
                                               self.SLIDE_WIDTH - 2*LayoutConfig.PADDING, 0.6)
                        p_tf = p_box.text_frame
                        p_tf.word_wrap = True
                        p = p_tf.paragraphs[0]
                        p.alignment = PP_ALIGN.LEFT

                        # Detect if this is a subheading paragraph
                        is_subheading = False
                        strong_elem = elem.find('.//strong')
                        if strong_elem is not None:
                            if len(p_text) < 100 and p_text.endswith(':'):
                                is_subheading = True

                        # Apply appropriate spacing
                        p.line_spacing = SpacingConfig.LINE_SPACING
                        if is_subheading:
                            p.space_before = SpacingConfig.SUBHEADING_SPACE_BEFORE
                            p.space_after = SpacingConfig.SUBHEADING_SPACE_AFTER
                        else:
                            p.space_before = SpacingConfig.PARAGRAPH_SPACE_BEFORE
                            p.space_after = SpacingConfig.PARAGRAPH_SPACE_AFTER

                        self.add_formatted_text(p, elem)

                        # Apply styling
                        for run in p.runs:
                            if not run.font.name:
                                run.font.name = FontConfig.BODY_FONT
                            if not run.font.size:
                                run.font.size = font_size
                            run.font.color.rgb = ColorConfig.CREAM if is_dark else ColorConfig.DARK_GRAY

                        y_pos += 0.7

                # Handle styled divs (like white cards with borders)
                elif elem.tag == 'div':
                    style = elem.get('style', '')

                    # Check for white card styling
                    has_border_left = 'border-left:' in style
                    has_background = 'background:' in style

                    # If this is a wrapper div without styling, process its children
                    if not has_border_left and not has_background:
                        # Process children of wrapper div
                        for child_elem in elem:
                            if child_elem.tag == 'p':
                                p_text = self.extract_text_content(child_elem).strip()
                                if p_text:
                                    child_style = child_elem.get('style', '')
                                    # Parse font size from style
                                    font_size = Pt(16)
                                    if 'font-size:' in child_style:
                                        size_match = re.search(r'font-size:\s*(\d+)px', child_style)
                                        if size_match:
                                            font_size = Pt(int(size_match.group(1)) * 0.75)

                                    # Create text box for paragraph
                                    p_box = self.add_textbox(slide, LayoutConfig.PADDING, y_pos,
                                                           self.SLIDE_WIDTH - 2*LayoutConfig.PADDING, 0.6)
                                    p_tf = p_box.text_frame
                                    p_tf.word_wrap = True
                                    p = p_tf.paragraphs[0]
                                    p.alignment = PP_ALIGN.LEFT

                                    self.add_formatted_text(p, child_elem)

                                    # Apply styling
                                    for run in p.runs:
                                        if not run.font.name:
                                            run.font.name = FontConfig.BODY_FONT
                                        if not run.font.size:
                                            run.font.size = font_size
                                        run.font.color.rgb = ColorConfig.CREAM if is_dark else ColorConfig.DARK_GRAY

                                    y_pos += 1.0  # More spacing before card

                            elif child_elem.tag == 'div':
                                child_style = child_elem.get('style', '')
                                child_has_border = 'border-left:' in child_style
                                child_has_bg = 'background:' in child_style

                                if child_has_border or child_has_bg:
                                    # This is the styled card
                                    has_border_left = child_has_border
                                    has_background = child_has_bg
                                    elem = child_elem  # Process this element
                                    style = child_style
                                    # Continue to card rendering below
                                    break

                        # If we processed a wrapper, check if we need to render a card
                        if not (has_border_left or has_background):
                            continue

                    if has_border_left or has_background:
                        # Calculate card height based on content
                        h3_elem = elem.find('./h3')
                        ul_elem = elem.find('./ul')

                        card_height = 0.3  # Top padding
                        if h3_elem is not None:
                            card_height += 0.5  # h3 height
                        if ul_elem is not None:
                            li_count = len(ul_elem.findall('./li'))
                            card_height += li_count * 0.35 + 0.2  # List items + spacing
                        card_height += 0.3  # Bottom padding

                        # Create a card-like shape for this div
                        card_y = y_pos
                        card_width = self.SLIDE_WIDTH - 2*LayoutConfig.PADDING

                        # Add white background card (no shadow)
                        from pptx.enum.shapes import MSO_SHAPE

                        # White card background
                        card_bg = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(LayoutConfig.PADDING),
                            Inches(card_y),
                            Inches(card_width),
                            Inches(card_height)
                        )
                        card_bg.fill.solid()
                        card_bg.fill.fore_color.rgb = ColorConfig.WHITE
                        card_bg.line.fill.background()
                        # Explicitly disable shadow
                        card_bg.shadow.inherit = False

                        # Add orange left border inside the card (avoiding rounded corners)
                        if has_border_left:
                            # Create border as rectangle, inset to avoid rounded corners
                            border_height = card_height - 0.24  # Avoid top and bottom rounded corners
                            border_y = card_y + 0.12  # Start below top corner radius

                            border_line = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE,
                                Inches(LayoutConfig.PADDING + 0.02),  # Small inset from card edge
                                Inches(border_y),
                                Inches(0.03),  # 4px width
                                Inches(border_height)
                            )
                            border_line.fill.solid()
                            border_line.fill.fore_color.rgb = ColorConfig.ORANGE
                            border_line.line.fill.background()
                            border_line.shadow.inherit = False

                        # Process div content (h3 + ul) with proper padding
                        div_y = card_y + 0.3  # Top padding inside card

                        # Find h3
                        h3_elem = elem.find('./h3')
                        if h3_elem is not None:
                            h3_text = self.extract_text_content(h3_elem).strip()
                            if h3_text:
                                h3_box = self.add_textbox(
                                    slide,
                                    LayoutConfig.PADDING + 0.35,  # Left padding (after border)
                                    div_y,
                                    self.SLIDE_WIDTH - 2*LayoutConfig.PADDING - 0.5,  # Width minus padding
                                    0.4
                                )
                                h3_tf = h3_box.text_frame
                                h3_p = h3_tf.paragraphs[0]
                                h3_p.alignment = PP_ALIGN.LEFT

                                self.add_formatted_text(h3_p, h3_elem)

                                # h3 styling from inline style
                                h3_style = h3_elem.get('style', '')
                                for run in h3_p.runs:
                                    run.font.name = FontConfig.HEADER_FONT
                                    run.font.size = Pt(16)  # 24px * 0.75
                                    run.font.bold = True
                                    # Parse color from style (default to orange for h3 in cards)
                                    if 'color:' in h3_style:
                                        color_match = re.search(r'color:\s*([^;]+)', h3_style)
                                        if color_match:
                                            color = self.parse_color(color_match.group(1).strip())
                                            if color:
                                                run.font.color.rgb = color
                                            else:
                                                run.font.color.rgb = ColorConfig.ORANGE
                                    else:
                                        run.font.color.rgb = ColorConfig.ORANGE

                                div_y += 0.5

                        # Find ul
                        ul_elem = elem.find('./ul')
                        if ul_elem is not None:
                            li_items = ul_elem.findall('./li')
                            list_height = len(li_items) * 0.35 + 0.2

                            # Create single text box for all list items
                            list_box = self.add_textbox(
                                slide,
                                LayoutConfig.PADDING + 0.4,  # Left padding (after border)
                                div_y,
                                self.SLIDE_WIDTH - 2*LayoutConfig.PADDING - 0.6,
                                list_height
                            )
                            list_tf = list_box.text_frame
                            list_tf.word_wrap = True
                            list_tf.clear()

                            # Detect bullet marker (checkmark)
                            marker = self.detect_list_marker(li_items[0]) if li_items else '✓'

                            # Add each list item as a paragraph
                            first_item = True
                            for li in li_items:
                                if first_item:
                                    li_p = list_tf.paragraphs[0] if len(list_tf.paragraphs) > 0 else list_tf.add_paragraph()
                                    first_item = False
                                else:
                                    li_p = list_tf.add_paragraph()

                                li_p.level = 0
                                li_p.alignment = PP_ALIGN.LEFT
                                li_p.space_after = Pt(6)

                                # Enable bullet formatting with checkmark
                                XMLHelper.add_bullet_with_indent(li_p, marker=marker, font_name=FontConfig.BODY_FONT)

                                # Get pPr for adding custom bullet color
                                pPr = li_p._element.get_or_add_pPr()

                                # Add bullet color (orange for checkmark)
                                buClr_xml = f'''<a:buClr {nsdecls("a")}>
                                    <a:srgbClr val="ed5e29"/>
                                </a:buClr>'''
                                buClr = parse_xml(buClr_xml)
                                pPr.append(buClr)

                                # Add formatted text (skip the positioned span with checkmark)
                                # The text is in the tail of the first span
                                children = list(li)
                                if children and children[0].tag == 'span':
                                    first_span = children[0]
                                    style = first_span.get('style', '')
                                    # Check if this is the absolutely positioned bullet
                                    if 'position: absolute' in style and first_span.text and first_span.text.strip() in ['✓', '•', '○', '■']:
                                        # Skip the span, but get its tail text (the actual content)
                                        if first_span.tail and first_span.tail.strip():
                                            li_p.add_run().text = first_span.tail.strip()
                                        # Add any remaining children
                                        for child in children[1:]:
                                            self.add_formatted_text(li_p, child)
                                    else:
                                        # Not a bullet span, process normally
                                        self.add_formatted_text(li_p, li, skip_leading_marker=True)
                                else:
                                    # No children or different structure
                                    self.add_formatted_text(li_p, li, skip_leading_marker=True)

                                # Apply list item styling
                                li_style = li.get('style', '')
                                for run in li_p.runs:
                                    if not run.font.name:
                                        run.font.name = FontConfig.BODY_FONT
                                    if 'font-size:' in li_style:
                                        size_match = re.search(r'font-size:\s*(\d+)px', li_style)
                                        if size_match:
                                            run.font.size = Pt(int(size_match.group(1)) * 0.75)
                                    elif not run.font.size:
                                        run.font.size = Pt(14)
                                    run.font.color.rgb = ColorConfig.DARK_GRAY

                            div_y += list_height

                        y_pos = div_y + 0.3

            return  # Styled div layout handled, exit method

        # Standard handling (no grid layout, no styled divs)
        # Always create our own text box for content
        content_box = self.add_textbox(slide, self.PADDING, LayoutConfig.CONTENT_START_Y,
                                       self.SLIDE_WIDTH - 2*self.PADDING, 5.5)
        tf = content_box.text_frame
        tf.word_wrap = True

        # Handle bullets and paragraphs
        ul_elems = content_body.findall('.//ul')
        ol_elems = content_body.findall('.//ol')
        p_elems = content_body.findall('.//p')

        # First pass: Identify paragraphs (numbered or heading-style) and their associated lists
        import re
        numbered_para_map = {}  # Maps p_elem index to following ul/ol element
        heading_para_map = {}  # Maps p_elem index to following ul/ol element (non-numbered headings)
        ul_skip_set = set()  # UL elements to skip in main processing

        for idx, p_elem in enumerate(p_elems):
            if p_elem.getparent().tag not in ['li', 'ul', 'ol']:
                p_text = self.extract_text_content(p_elem).strip()
                match = re.match(r'^(\d+)[\.\)]\s+', p_text)

                # Check if followed by a list
                next_sibling = p_elem.getnext()

                if match:
                    # This is a numbered paragraph. Check if followed by a list
                    if next_sibling is not None and next_sibling.tag in ['ul', 'ol']:
                        numbered_para_map[idx] = next_sibling
                        ul_skip_set.add(next_sibling)
                elif next_sibling is not None and next_sibling.tag in ['ul', 'ol']:
                    # This is a heading paragraph (bold, ends with colon, or short) followed by a list
                    style = p_elem.get('style', '')
                    is_heading = (
                        'font-weight: 600' in style or
                        'font-weight: bold' in style or
                        p_text.endswith(':') or
                        (len(p_text) < 60 and not p_text.endswith('.'))  # Short paragraph, not a sentence
                    )
                    if is_heading:
                        heading_para_map[idx] = next_sibling
                        ul_skip_set.add(next_sibling)

        # Track which paragraphs and lists have been processed
        processed_para_indices = set()
        processed_lists = set()

        # Process unordered lists that are NOT associated with paragraphs
        first_item = True
        for ul in ul_elems:
            if ul in ul_skip_set:
                continue  # Skip - will be processed with its paragraph

            # Skip ULs that are children of OLs (nested lists within ordered lists)
            if ul.getparent() is not None and ul.getparent().tag == 'ol':
                continue  # Will be processed as part of OL processing

            processed_lists.add(ul)

            # Check if this should be treated as a numbered list
            is_numbered = self.is_numbered_list(ul)

            for li in ul.findall('./li'):
                if first_item:
                    # Use the first paragraph
                    p = tf.paragraphs[0]
                    first_item = False
                else:
                    p = tf.add_paragraph()

                p.level = 0
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(0)  # No spacing between list items

                if is_numbered:
                    # Use XMLHelper for numbering with indentation
                    XMLHelper.add_numbering_with_indent(p, style='arabicPeriod')

                    # Add the actual content with formatting (skip leading number in original text)
                    self.add_formatted_text(p, li, skip_leading_number=True)
                else:
                    # Detect marker type for bullets
                    marker = self.detect_list_marker(li)

                    # Use XMLHelper for bullet with indentation
                    XMLHelper.add_bullet_with_indent(p, marker=marker, font_name=self.BODY_FONT)

                    # Add the actual content with formatting (skip marker in original text)
                    self.add_formatted_text(p, li, skip_leading_marker=True)

                # Apply font styling to all runs using helper
                FontStyler.apply_body_font_styling(p.runs, is_dark_bg=is_dark)

        # Process ordered lists (numbered) that are NOT associated with paragraphs
        for ol in ol_elems:
            if ol in ul_skip_set:
                continue  # Skip - will be processed with its paragraph

            processed_lists.add(ol)

            # Process all children of OL (both li and nested ul elements)
            item_num = 1
            for child in ol:
                # Handle <li> elements as numbered items
                if child.tag == 'li':
                    if first_item:
                        p = tf.paragraphs[0]
                        first_item = False
                    else:
                        p = tf.add_paragraph()

                    p.level = 0
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(0)  # No spacing between list items

                    # Enable numbering with proper indentation
                    XMLHelper.add_numbering_with_indent(p)

                    # Add the actual content with formatting
                    self.add_formatted_text(p, child)

                    # Apply font styling to all runs
                    for run in p.runs:
                        if not run.font.name:
                            run.font.name = self.BODY_FONT
                        if not run.font.size:
                            run.font.size = Pt(16)
                        # Set colors: orange for bold, appropriate color for regular based on background
                        if run.font.bold:
                            run.font.color.rgb = ColorConfig.ORANGE
                        else:
                            run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

                    item_num += 1

                # Handle nested <ul> elements as sub-bullets
                elif child.tag == 'ul':
                    processed_lists.add(child)  # Mark as processed
                    for nested_li in child.findall('./li'):
                        if first_item:
                            p = tf.paragraphs[0]
                            first_item = False
                        else:
                            p = tf.add_paragraph()

                        p.alignment = PP_ALIGN.LEFT
                        p.space_after = Pt(0)

                        # Detect marker type for bullets
                        marker = self.detect_list_marker(nested_li)
                        XMLHelper.add_bullet_with_indent(
                            p,
                            marker=marker,
                            font_name=self.BODY_FONT,
                            left_margin_inches=0.75,  # Indent for level 1
                            hanging_indent_inches=-0.25
                        )

                        # Set level after add_bullet_with_indent to avoid being overridden
                        p.level = 1

                        # Add the actual content with formatting
                        self.add_formatted_text(p, nested_li, skip_leading_marker=True)

                        # Apply font styling to all runs
                        for run in p.runs:
                            if not run.font.name:
                                run.font.name = self.BODY_FONT
                            if not run.font.size:
                                run.font.size = Pt(16)
                            # Set colors: orange for bold, appropriate color for regular based on background
                            if run.font.bold:
                                run.font.color.rgb = ColorConfig.ORANGE
                            else:
                                run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

        # Process all paragraphs in document order
        # This ensures content appears in the correct sequence
        for idx, p_elem in enumerate(p_elems):
            if p_elem.getparent().tag in ['li', 'ul', 'ol']:
                continue  # Skip paragraphs inside list items

            if idx in processed_para_indices:
                continue  # Already processed

            p_text = self.extract_text_content(p_elem).strip()
            if not p_text:
                continue

            processed_para_indices.add(idx)

            # Determine paragraph type and process accordingly
            style = p_elem.get('style', '')

            # Check if this is a numbered paragraph with an associated list
            match = re.match(r'^(\d+)[\.\)]\s+', p_text)
            if match and idx in numbered_para_map:
                # Numbered paragraph with auto-numbering
                if first_item:
                    p = tf.paragraphs[0]
                    first_item = False
                else:
                    p = tf.add_paragraph()
                p.level = 0
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(0)  # No spacing between list items

                # Enable numbering with proper indentation
                pPr = p._element.get_or_add_pPr()
                pPr.set('marL', '457200')
                pPr.set('indent', '-228600')
                pPr.set('lvl', '0')

                # Add autonumbering
                buAutoNum_xml = f'<a:buAutoNum {nsdecls("a")} type="arabicPeriod"/>'
                buAutoNum = parse_xml(buAutoNum_xml)
                pPr.append(buAutoNum)

                # Add content, stripping the number
                self.add_formatted_text(p, p_elem, skip_leading_number=True)

                # Apply styling
                for run in p.runs:
                    if 'font-family:' in style:
                        if 'var(--font-header)' in style or self.HEADER_FONT.lower() in style.lower():
                            run.font.name = self.HEADER_FONT
                        elif not run.font.name:
                            run.font.name = self.BODY_FONT
                    elif not run.font.name:
                        run.font.name = self.BODY_FONT

                    if 'font-size:' in style:
                        size_match = re.search(r'font-size:\s*(\d+)px', style)
                        if size_match:
                            run.font.size = Pt(int(size_match.group(1)) * 0.75)
                    elif not run.font.size:
                        run.font.size = Pt(16)

                    if 'font-weight: 600' in style or 'font-weight: bold' in style:
                        run.font.bold = True

                    if 'color:' in style:
                        color_match = re.search(r'color:\s*([^;]+)', style)
                        if color_match:
                            color = self.parse_color(color_match.group(1).strip())
                            if color:
                                run.font.color.rgb = color
                            else:
                                # Set colors: orange for bold, appropriate color for regular based on background
                                if run.font.bold:
                                    run.font.color.rgb = ColorConfig.ORANGE
                                else:
                                    run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']
                        else:
                            # Set colors: orange for bold, appropriate color for regular based on background
                            if run.font.bold:
                                run.font.color.rgb = ColorConfig.ORANGE
                            else:
                                run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']
                    else:
                        # Set colors: orange for bold, appropriate color for regular based on background
                        if run.font.bold:
                            run.font.color.rgb = ColorConfig.ORANGE
                        else:
                            run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

                # Process associated list
                list_elem = numbered_para_map[idx]
                for li in list_elem.findall('./li'):
                    li_p = tf.add_paragraph()
                    li_p.level = 0
                    li_p.alignment = PP_ALIGN.LEFT
                    li_p.space_after = Pt(0)  # No spacing between list items

                    marker = self.detect_list_marker(li)
                    XMLHelper.add_bullet_with_indent(
                        li_p,
                        marker=marker,
                        font_name=self.BODY_FONT,
                        left_margin_inches=0.75,  # 685800 EMUs
                        hanging_indent_inches=-0.25
                    )

                    self.add_formatted_text(li_p, li, skip_leading_marker=True)

                    for run in li_p.runs:
                        if not run.font.name:
                            run.font.name = self.BODY_FONT
                        if not run.font.size:
                            run.font.size = Pt(16)
                        # Set colors: orange for bold, appropriate color for regular based on background
                        if run.font.bold:
                            run.font.color.rgb = ColorConfig.ORANGE
                        else:
                            run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

            # Check if this is a heading paragraph with an associated list
            elif idx in heading_para_map:
                # Heading paragraph (bold, ends with colon, etc.)
                if first_item:
                    p = tf.paragraphs[0]
                    first_item = False
                else:
                    p = tf.add_paragraph()
                p.level = 0
                p.alignment = PP_ALIGN.LEFT

                # Detect if this is a subheading (bold, short, ends with colon)
                is_subheading = False
                strong_elem = p_elem.find('.//strong')
                if strong_elem is not None:
                    p_text = self.extract_text_content(p_elem).strip()
                    if len(p_text) < 100 and p_text.endswith(':'):
                        is_subheading = True

                # Apply appropriate spacing
                p.line_spacing = SpacingConfig.LINE_SPACING
                if is_subheading:
                    p.space_before = SpacingConfig.SUBHEADING_SPACE_BEFORE
                    p.space_after = SpacingConfig.SUBHEADING_SPACE_AFTER
                else:
                    p.space_before = SpacingConfig.PARAGRAPH_SPACE_BEFORE
                    p.space_after = SpacingConfig.PARAGRAPH_SPACE_AFTER

                self.add_formatted_text(p, p_elem)

                for run in p.runs:
                    if not run.font.name:
                        run.font.name = self.BODY_FONT

                    if 'font-size:' in style:
                        size_match = re.search(r'font-size:\s*(\d+)px', style)
                        if size_match:
                            run.font.size = Pt(int(size_match.group(1)) * 0.75)
                    elif not run.font.size:
                        run.font.size = Pt(16)

                    if 'font-weight: 600' in style or 'font-weight: bold' in style:
                        run.font.bold = True

                    if 'color:' in style:
                        color_match = re.search(r'color:\s*([^;]+)', style)
                        if color_match:
                            color = self.parse_color(color_match.group(1).strip())
                            if color:
                                run.font.color.rgb = color
                            else:
                                # Set colors: orange for bold, appropriate color for regular based on background
                                if run.font.bold:
                                    run.font.color.rgb = ColorConfig.ORANGE
                                else:
                                    run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']
                    else:
                        # Set colors: orange for bold, appropriate color for regular based on background
                        if run.font.bold:
                            run.font.color.rgb = ColorConfig.ORANGE
                        else:
                            run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

                # Process associated list (both li and nested ul elements)
                list_elem = heading_para_map[idx]
                for child in list_elem:
                    # Handle <li> elements
                    if child.tag == 'li':
                        li_p = tf.add_paragraph()
                        li_p.level = 0
                        li_p.alignment = PP_ALIGN.LEFT
                        li_p.space_after = Pt(0)  # No spacing between list items

                        marker = self.detect_list_marker(child)
                        XMLHelper.add_bullet_with_indent(li_p, marker=marker, font_name=self.BODY_FONT)

                        self.add_formatted_text(li_p, child, skip_leading_marker=True)

                        for run in li_p.runs:
                            if not run.font.name:
                                run.font.name = self.BODY_FONT
                            if not run.font.size:
                                run.font.size = Pt(16)
                            # Set colors: orange for bold, appropriate color for regular based on background
                            if run.font.bold:
                                run.font.color.rgb = ColorConfig.ORANGE
                            else:
                                run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

                    # Handle nested <ul> elements as sub-bullets
                    elif child.tag == 'ul':
                        for nested_li in child.findall('./li'):
                            nested_p = tf.add_paragraph()
                            nested_p.alignment = PP_ALIGN.LEFT
                            nested_p.space_after = Pt(0)

                            marker = self.detect_list_marker(nested_li)
                            XMLHelper.add_bullet_with_indent(
                                nested_p,
                                marker=marker,
                                font_name=self.BODY_FONT,
                                left_margin_inches=0.75,  # Indent for level 1
                                hanging_indent_inches=-0.25
                            )

                            # Set level after add_bullet_with_indent to avoid being overridden
                            nested_p.level = 1

                            self.add_formatted_text(nested_p, nested_li, skip_leading_marker=True)

                            for run in nested_p.runs:
                                if not run.font.name:
                                    run.font.name = self.BODY_FONT
                                if not run.font.size:
                                    run.font.size = Pt(16)
                                if run.font.bold:
                                    run.font.color.rgb = ColorConfig.ORANGE
                                else:
                                    run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

            # Regular paragraph
            else:
                if first_item:
                    p = tf.paragraphs[0]
                    first_item = False
                else:
                    p = tf.add_paragraph()

                if 'text-align: center' in style:
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.alignment = PP_ALIGN.LEFT

                # Detect if this is a subheading paragraph (bold text, typically ending with colon)
                is_subheading = False
                strong_elem = p_elem.find('.//strong')
                if strong_elem is not None:
                    # Get text content to check if entire paragraph is bold and short
                    p_text = self.extract_text_content(p_elem).strip()
                    # Subheading criteria: has strong tag, short text (<100 chars), and ends with colon
                    if len(p_text) < 100 and p_text.endswith(':'):
                        is_subheading = True

                # Apply appropriate spacing
                if is_subheading:
                    p.space_before = SpacingConfig.SUBHEADING_SPACE_BEFORE
                    p.space_after = SpacingConfig.SUBHEADING_SPACE_AFTER
                else:
                    p.space_before = SpacingConfig.PARAGRAPH_SPACE_BEFORE
                    p.space_after = SpacingConfig.PARAGRAPH_SPACE_AFTER

                self.add_formatted_text(p, p_elem)

                for run in p.runs:
                    if 'font-family:' in style:
                        if 'var(--font-header)' in style or self.HEADER_FONT.lower() in style.lower():
                            run.font.name = self.HEADER_FONT
                        elif not run.font.name:
                            run.font.name = self.BODY_FONT
                    elif not run.font.name:
                        run.font.name = self.BODY_FONT

                    if 'font-size:' in style:
                        size_match = re.search(r'font-size:\s*(\d+)px', style)
                        if size_match:
                            run.font.size = Pt(int(size_match.group(1)) * 0.75)
                    elif not run.font.size:
                        run.font.size = Pt(16)

                    if 'font-weight: 600' in style or 'font-weight: bold' in style or 'font-weight:600' in style:
                        run.font.bold = True

                    if 'color:' in style:
                        color_match = re.search(r'color:\s*([^;]+)', style)
                        if color_match:
                            color = self.parse_color(color_match.group(1).strip())
                            if color:
                                run.font.color.rgb = color
                            else:
                                # Set colors: orange for bold, appropriate color for regular based on background
                                if run.font.bold:
                                    run.font.color.rgb = ColorConfig.ORANGE
                                else:
                                    run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']
                        else:
                            # Set colors: orange for bold, appropriate color for regular based on background
                            if run.font.bold:
                                run.font.color.rgb = ColorConfig.ORANGE
                            else:
                                run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']
                    else:
                        # Set colors: orange for bold, appropriate color for regular based on background
                        if run.font.bold:
                            run.font.color.rgb = ColorConfig.ORANGE
                        else:
                            run.font.color.rgb = ColorConfig.CREAM if is_dark else self.COLORS['#475569']

        # Handle citations
        citation_elem = html_slide.find('.//*[@class="citation"]')
        if citation_elem is not None:
            citation_box = self.add_textbox(slide, self.PADDING, self.SLIDE_HEIGHT - 0.7,
                                           self.SLIDE_WIDTH - 2*self.PADDING, 0.4)
            self.add_text(citation_box.text_frame, self.extract_text_content(citation_elem),
                         font_name=self.BODY_FONT, font_size=Pt(13),
                         color=self.COLORS['#64748b'], alignment=PP_ALIGN.LEFT,
                         italic=True)

    def handle_vocab_table_slide(self, slide, html_slide):
        """Handle vocabulary table slide layout"""
        # Add title
        title_elem = html_slide.find('.//h2')
        if title_elem is not None:
            title_box = self.add_textbox(slide, self.PADDING, self.PADDING,
                                        self.SLIDE_WIDTH - 2*self.PADDING, 0.6)
            self.add_text(title_box.text_frame, self.extract_text_content(title_elem),
                         font_name=self.HEADER_FONT, font_size=Pt(28), bold=True,
                         color=self.COLORS['#1e293b'], alignment=PP_ALIGN.LEFT)

        # Find table
        table_elem = html_slide.find('.//table')
        if table_elem is not None:
            rows = table_elem.findall('.//tr')

            # Calculate table dimensions
            num_rows = len(rows)
            num_cols = len(rows[0].findall('.//th')) if len(rows) > 0 else 2

            # Add table shape
            x, y, cx, cy = Inches(self.PADDING), Inches(1.3), Inches(self.SLIDE_WIDTH - 2*self.PADDING), Inches(5.5)
            shape = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy)
            table = shape.table

            # Distribute table width evenly across columns for autofit
            table_width = self.SLIDE_WIDTH - 2*self.PADDING
            col_width = table_width / num_cols
            for col in table.columns:
                col.width = Inches(col_width)

            # Populate table
            for row_idx, row in enumerate(rows):
                cells = row.findall('.//th') + row.findall('.//td')
                for col_idx, cell in enumerate(cells):
                    if col_idx < num_cols:
                        cell_text = self.extract_text_content(cell)
                        table.cell(row_idx, col_idx).text = cell_text

                        # Format header row
                        if row_idx == 0:
                            table.cell(row_idx, col_idx).fill.solid()
                            table.cell(row_idx, col_idx).fill.fore_color.rgb = self.COLORS['#f97316']  # Orange accent color
                            for paragraph in table.cell(row_idx, col_idx).text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = RGBColor(255, 255, 255)
                                    run.font.size = Pt(16)
                                    run.font.name = self.HEADER_FONT
                                    run.font.bold = True
                        else:
                            # Data rows - white background with light gray border
                            table.cell(row_idx, col_idx).fill.solid()
                            table.cell(row_idx, col_idx).fill.fore_color.rgb = RGBColor(255, 255, 255)

                            # Add bottom border
                            tc = table.cell(row_idx, col_idx)._tc
                            tcPr = tc.get_or_add_tcPr()
                            # Add bottom border
                            from pptx.oxml import parse_xml
                            border_xml = '<a:lnB xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="12700"><a:solidFill><a:srgbClr val="E2E8F0"/></a:solidFill></a:lnB>'
                            tcPr.append(parse_xml(border_xml))

                            for paragraph in table.cell(row_idx, col_idx).text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(16)
                                    run.font.name = self.BODY_FONT

    def handle_comparison_slide(self, slide, html_slide):
        """Handle comparison slide layout (two columns)"""
        # Add title
        title_elem = html_slide.find('.//h2')
        if title_elem is not None:
            title_box = self.add_textbox(slide, self.PADDING, self.PADDING,
                                        self.SLIDE_WIDTH - 2*self.PADDING, 0.6)
            p = title_box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            self.add_formatted_text(p, title_elem)
            for run in p.runs:
                run.font.name = self.HEADER_FONT
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = self.COLORS['#1e293b']

        # Find comparison boxes
        comparison_boxes = html_slide.findall('.//*[@class="comparison-box"]')

        # Two-column layout
        col_width = (self.SLIDE_WIDTH - 2*self.PADDING - 0.4) / 2
        y_start = LayoutConfig.CONTENT_START_Y

        for idx, box in enumerate(comparison_boxes[:2]):  # Handle first two boxes
            x_pos = self.PADDING if idx == 0 else self.PADDING + col_width + 0.4

            # Add background box
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_start),
                Inches(col_width), Inches(5.5)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = self.COLORS['#f8fafc']
            bg_shape.line.color.rgb = self.COLORS['#e2e8f0']
            bg_shape.line.width = Pt(2)
            bg_shape.shadow.inherit = False

            # Add heading
            heading = box.find('.//h3')
            if heading is not None:
                heading_box = self.add_textbox(slide, x_pos + 0.1, y_start + 0.2,
                                              col_width - 0.2, 0.5)
                p = heading_box.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                self.add_formatted_text(p, heading)
                for run in p.runs:
                    run.font.name = self.HEADER_FONT
                    run.font.size = Pt(16)
                    run.font.bold = True
                    run.font.color.rgb = self.COLORS['#7373b0']

            # Create single text box for all content in this column
            content_box = self.add_textbox(slide, x_pos + 0.15, y_start + 0.8,
                                          col_width - 0.3, 4.5)
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.clear()

            # Add paragraphs
            paragraphs = box.findall('.//p')
            for p_elem in paragraphs:
                p_text = self.extract_text_content(p_elem)
                if p_text:
                    p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(10)

                    # Add formatted text
                    self.add_formatted_text(p, p_elem)

                    # Apply styling
                    for run in p.runs:
                        if not run.font.name:
                            run.font.name = self.BODY_FONT
                        if not run.font.size:
                            # Check for italic (smaller intro text)
                            if 'italic' in p_elem.get('style', '') or p_elem.tag == 'em':
                                run.font.size = Pt(16)
                                run.font.italic = True
                            else:
                                run.font.size = Pt(16)
                        run.font.color.rgb = self.COLORS['#475569']

            # Add bullets if present
            ul = box.find('.//ul')
            if ul is not None:
                # Check if this should be treated as a numbered list
                is_numbered = self.is_numbered_list(ul)

                for li in ul.findall('.//li'):
                    p = tf.add_paragraph()
                    p.level = 0
                    p.alignment = PP_ALIGN.LEFT

                    # Enable formatting
                    if is_numbered:
                        XMLHelper.add_numbering_with_indent(p)
                        # Add formatted content (skip leading number in original text)
                        self.add_formatted_text(p, li, skip_leading_number=True)
                    else:
                        # Detect marker type and add bullet
                        marker = self.detect_list_marker(li)
                        XMLHelper.add_bullet_with_indent(p, marker=marker, font_name=self.BODY_FONT)
                        # Add formatted content (skip marker in original text)
                        self.add_formatted_text(p, li, skip_leading_marker=True)

                    # Apply styling to all runs
                    for run in p.runs:
                        if not run.font.name:
                            run.font.name = self.BODY_FONT
                        if not run.font.size:
                            run.font.size = Pt(16)
                        run.font.color.rgb = self.COLORS['#475569']

        # Add any additional content below comparison boxes (like Vietnamese proverb)
        extra_paragraphs = html_slide.xpath('./p[not(ancestor::*[@class="comparison-box"])]')
        if extra_paragraphs:
            y_bottom = 7.0
            for p_elem in extra_paragraphs:
                p_text = self.extract_text_content(p_elem)
                if p_text:
                    extra_box = self.add_textbox(slide, self.PADDING, y_bottom,
                                                 self.SLIDE_WIDTH - 2*self.PADDING, 0.4)
                    p = extra_box.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    self.add_formatted_text(p, p_elem)
                    for run in p.runs:
                        run.font.name = self.BODY_FONT
                        run.font.size = Pt(16)
                        if 'font-weight: 600' in p_elem.get('style', ''):
                            run.font.bold = True
                        run.font.color.rgb = self.COLORS['#475569']
                    y_bottom += 0.5

    def handle_activity_slide(self, slide, html_slide):
        """Handle activity/instruction slide layout"""
        # Add activity header with colored background
        header = html_slide.find('.//*[@class="activity-header"]')
        if header is not None:
            # Add purple background
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(self.PADDING), Inches(self.PADDING),
                Inches(self.SLIDE_WIDTH - 2*self.PADDING), Inches(1.2)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = self.COLORS['#7373b0']
            bg_shape.line.width = 0
            bg_shape.shadow.inherit = False

            # Add title text
            title = header.find('.//*[@class="activity-title"]')
            if title is not None:
                title_box = self.add_textbox(slide, self.PADDING + 0.2, self.PADDING + 0.2,
                                            self.SLIDE_WIDTH - 2*self.PADDING - 0.4, 0.8)
                self.add_text(title_box.text_frame, self.extract_text_content(title),
                             font_name=self.HEADER_FONT, font_size=Pt(32), bold=True,
                             color=RGBColor(255, 255, 255), alignment=PP_ALIGN.CENTER)

        # Handle regular content below
        y_pos = 2.0

        # Find content body
        content_body = html_slide.find('.//*[@class="content-body"]')
        if content_body is None:
            content_body = html_slide

        # Get all children from content body
        content_elements = list(content_body)

        for elem in content_elements:
            if elem.tag == 'p':
                p_text = self.extract_text_content(elem)
                if p_text and len(p_text.strip()) > 0:
                    para_box = self.add_textbox(slide, self.PADDING, y_pos,
                                               self.SLIDE_WIDTH - 2*self.PADDING, 0.6)
                    self.add_text(para_box.text_frame, p_text,
                                 font_name=self.BODY_FONT, font_size=Pt(18),
                                 color=self.COLORS['#475569'])
                    y_pos += 0.7

            elif elem.tag == 'div' and 'grid' in elem.get('style', ''):
                # Handle grid layouts (e.g., three-column reflections)
                grid_items = elem.findall('./div')
                if grid_items:
                    num_cols = len(grid_items)
                    col_width = (self.SLIDE_WIDTH - 2*self.PADDING - (num_cols-1)*0.2) / num_cols

                    for idx, item in enumerate(grid_items):
                        x_pos = self.PADDING + idx * (col_width + 0.2)

                        # Add background box
                        bg_shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x_pos), Inches(y_pos),
                            Inches(col_width), Inches(1.5)
                        )
                        bg_shape.fill.solid()
                        bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        bg_shape.line.color.rgb = self.COLORS['#e2e8f0']
                        bg_shape.shadow.inherit = False

                        # Add content from the grid item
                        item_box = self.add_textbox(slide, x_pos + 0.1, y_pos + 0.1,
                                                    col_width - 0.2, 1.3)
                        tf = item_box.text_frame
                        tf.word_wrap = True
                        tf.clear()

                        # Add all paragraphs from this grid item
                        first_p = True
                        for p_elem in item.findall('.//p'):
                            if first_p:
                                p = tf.paragraphs[0]
                                first_p = False
                            else:
                                p = tf.add_paragraph()

                            p.alignment = PP_ALIGN.LEFT
                            p.space_after = Pt(8)
                            self.add_formatted_text(p, p_elem)

                            # Apply styling
                            for run in p.runs:
                                if not run.font.name:
                                    run.font.name = self.BODY_FONT
                                if not run.font.size:
                                    # Check for larger font
                                    style = p_elem.get('style', '')
                                    if 'font-size: 26px' in style:
                                        run.font.size = Pt(16)
                                        run.font.bold = True
                                    elif 'font-size: 22px' in style:
                                        run.font.size = Pt(16)
                                    else:
                                        run.font.size = Pt(16)

                                # Check for color in style
                                if 'color:' in p_elem.get('style', ''):
                                    color = self.parse_color(p_elem.get('style').split('color:')[1].split(';')[0])
                                    if color:
                                        run.font.color.rgb = color
                                    else:
                                        run.font.color.rgb = self.COLORS['#475569']
                                else:
                                    run.font.color.rgb = self.COLORS['#475569']

                    y_pos += 1.7

            elif elem.tag == 'ul':
                # Create ONE text box for all bullets in this list
                li_items = elem.findall('./li')
                if li_items:
                    # Check if this should be treated as a numbered list
                    is_numbered = self.is_numbered_list(elem)

                    # Calculate height needed for all bullets
                    bullet_height = len(li_items) * 0.45 + 0.2
                    bullet_box = self.add_textbox(slide, self.PADDING + 0.3, y_pos,
                                                  self.SLIDE_WIDTH - 2*self.PADDING - 0.3, bullet_height)
                    tf = bullet_box.text_frame
                    tf.word_wrap = True
                    tf.clear()

                    # Add all bullets as paragraphs in single text box
                    first_item = True
                    for li in li_items:
                        if first_item:
                            p = tf.paragraphs[0] if len(tf.paragraphs) > 0 else tf.add_paragraph()
                            first_item = False
                        else:
                            p = tf.add_paragraph()

                        p.level = 0
                        p.alignment = PP_ALIGN.LEFT
                        p.space_after = Pt(6)

                        # Enable formatting
                        if is_numbered:
                            XMLHelper.add_numbering_with_indent(p)
                            # Add formatted text (skip leading number in original text)
                            self.add_formatted_text(p, li, skip_leading_number=True)
                        else:
                            # Detect marker type and add bullet
                            marker = self.detect_list_marker(li)
                            XMLHelper.add_bullet_with_indent(p, marker=marker, font_name=self.BODY_FONT)
                            # Add formatted text (skip marker in original text)
                            self.add_formatted_text(p, li, skip_leading_marker=True)

                        # Apply font styling to all runs
                        for run in p.runs:
                            if not run.font.name:
                                run.font.name = self.BODY_FONT
                            if not run.font.size:
                                run.font.size = Pt(16)
                            run.font.color.rgb = self.COLORS['#475569']

                    y_pos += bullet_height

    def handle_quote_slide(self, slide, html_slide):
        """
        LAYOUT TYPE: quote-slide

        Handle quote slide with title, then quote text and attribution.

        Structure:
            <div class="slide quote-slide">
              <h2 class="slide-title">Title</h2>
              <div class="slide-content">
                <p>Quote text here</p>
                <p>— Attribution</p>
              </div>
            </div>
        """
        # Apply cream background
        self.apply_slide_background(slide, ColorConfig.CREAM)

        # Add title
        title_elem = html_slide.find('.//h2[@class="slide-title"]')
        if title_elem is None:
            title_elem = html_slide.find('.//h2')

        if title_elem is not None:
            title_box = self.add_textbox(
                slide,
                LayoutConfig.PADDING,
                LayoutConfig.TITLE_Y,
                self.SLIDE_WIDTH - 2*LayoutConfig.PADDING,
                LayoutConfig.TITLE_HEIGHT
            )
            p = title_box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            self.add_formatted_text(p, title_elem)

            for run in p.runs:
                run.font.name = FontConfig.HEADER_FONT
                run.font.size = FontConfig.HEADING_MEDIUM
                run.font.bold = True
                run.font.color.rgb = ColorConfig.DARK_GRAY
                # Apply condensed letter spacing
                spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
                run.font._element.set('spc', str(spacing_hundredths))

        # Find quote-content div (new structure) or slide-content (legacy)
        content_div = html_slide.find('.//*[@class="quote-content"]')
        if content_div is None:
            content_div = html_slide.find('.//*[@class="slide-content"]')
        if content_div is None:
            return

        y_pos = LayoutConfig.CONTENT_START_Y

        # Extract quote from blockquote or first paragraph
        quote_elem = content_div.find('.//blockquote[@class="large-quote"]')
        if quote_elem is None:
            # Fallback to first <p> tag (legacy structure)
            paragraphs = content_div.findall('.//p')
            if paragraphs:
                quote_elem = paragraphs[0]

        if quote_elem is not None:
            quote_text = self.extract_text_content(quote_elem).strip().strip('"\'""''')

            quote_box = slide.shapes.add_textbox(
                Inches(self.PADDING),
                Inches(y_pos),
                Inches(self.SLIDE_WIDTH - 2*self.PADDING),
                Inches(3.0)
            )
            quote_frame = quote_box.text_frame
            quote_frame.word_wrap = True
            quote_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            # Set margins (.25cm = ~0.1 inches)
            quote_frame.margin_left = Inches(0.1)
            quote_frame.margin_right = Inches(0.1)
            quote_frame.margin_top = 0
            quote_frame.margin_bottom = 0

            p = quote_frame.paragraphs[0]
            p.text = quote_text
            p.font.name = self.BODY_FONT
            p.font.size = Pt(26)  # Reduced by ~20% from 32
            p.font.italic = True
            p.font.color.rgb = self.COLORS['#1e293b']
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.0  # Single spacing

            y_pos += 3.2

        # Extract attribution from quote-attribution or second paragraph
        attr_elem = content_div.find('.//*[@class="quote-attribution"]')
        if attr_elem is None:
            # Fallback to second <p> tag (legacy structure)
            paragraphs = content_div.findall('.//p')
            if len(paragraphs) > 1:
                attr_elem = paragraphs[1]

        if attr_elem is not None:
            attribution = self.extract_text_content(attr_elem).strip()

            attr_box = slide.shapes.add_textbox(
                Inches(self.PADDING),
                Inches(y_pos),
                Inches(self.SLIDE_WIDTH - 2*self.PADDING),
                Inches(0.6)
            )
            attr_frame = attr_box.text_frame
            attr_frame.word_wrap = True

            p = attr_frame.paragraphs[0]
            self.add_formatted_text(p, attr_elem)

            for run in p.runs:
                if not run.font.name:
                    run.font.name = self.BODY_FONT
                if not run.font.size:
                    run.font.size = Pt(16)
                # Only set color if not already set by formatting
                try:
                    if not run.font.color.rgb:
                        run.font.color.rgb = self.COLORS['#64748b']
                except AttributeError:
                    # Color is None, set default
                    run.font.color.rgb = self.COLORS['#64748b']

            p.alignment = PP_ALIGN.LEFT

    def handle_framework_slide(self, slide, html_slide):
        """
        LAYOUT TYPE: framework-slide

        Handle framework slide with component grid.

        Structure:
            <div class="slide framework-slide">
              <h2 class="framework-title">Framework Name</h2>
              <div class="framework-components">
                <div class="component">
                  <h3>Component Title</h3>
                  <p>Component description</p>
                </div>
              </div>
              <div class="framework-commentary">
                <p><strong>General statement</strong></p>
              </div>
            </div>
        """
        print(f"DEBUG: handle_framework_slide called")

        # Apply cream background
        self.apply_slide_background(slide, ColorConfig.CREAM)

        # Extract framework title
        title_elem = html_slide.find('.//*[@class="framework-title"]')
        if title_elem is None:
            title_elem = html_slide.find('.//h2')

        if title_elem is not None:
            title_text = self.extract_text_content(title_elem)
            # Use consistent title positioning from LayoutConfig
            title_box = slide.shapes.add_textbox(
                Inches(LayoutConfig.PADDING),
                Inches(LayoutConfig.TITLE_Y),
                Inches(self.SLIDE_WIDTH - 2*LayoutConfig.PADDING),
                Inches(LayoutConfig.TITLE_HEIGHT)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.margin_top = 0
            title_frame.margin_bottom = 0
            title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            p = title_frame.paragraphs[0]
            p.text = title_text
            p.line_spacing = 1.0  # Single spacing
            p.font.name = self.HEADER_FONT
            p.font.size = Pt(36)
            p.font.color.rgb = ColorConfig.DARK_GRAY
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
            # Apply condensed letter spacing
            spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
            p.font._element.set('spc', str(spacing_hundredths))

        # Use consistent content start position from LayoutConfig
        y_start = LayoutConfig.CONTENT_START_Y

        # Extract and render framework subtitle if present
        subtitle_elem = html_slide.find('.//*[@class="framework-subtitle"]')
        if subtitle_elem is not None:
            subtitle_text = self.extract_text_content(subtitle_elem)
            subtitle_box = slide.shapes.add_textbox(
                Inches(LayoutConfig.PADDING),
                Inches(y_start),
                Inches(self.SLIDE_WIDTH - 2*LayoutConfig.PADDING),
                Inches(0.6)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            p = subtitle_frame.paragraphs[0]
            p.text = subtitle_text
            p.font.name = self.HEADER_FONT
            p.font.size = Pt(24)  # 32px * 0.75
            p.font.bold = True
            p.font.color.rgb = ColorConfig.DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.3

            # Adjust y_start to account for subtitle
            y_start += 0.7

        # Find framework-components container
        components_container = html_slide.find('.//*[@class="framework-components"]')
        if components_container is None:
            return

        # Find all component divs
        components = components_container.findall('./div[@class="component"]')
        if not components:
            return

        # Layout components in grid
        num_components = len(components)
        if num_components <= 3:
            cols = num_components
        elif num_components == 4:
            cols = 2
        else:
            cols = 3

        box_width = (self.SLIDE_WIDTH - 2*self.PADDING - (cols-1)*0.25) / cols
        box_height = 1.3

        for idx, component in enumerate(components):
            row = idx // cols
            col = idx % cols

            x = self.PADDING + col * (box_width + 0.25)
            y = y_start + row * (box_height + 0.25)

            # Extract component heading
            heading_elem = component.find('.//h3')
            if heading_elem is None:
                continue

            heading_text = self.extract_text_content(heading_elem)

            # Extract component parts: question, bullets, description
            question_elem = component.find('.//*[@class="component-question"]')
            bullets_elem = component.find('.//*[@class="component-bullets"]')

            print(f"DEBUG EXTRACTION for '{heading_text}':")
            print(f"  question_elem: {question_elem is not None}")
            print(f"  bullets_elem: {bullets_elem is not None}")
            if bullets_elem is not None:
                print(f"  bullets_elem tag: {bullets_elem.tag}")
                print(f"  bullets_elem class: {bullets_elem.get('class')}")

            # Get all paragraph descriptions (excluding question)
            desc_paragraphs = [p for p in component.findall('.//p')
                             if 'component-question' not in p.get('class', '')]

            # Create white component box
            comp_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(box_width), Inches(box_height)
            )
            comp_shape.fill.solid()
            comp_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            comp_shape.line.color.rgb = self.COLORS['#cbd5e1']  # Light gray border
            comp_shape.line.width = Pt(1)
            comp_shape.shadow.inherit = False

            # Add text
            text_frame = comp_shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.15)
            text_frame.margin_right = Inches(0.15)
            text_frame.margin_top = Inches(0.12)
            text_frame.margin_bottom = Inches(0.12)
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            # Enable autofit - resize shape to fit text
            from pptx.enum.text import MSO_AUTO_SIZE
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            # Add heading
            p = text_frame.paragraphs[0]
            p.text = heading_text
            p.font.name = self.HEADER_FONT
            p.font.size = Pt(20)  # 27px * 0.75
            p.font.bold = True
            p.font.color.rgb = ColorConfig.DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(6)

            # Add question if present
            if question_elem is not None:
                question_text = self.extract_text_content(question_elem)
                p_question = text_frame.add_paragraph()
                p_question.text = question_text
                p_question.font.name = self.BODY_FONT
                p_question.font.size = Pt(12)
                p_question.font.italic = True
                p_question.font.color.rgb = ColorConfig.DARK_GRAY
                p_question.alignment = PP_ALIGN.CENTER
                p_question.space_after = Pt(4)

            # Add bullets if present
            if bullets_elem is not None:
                bullet_items = bullets_elem.findall('.//li')
                print(f"DEBUG: Found {len(bullet_items)} bullets for component '{heading_text}'")
                for li in bullet_items:
                    bullet_text = self.extract_text_content(li)
                    print(f"DEBUG: Adding bullet: {bullet_text}")
                    if bullet_text:
                        p_bullet = text_frame.add_paragraph()
                        p_bullet.text = '• ' + bullet_text
                        p_bullet.font.name = self.BODY_FONT
                        p_bullet.font.size = Pt(11)
                        p_bullet.font.color.rgb = self.COLORS['#64748b']
                        p_bullet.alignment = PP_ALIGN.LEFT
                        p_bullet.level = 0
            else:
                print(f"DEBUG: No bullets found for component '{heading_text}'")

            # Add description paragraphs
            for desc_p in desc_paragraphs:
                desc_text = self.extract_text_content(desc_p)
                if desc_text:
                    p_desc = text_frame.add_paragraph()
                    p_desc.text = desc_text
                    p_desc.font.name = self.BODY_FONT
                    p_desc.font.size = Pt(13)
                    p_desc.font.color.rgb = self.COLORS['#64748b']
                    p_desc.alignment = PP_ALIGN.CENTER
                    p_desc.line_spacing = 1.3

        # Add auxiliary content if present (formerly commentary)
        commentary_div = html_slide.find('.//*[@class="framework-auxiliary"]')
        if commentary_div is None:
            # Fallback to old class name for backward compatibility
            commentary_div = html_slide.find('.//*[@class="framework-commentary"]')
        if commentary_div is not None:
            y_commentary = y_start + ((num_components + cols - 1) // cols) * (box_height + 0.25) + 0.2

            commentary_box = slide.shapes.add_textbox(
                Inches(self.PADDING + 1.0),
                Inches(y_commentary),
                Inches(self.SLIDE_WIDTH - 2*self.PADDING - 2.0),
                Inches(1.0)
            )
            commentary_frame = commentary_box.text_frame
            commentary_frame.word_wrap = True
            commentary_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            # Get all paragraphs in commentary
            for p_elem in commentary_div.findall('.//p'):
                p = commentary_frame.add_paragraph() if len(commentary_frame.paragraphs) > 1 else commentary_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                self.add_formatted_text(p, p_elem)

                # Apply styling
                for run in p.runs:
                    if not run.font.name:
                        run.font.name = self.BODY_FONT
                    if not run.font.size:
                        run.font.size = Pt(16)  # 24px * 0.75
                    # Bold text should be orange
                    if run.font.bold:
                        run.font.color.rgb = ColorConfig.ORANGE
                    else:
                        # Only set color if not already set by formatting
                        try:
                            if not run.font.color.rgb:
                                run.font.color.rgb = self.COLORS['#64748b']
                        except AttributeError:
                            # Color is None, set default
                            run.font.color.rgb = self.COLORS['#64748b']

    def handle_comparison_table_slide(self, slide, html_slide):
        """
        LAYOUT TYPE: comparison-table-slide

        Handle 2-column comparison table layout with actual PowerPoint table.
        Creates a formatted table with orange header row and light grey borders between rows.

        Usage in HTML:
            <div class="slide comparison-table-slide">
              <h2 class="slide-title">Table Title</h2>
              <table>
                <thead>
                  <tr>
                    <th>Column 1</th>
                    <th>Column 2</th>
                  </tr>
                </thead>
                <tbody>
                  <tr>
                    <td>Cell content</td>
                    <td>Cell content</td>
                  </tr>
                  <!-- More rows -->
                </tbody>
              </table>
            </div>
        """
        # Apply cream background
        self.apply_slide_background(slide, ColorConfig.CREAM)

        # Extract slide title
        title_elem = html_slide.find('.//*[@class="slide-title"]')
        if title_elem is None:
            title_elem = html_slide.find('.//h2')

        y_start = self.PADDING
        if title_elem is not None:
            title_text = ''.join(title_elem.itertext()).strip()
            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(self.PADDING),
                Inches(self.PADDING),
                Inches(self.SLIDE_WIDTH - 2*self.PADDING),
                Inches(0.7)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.margin_top = 0
            title_frame.margin_bottom = 0
            p = title_frame.paragraphs[0]
            p.text = title_text
            p.line_spacing = 1.0  # Single spacing
            p.font.name = self.HEADER_FONT
            p.font.size = Pt(36)
            p.font.color.rgb = self.COLORS['#1e293b']
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
            # Apply condensed letter spacing
            spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
            p.font._element.set('spc', str(spacing_hundredths))
            y_start = self.PADDING + 0.9  # After title

        # Find slide-content container
        content_div = html_slide.find('.//*[@class="slide-content"]')
        if content_div is None:
            content_div = html_slide

        # Find table element
        table_elem = content_div.find('.//table')

        if table_elem is None:
            return

        # Extract paragraphs before and after table (in document order)
        paras_before = []
        paras_after = []
        found_table = False

        for child in content_div:
            if child.tag == 'table':
                found_table = True
            elif child.tag == 'p':
                if not found_table:
                    paras_before.append(child)
                else:
                    paras_after.append(child)

        # Render paragraphs before table
        if paras_before:
            for para_elem in paras_before:
                para_text = self.extract_text_content(para_elem).strip()
                if para_text:
                    para_box = self.add_textbox(slide, self.PADDING, y_start,
                                               self.SLIDE_WIDTH - 2*self.PADDING, 0.4)
                    p_tf = para_box.text_frame
                    p_tf.word_wrap = True
                    p = p_tf.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT

                    self.add_formatted_text(p, para_elem)

                    for run in p.runs:
                        if not run.font.name:
                            run.font.name = self.BODY_FONT
                        if not run.font.size:
                            run.font.size = Pt(16)
                        # Bold text in orange
                        if run.font.bold:
                            run.font.color.rgb = ColorConfig.ORANGE
                        else:
                            run.font.color.rgb = self.COLORS['#475569']

                    y_start += 0.45

        # Extract table headers (column titles)
        headers = []
        thead = table_elem.find('.//thead')
        if thead is not None:
            header_cells = thead.findall('.//th')
            headers = [''.join(cell.itertext()).strip() for cell in header_cells]

        # Extract table rows
        rows_data = []
        tbody = table_elem.find('.//tbody')
        if tbody is not None:
            rows = tbody.findall('.//tr')
        else:
            rows = table_elem.findall('.//tr')
            # Skip header row if no thead
            if not headers and rows:
                first_row = rows[0]
                header_cells = first_row.findall('.//th')
                if header_cells:
                    headers = [''.join(cell.itertext()).strip() for cell in header_cells]
                    rows = rows[1:]

        for row in rows:
            cells = row.findall('.//td')
            if not cells:
                cells = row.findall('.//th')
            row_data = [''.join(cell.itertext()).strip() for cell in cells]
            if row_data:
                rows_data.append(row_data)

        if not rows_data:
            return

        # Determine number of columns
        num_cols = len(headers) if headers else len(rows_data[0]) if rows_data else 0

        if num_cols >= 2:
            # Create actual PowerPoint table
            num_rows = len(rows_data) + (1 if headers else 0)

            # Calculate table dimensions
            table_width = self.SLIDE_WIDTH - 2*self.PADDING
            table_height = min(5.5, 0.5 + (len(rows_data) * 0.5))  # Header + rows

            # Add table shape
            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                Inches(self.PADDING),
                Inches(y_start),
                Inches(table_width),
                Inches(table_height)
            ).table

            # Distribute table width evenly across columns
            col_width = table_width / num_cols
            for col in table_shape.columns:
                col.width = Inches(col_width)

            # Format header row if headers exist
            if headers:
                for col_idx, header_text in enumerate(headers):
                    cell = table_shape.cell(0, col_idx)

                    # Orange background with white text
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLORS['#ed5e29']  # Orange accent

                    # Remove all borders (python-pptx doesn't support cell borders directly)
                    # Borders are controlled at table level, set to transparent

                    # Set text
                    text_frame = cell.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = header_text
                    p.font.name = self.HEADER_FONT
                    p.font.size = Pt(20)
                    p.font.color.rgb = self.COLORS['#ffffff']  # White text
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                    # Center vertically
                    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    text_frame.margin_top = Inches(0.1)
                    text_frame.margin_bottom = Inches(0.1)

            # Format data rows
            row_offset = 1 if headers else 0
            for row_idx, row_data in enumerate(rows_data):
                for col_idx, cell_text in enumerate(row_data):
                    cell = table_shape.cell(row_idx + row_offset, col_idx)

                    # White background with grey text
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLORS['#ffffff']  # White

                    # Set text
                    text_frame = cell.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = cell_text
                    p.font.name = self.BODY_FONT
                    p.font.size = Pt(14)
                    p.font.color.rgb = self.COLORS['#475569']  # Grey text
                    p.alignment = PP_ALIGN.LEFT

                    # Set margins
                    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                    text_frame.margin_left = Inches(0.15)
                    text_frame.margin_right = Inches(0.15)
                    text_frame.margin_top = Inches(0.1)
                    text_frame.margin_bottom = Inches(0.1)

            # Set table-level formatting for borders using XML
            # Remove all borders except horizontal lines between data rows
            from pptx.oxml import parse_xml

            tbl = table_shape._tbl
            tblPr = tbl.tblPr

            # Set table borders to none by default
            if tblPr is not None:
                # Try to set first row as special (header style)
                try:
                    if tblPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}firstRow') is None:
                        firstRow = parse_xml('<a:firstRow xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="1"/>')
                        tblPr.append(firstRow)
                except:
                    pass

            # Set individual cell borders via XML
            for row_idx in range(num_rows):
                for col_idx in range(num_cols):
                    tc = tbl.tr_lst[row_idx].tc_lst[col_idx]
                    tcPr = tc.get_or_add_tcPr()

                    # Remove left, right, top borders
                    for border_name in ['lnL', 'lnR', 'lnT']:
                        ln = tcPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}%s' % border_name)
                        if ln is not None:
                            tcPr.remove(ln)
                        # Add explicit "no line"
                        no_line = parse_xml('<a:%s xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:noFill/></a:%s>' % (border_name, border_name))
                        tcPr.append(no_line)

                    # Set bottom border
                    ln_b = tcPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lnB')
                    if ln_b is not None:
                        tcPr.remove(ln_b)

                    # Add light grey bottom border for data rows (not last row)
                    if row_idx < num_rows - 1:
                        bottom_border = parse_xml(
                            '<a:lnB xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="9525">'
                            '<a:solidFill><a:srgbClr val="E2E8F0"/></a:solidFill>'
                            '</a:lnB>'
                        )
                        tcPr.append(bottom_border)
                    else:
                        # No border on last row
                        no_line = parse_xml('<a:lnB xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:noFill/></a:lnB>')
                        tcPr.append(no_line)

            # Render paragraphs after table
            if paras_after:
                y_after = y_start + table_height + 0.2
                for para_elem in paras_after:
                    para_text = self.extract_text_content(para_elem).strip()
                    if para_text:
                        para_box = self.add_textbox(slide, self.PADDING, y_after,
                                                   self.SLIDE_WIDTH - 2*self.PADDING, 0.5)
                        p_tf = para_box.text_frame
                        p_tf.word_wrap = True
                        p = p_tf.paragraphs[0]
                        p.alignment = PP_ALIGN.LEFT

                        self.add_formatted_text(p, para_elem)

                        for run in p.runs:
                            if not run.font.name:
                                run.font.name = self.BODY_FONT
                            if not run.font.size:
                                run.font.size = Pt(16)
                            # Bold text in orange
                            if run.font.bold:
                                run.font.color.rgb = ColorConfig.ORANGE
                            else:
                                run.font.color.rgb = self.COLORS['#475569']

                        y_after += 0.5
        else:
            # Fall back to standard table handling for non-2-column tables
            self.handle_content_slide(slide, html_slide)

    def handle_card_layout(self, slide, cards_container, y_start=2.2):
        """Handle card layout with colored headers and white backgrounds using CSS grid"""
        cards = cards_container.findall('.//div[@class="card"]') if cards_container is not None else []
        if not cards:
            return y_start

        # Check if container uses grid layout (from CSS classes or inline styles)
        computed = self.css_parser.get_computed_style(cards_container)
        is_grid = computed.get('display') == 'grid'

        if is_grid:
            # Use CSS grid layout engine
            available_height = self.SLIDE_HEIGHT - y_start - 1.0  # Leave space for footer
            card_height = 4.5  # Default card height

            positions = self.layout_grid_container(
                cards_container,
                x_start=self.PADDING,
                y_start=y_start,
                width=self.SLIDE_WIDTH - 2*self.PADDING,
                height=available_height,
                row_height=card_height
            )

            # Create cards at calculated positions
            for pos in positions:
                if pos['element'].get('class', '') == 'card':
                    self._create_card(slide, pos['element'],
                                    pos['x'], pos['y'], pos['width'], pos['height'])

            return y_start + card_height
        else:
            # Stacked cards (vertical layout)
            for card in cards:
                y_start = self._create_card(slide, card, self.PADDING, y_start,
                                           self.SLIDE_WIDTH - 2*self.PADDING, 2.5)
                y_start += 0.3  # Gap between cards

            return y_start

    def _create_card(self, slide, card, x, y, width, height):
        """Create a single card with header and body using CSS styling"""
        # Get computed styles (CSS rules + inline styles)
        card_styles = self.css_parser.get_computed_style(card)

        # Get background color (default white)
        bg_color = self.css_parser.parse_color(
            card_styles.get('background', card_styles.get('background-color', 'var(--color-white)'))
        ) or self.COLORS['#ffffff']

        # Get border
        border_value = card_styles.get('border', None)
        border_params = self.css_parser.parse_border(border_value) if border_value else None

        # Get box shadow
        shadow_value = card_styles.get('box-shadow', None)
        shadow_params = self.css_parser.parse_box_shadow(shadow_value) if shadow_value else None

        # Card background
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = bg_color

        # Apply border
        self.apply_border(card_bg, border_params)

        # Explicitly disable shadow
        card_bg.shadow.inherit = False

        # Apply shadow (disabled - no shadows on cards)
        # self.apply_shadow(card_bg, shadow_params)

        # Card header - support both wrapped and unwrapped h3
        header = card.find('.//div[@class="card-header"]')
        header_h3 = None

        if header is not None:
            # Wrapped structure: <div class="card-header"><h3>...</h3></div>
            header_h3 = header.find('.//h3')
            header_style = header.get('style', '')
        else:
            # Unwrapped structure: <h3> directly in card
            header_h3 = card.find('./h3')  # Direct child only
            header_style = ''

        if header_h3 is not None:
            # Parse header styles
            header_styles = self.css_parser.parse_inline_style(header_style)

            # Get header background color (default orange)
            header_bg_color = self.css_parser.parse_color(
                header_styles.get('background', header_styles.get('background-color', 'var(--color-accent)'))
            ) or self.COLORS['#ed5e29']

            # Get header padding
            header_padding = self.css_parser.parse_spacing(header_styles.get('padding', '16px 24px'))
            if isinstance(header_padding, dict):
                header_padding_h = header_padding['left']
                header_padding_v = header_padding['top']
            else:
                header_padding_h = header_padding or 0.25
                header_padding_v = header_padding or 0.16

            if True:  # Keep indentation for existing code
                # Parse h3 styles
                h3_style = header_h3.get('style', '')
                h3_styles = self.css_parser.parse_inline_style(h3_style)

                # Header bar
                header_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x), Inches(y), Inches(width), Inches(0.5)
                )
                header_bar.fill.solid()
                header_bar.fill.fore_color.rgb = header_bg_color
                header_bar.line.fill.background()
                header_bar.shadow.inherit = False

                # Header text
                header_text = self.extract_text_content(header_h3)
                tf = header_bar.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT

                run = p.add_run()
                run.text = header_text

                # Apply CSS font properties (h3 should use header font)
                font_family = h3_styles.get('font-family', 'var(--font-header)')
                # Resolve CSS variable if needed
                if font_family.startswith('var(--'):
                    var_name = font_family.replace('var(', '').replace(')', '').strip()
                    font_family = ColorConfig.CSS_VARIABLES.get(var_name, font_family)
                # Use Cal Sans for header font
                if font_family == 'var(--font-header)' or 'Cal Sans' in font_family:
                    run.font.name = self.HEADER_FONT
                else:
                    run.font.name = font_family

                font_size = h3_styles.get('font-size', '22px')
                run.font.size = self.css_parser.parse_size(font_size, context='font') or Pt(16)

                font_weight = h3_styles.get('font-weight', '700')
                run.font.bold = font_weight in ['bold', '600', '700', '800', '900']

                # Apply condensed letter spacing to all headings via XML (1/100th of a point)
                spacing_hundredths = int(FontConfig.TITLE_LETTER_SPACING.pt * 100)
                run.font._element.set('spc', str(spacing_hundredths))

                # Header text color (default white)
                text_color = self.css_parser.parse_color(
                    h3_styles.get('color', 'var(--color-white)')
                ) or self.COLORS['#ffffff']
                run.font.color.rgb = text_color

                tf.margin_left = Inches(header_padding_h)
                tf.margin_top = Inches(header_padding_v)
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Card body - support both wrapped and unwrapped content
        body = card.find('.//div[@class="card-body"]')

        # Determine body start position based on whether header exists
        if header_h3 is not None:
            body_y = y + 0.6  # After header bar
            body_height = height - 0.7
        else:
            body_y = y + 0.25  # No header, start from top with padding
            body_height = height - 0.5

        body_box = self.add_textbox(slide, x + 0.25, body_y, width - 0.5, body_height)
        tf = body_box.text_frame
        tf.word_wrap = True
        tf.clear()

        # Get body content
        if body is not None:
            # Wrapped structure: <div class="card-body">
            body_content = body
        else:
            # Unwrapped structure: content is directly in card (skip h3)
            body_content = card

        # Process body content (usually p and/or ul)
        # Process any p elements first
        first_element = True
        for p_elem in body_content.findall('./p'):
            if first_element:
                p = tf.paragraphs[0] if len(tf.paragraphs) > 0 else tf.add_paragraph()
                first_element = False
            else:
                p = tf.add_paragraph()

            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)

            self.add_formatted_text(p, p_elem)

            for run in p.runs:
                if not run.font.name:
                    run.font.name = self.BODY_FONT
                if not run.font.size:
                    run.font.size = Pt(14)
                # Check if color needs to be set
                try:
                    _ = run.font.color.rgb
                except AttributeError:
                    # Set colors: orange for bold, gray for regular
                    if run.font.bold:
                        run.font.color.rgb = ColorConfig.ORANGE
                    else:
                        run.font.color.rgb = self.COLORS['#131313']

        # Then process ul if present
        ul = body_content.find('.//ul')
        if ul is not None:
            li_items = ul.findall('./li')

            for li in li_items:
                if first_element:
                    p = tf.paragraphs[0] if len(tf.paragraphs) > 0 else tf.add_paragraph()
                    first_element = False
                else:
                    p = tf.add_paragraph()

                p.level = 0
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(4)

                # Enable bullet formatting with custom margins
                XMLHelper.add_bullet_with_indent(
                    p,
                    marker='•',
                    font_name=self.BODY_FONT,
                    left_margin_inches=0.25,  # 228600 EMUs
                    hanging_indent_inches=-0.125  # -114300 EMUs
                )

                # Add orange bullet color
                pPr = p._element.get_or_add_pPr()
                buClr_xml = f'<a:buClr {nsdecls("a")}><a:srgbClr val="ed5e29"/></a:buClr>'
                buClr = parse_xml(buClr_xml)
                pPr.append(buClr)

                # Add text content
                self.add_formatted_text(p, li)

                # Format text
                for run in p.runs:
                    if not run.font.name:
                        run.font.name = self.BODY_FONT
                    if not run.font.size:
                        run.font.size = Pt(13)
                    # Check if color needs to be set
                    try:
                        _ = run.font.color.rgb
                    except AttributeError:
                        # Set colors: orange for bold, gray for regular
                        if run.font.bold:
                            run.font.color.rgb = ColorConfig.ORANGE
                        else:
                            run.font.color.rgb = self.COLORS['#131313']

        return y + height

    def _create_grid_item(self, slide, grid_item_elem, x, y, width, height):
        """Create a styled grid item container with icon, heading, and description"""
        # Get computed styles for grid-item
        item_styles = self.css_parser.get_computed_style(grid_item_elem)

        # Get background color (default white)
        bg_color = self.css_parser.parse_color(
            item_styles.get('background', item_styles.get('background-color', 'var(--color-white)'))
        ) or self.COLORS['#ffffff']

        # Create background container with rounded corners
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        container.fill.solid()
        container.fill.fore_color.rgb = bg_color

        # No border (invisible on cream background)
        container.line.fill.background()
        container.shadow.inherit = False

        # Add rounded corners by adjusting the shape
        # In python-pptx, ROUNDED_RECTANGLE has an adjustment for corner radius
        try:
            # Set corner radius (12px = 0.125 inches, scaled to shape adjustment)
            # Adjustment value is a fraction of the smaller dimension
            # Typical range: 0-50000 (percentage * 1000)
            container.adjustments[0] = 0.05  # 5% corner radius
        except:
            pass  # If adjustments not available, continue without rounded corners

        # Add orange accent bar at top (grid-item::before)
        accent_height = 0.04  # 4px converted to inches
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(width), Inches(accent_height)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.COLORS['#ed5e29']  # Orange
        accent_bar.line.fill.background()
        accent_bar.shadow.inherit = False

        # Start content below accent bar with padding
        padding = 0.25  # 24px converted to inches
        content_y = y + accent_height + padding
        content_x = x + padding
        content_width = width - 2 * padding

        # Check for icon-container (render placeholder)
        icon_container = grid_item_elem.find('.//*[@class="icon-container"]')
        if icon_container is not None:
            # Add icon placeholder (simple orange circle/shape)
            icon_size = 0.6  # 60px
            icon_x = x + (width - icon_size) / 2  # Center horizontally
            icon_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(icon_x), Inches(content_y), Inches(icon_size), Inches(icon_size)
            )
            icon_circle.fill.solid()
            icon_circle.fill.fore_color.rgb = RGBColor(237, 94, 41)  # Orange
            icon_circle.fill.transparency = 0.2  # 20% opacity for placeholder
            icon_circle.line.fill.background()
            icon_circle.shadow.inherit = False

            content_y += icon_size + 0.16  # Icon + margin-bottom

        # Extract and render heading (h3)
        heading_elem = grid_item_elem.find('.//h3')
        if heading_elem is not None:
            heading_text = self.extract_text_content(heading_elem).strip()
            if heading_text:
                heading_height = 0.4
                heading_box = self.add_textbox(slide, content_x, content_y,
                                              content_width, heading_height)
                h_tf = heading_box.text_frame
                h_tf.word_wrap = True
                h_p = h_tf.paragraphs[0]
                h_p.alignment = PP_ALIGN.CENTER  # Center text

                self.add_formatted_text(h_p, heading_elem)

                # Apply heading styles from CSS
                heading_styles = self.css_parser.get_computed_style(heading_elem)
                for run in h_p.runs:
                    run.font.name = self.BODY_FONT

                    # Font size from CSS (default 22px)
                    font_size = heading_styles.get('font-size', '22px')
                    run.font.size = self.css_parser.parse_size(font_size, context='font') or Pt(16)

                    run.font.bold = True

                    # Color from CSS (default primary color)
                    color_value = heading_styles.get('color', 'var(--color-primary)')
                    color = self.css_parser.parse_color(color_value)
                    if color:
                        run.font.color.rgb = color

                content_y += heading_height + 0.12  # Heading + margin-bottom

        # Extract and render paragraph description
        para_elem = None
        for child in grid_item_elem:
            if child.tag == 'p':
                para_elem = child
                break

        if para_elem is not None:
            para_text = self.extract_text_content(para_elem).strip()
            if para_text:
                para_height = 0.8
                para_box = self.add_textbox(slide, content_x, content_y,
                                           content_width, para_height)
                p_tf = para_box.text_frame
                p_tf.word_wrap = True
                p = p_tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER  # Center text

                self.add_formatted_text(p, para_elem)

                # Apply paragraph styles
                for run in p.runs:
                    if not run.font.name:
                        run.font.name = self.BODY_FONT
                    if not run.font.size:
                        run.font.size = Pt(13)

                    # Default text color
                    run.font.color.rgb = self.COLORS['#475569']

    def handle_stats_banner(self, slide, stats_container, y_start=2.2, is_dark_bg=False):
        """Handle stats banner with large numbers and descriptions

        Args:
            slide: PowerPoint slide object
            stats_container: HTML element containing stat items
            y_start: Y position to start rendering stats
            is_dark_bg: Whether slide has dark background (affects text color)
        """
        stat_items = stats_container.findall('.//div[@class="stat-item"]')
        if not stat_items:
            return y_start

        num_stats = len(stat_items)
        stat_width = (self.SLIDE_WIDTH - 2*self.PADDING) / num_stats

        for idx, stat_item in enumerate(stat_items):
            x_pos = self.PADDING + idx * stat_width

            # Extract stat parts
            stat_number = stat_item.find('.//div[@class="stat-number"]')
            stat_label = stat_item.find('.//div[@class="stat-label"]')
            stat_desc = stat_item.find('.//div[@class="stat-description"]')

            current_y = y_start

            # Add icon above stat (small geometric shape)
            icon_size = 0.5
            icon_x = x_pos + (stat_width - icon_size) / 2
            icon_y = current_y - 0.65
            # Vary icon types: circle, square, circle pattern
            icon_types = ['circle', 'square', 'circle']
            icon_type = icon_types[idx % len(icon_types)]
            self.add_icon_above_stat(slide, icon_x, icon_y, icon_type, icon_size)

            # Add number (large, orange)
            if stat_number is not None:
                num_text = self.extract_text_content(stat_number)
                num_box = self.add_textbox(slide, x_pos, current_y, stat_width, 1.2)
                num_tf = num_box.text_frame
                num_p = num_tf.paragraphs[0]
                num_p.alignment = PP_ALIGN.CENTER

                run = num_p.add_run()
                run.text = num_text
                run.font.name = self.HEADER_FONT
                run.font.size = Pt(48)
                run.font.bold = True
                run.font.color.rgb = self.COLORS['#ed5e29']

                current_y += 1.2

            # Add label (bold, color depends on background)
            if stat_label is not None:
                label_text = self.extract_text_content(stat_label)
                label_box = self.add_textbox(slide, x_pos, current_y, stat_width, 0.4)
                label_tf = label_box.text_frame
                label_p = label_tf.paragraphs[0]
                label_p.alignment = PP_ALIGN.CENTER

                run = label_p.add_run()
                run.text = label_text
                run.font.name = self.BODY_FONT
                run.font.size = Pt(14)
                run.font.bold = True
                # Use cream text on dark backgrounds, dark text on light backgrounds
                run.font.color.rgb = ColorConfig.CREAM if is_dark_bg else ColorConfig.DARK_GRAY

                current_y += 0.4

            # Add description
            if stat_desc is not None:
                desc_text = self.extract_text_content(stat_desc)
                desc_box = self.add_textbox(slide, x_pos + 0.1, current_y, stat_width - 0.2, 0.8)
                desc_tf = desc_box.text_frame
                desc_tf.word_wrap = True
                desc_p = desc_tf.paragraphs[0]
                desc_p.alignment = PP_ALIGN.CENTER

                run = desc_p.add_run()
                run.text = desc_text
                run.font.name = self.BODY_FONT
                run.font.size = Pt(11)
                # Use cream text on dark backgrounds, slate gray on light backgrounds
                run.font.color.rgb = ColorConfig.CREAM if is_dark_bg else self.COLORS['#64748b']

        return y_start + 3.5

    def extract_and_add_speaker_notes(self, slide, html_slide):
        """
        Extract speaker notes from HTML slide and add them to PPTX slide.

        Args:
            slide: python-pptx slide object
            html_slide: lxml HTML element for the slide

        Notes are extracted from <aside class="speaker-notes"> elements in the HTML.
        """
        # Find speaker notes element
        notes_element = html_slide.find('.//aside[@class="speaker-notes"]')

        if notes_element is None:
            # No speaker notes for this slide
            return

        # Extract text content from notes element
        notes_text = self.extract_text_content_recursive(notes_element)

        if not notes_text or not notes_text.strip():
            # Empty notes, skip
            return

        # Clean up notes text:
        # 1. Remove "Speaker Notes" header and following empty line
        notes_text = notes_text.strip()
        if notes_text.startswith("Speaker Notes"):
            # Remove first line (header) and any following empty lines
            lines = notes_text.split('\n')
            lines = lines[1:]  # Remove first line
            notes_text = '\n'.join(lines).strip()

        # 2. Remove consecutive empty lines (replace multiple newlines with single)
        notes_text = re.sub(r'\n\s*\n+', '\n\n', notes_text)

        if not notes_text:
            # After cleaning, notes are empty
            return

        # Add notes to PPTX slide
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame

            # Clear any existing placeholder text
            text_frame.clear()

            # Add the cleaned notes text
            text_frame.text = notes_text

            print(f"  ✓ Added speaker notes ({len(notes_text)} characters)")

        except Exception as e:
            print(f"  ⚠ Warning: Could not add speaker notes: {e}")

    def extract_text_content_recursive(self, element):
        """
        Recursively extract all text content from an HTML element, preserving structure.

        Args:
            element: lxml HTML element

        Returns:
            String with text content, including line breaks for structure
        """
        if element is None:
            return ""

        result = []

        # Get element's direct text
        if element.text:
            # Only strip spaces and tabs, preserve newlines
            result.append(element.text.strip(' \t'))

        # Process children
        for child in element:
            # Handle different HTML elements
            try:
                tag = child.tag.lower() if isinstance(child.tag, str) else ''
            except (AttributeError, TypeError):
                tag = ''

            if tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                # Headers - add extra line break before and after
                result.append('\n' + self.extract_text_content_recursive(child) + '\n')
            elif tag == 'p':
                # Check if paragraph contains only "---" separator
                child_text = self.extract_text_content_recursive(child).strip()
                if child_text == '---':
                    # Skip this separator paragraph
                    continue
                # Paragraphs - add line break after
                result.append(child_text + '\n')
            elif tag in ['ul', 'ol']:
                # Lists - add line break before
                result.append('\n' + self.extract_text_content_recursive(child))
            elif tag == 'li':
                # List items - add bullet and line break
                result.append('• ' + self.extract_text_content_recursive(child) + '\n')
            elif tag == 'br':
                # Line breaks
                result.append('\n')
            elif tag == 'strong' or tag == 'b':
                # Bold - just extract text (PPTX notes don't support formatting)
                result.append(self.extract_text_content_recursive(child))
            elif tag == 'em' or tag == 'i':
                # Italic - just extract text
                result.append(self.extract_text_content_recursive(child))
            else:
                # Other elements - just extract text
                result.append(self.extract_text_content_recursive(child))

            # Get tail text (text after the child element)
            if child.tail:
                # Only strip spaces and tabs, preserve newlines
                result.append(child.tail.strip(' \t'))

        # Join and clean up excessive whitespace
        text = ''.join(result)
        # Replace multiple consecutive newlines with just two (paragraph break)
        text = re.sub(r'\n\n+', '\n\n', text)
        # Strip leading/trailing whitespace
        text = text.strip()

        return text

    def apply_dark_background(self, slide):
        """Apply dark background to a slide"""
        # Add a rectangle shape covering the entire slide
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Inches(self.SLIDE_WIDTH),
            Inches(self.SLIDE_HEIGHT)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.COLORS['#131313']  # Dark gray
        bg_shape.line.fill.background()
        bg_shape.shadow.inherit = False

        # Move to back
        slide.shapes._spTree.remove(bg_shape._element)
        slide.shapes._spTree.insert(2, bg_shape._element)

    def convert(self, html_file, output_file):
        """Convert HTML file to PPTX"""
        print(f"Reading HTML file: {html_file}")

        # Parse HTML
        with open(html_file, 'r', encoding='utf-8') as f:
            html_content = f.read()

        tree = lhtml.fromstring(html_content)

        # Load CSS rules from <style> tags
        print("Loading CSS rules from <style> tags...")
        self.css_parser.load_css_rules_from_html(tree)
        print(f"  Loaded {len(self.css_parser.css_rules)} CSS rules")
        print(f"  Loaded {len(self.css_parser.css_vars)} CSS variables")

        # Extract course name from first title slide (for footer)
        course_name = ""
        title_slide = tree.xpath('//div[contains(@class, "title-slide")]//h1')
        if title_slide:
            # Replace <br> tags with spaces before extracting text
            h1_elem = title_slide[0]
            for br in h1_elem.xpath('.//br'):
                br.tail = ' ' + (br.tail or '')
            course_name = self.extract_text_content(h1_elem)

        # Find all slides (divs that have "slide" in their class attribute)
        slides = tree.xpath('//div[contains(concat(" ", normalize-space(@class), " "), " slide ")]')
        print(f"Found {len(slides)} slides")

        # Process each slide
        for i, html_slide in enumerate(slides, 1):
            print(f"Processing slide {i}/{len(slides)}")

            # Get slide classes to determine layout
            classes = html_slide.get('class', '').split()

            # Determine which layout to use
            # For content slides with bullets, use Title and Content layout for proper bullet support
            if 'content-slide' in classes or (len(classes) == 1 and classes[0] == 'slide'):
                has_bullets = html_slide.find('.//ul') is not None
                if has_bullets:
                    layout = self.prs.slide_layouts[1]  # Title and Content layout
                else:
                    layout = self.prs.slide_layouts[6]  # Blank layout
            else:
                layout = self.prs.slide_layouts[6]  # Blank layout for other types

            # Add slide with appropriate layout
            slide = self.prs.slides.add_slide(layout)

            # Apply dark background if needed
            if 'dark-bg' in classes:
                self.apply_dark_background(slide)

            # Extract and add speaker notes BEFORE processing slide content
            # (so we can remove them from the content)
            self.extract_and_add_speaker_notes(slide, html_slide)

            # Remove speaker notes from slide content before processing
            # (they've already been extracted to notes section)
            for notes_elem in html_slide.findall('.//aside[@class="speaker-notes"]'):
                notes_elem.getparent().remove(notes_elem)

            # Use handler registry to process slide
            # Handlers are tried in priority order until one can handle the slide
            handled = self.handler_registry.handle_slide(slide, html_slide)

            # Fallback to old methods for slide types not yet migrated to handlers
            if not handled:
                if 'vocab-table-slide' in classes:
                    self.handle_vocab_table_slide(slide, html_slide)
                elif 'comparison-slide' in classes:
                    self.handle_comparison_slide(slide, html_slide)
                elif 'objectives-slide' in classes:
                    self.handle_objectives_slide(slide, html_slide)
                elif 'activity-slide' in classes:
                    self.handle_activity_slide(slide, html_slide)
                elif 'assessment-checklist-slide' in classes:
                    self.handle_checklist_slide(slide, html_slide)
                # NEW ACADEMIC LAYOUTS
                elif 'quote-slide' in classes:
                    self.handle_quote_slide(slide, html_slide)
                elif 'reflection-slide' in classes:
                    self.handle_reflection_slide(slide, html_slide)
                elif 'comparison-table-slide' in classes:
                    self.handle_comparison_table_slide(slide, html_slide)
                elif 'framework-slide' in classes:
                    self.handle_framework_slide(slide, html_slide)
                elif 'references-slide' in classes:
                    self.handle_references_slide(slide, html_slide)
                else:
                    # Final fallback to content slide handler
                    self.handle_content_slide(slide, html_slide)

            # Add footer to all slides except title and section breaks
            if 'title-slide' not in classes and 'section-break-slide' not in classes:
                is_dark = 'dark-bg' in classes
                self.add_footer_system(slide, i, len(slides), course_name, is_dark)

        # Save presentation
        print(f"Saving PowerPoint to: {output_file}")
        self.prs.save(output_file)
        print(f"✓ Successfully created {output_file} with {len(slides)} slides")


def main():
    if len(sys.argv) < 3:
        print("Usage: python3 html_to_pptx_converter.py <input.html> <output.pptx> [--no-validate]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]
    skip_validation = '--no-validate' in sys.argv

    if not Path(input_file).exists():
        print(f"Error: Input file not found: {input_file}")
        sys.exit(1)

    converter = HTMLToPPTXConverter()
    converter.convert(input_file, output_file)

    # Run validation unless --no-validate flag
    if not skip_validation:
        try:
            from tools.validate_conversion import validate_conversion
            print("\nRunning quality validation...")
            report = validate_conversion(input_file, output_file)
            report.print_summary()

            # Exit with error code if validation failed
            if report.has_errors():
                sys.exit(1)
        except ImportError as e:
            print(f"Note: Validation skipped (module not found: {e})")


if __name__ == "__main__":
    main()
