#!/usr/bin/env python3
"""
Add Speaker Notes: Insert speaker notes into existing PPTX.

This script inserts speaker notes from lecture-content.md directly into
an existing Gemini-created PPTX file, saving in-place.

Usage:
    python tools/add_speaker_notes.py [course-code] [week-number]

Example:
    python tools/add_speaker_notes.py BCI2AU 1
"""

import os
import re
import sys
import glob
from pathlib import Path
from dataclasses import dataclass
from typing import Optional

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# =============================================================================
# DATA STRUCTURES
# =============================================================================

@dataclass
class SlideNotes:
    """Speaker notes for a single slide."""
    slide_number: int
    slide_title: str
    notes_content: str


@dataclass
class InsertionReport:
    """Result of speaker notes insertion."""
    success: bool
    pptx_file: str
    slides_total: int
    notes_inserted: int
    notes_skipped: int
    errors: list


# =============================================================================
# SPEAKER NOTES PARSER
# =============================================================================

def parse_speaker_notes(lecture_content_path: str) -> list[SlideNotes]:
    """
    Parse speaker notes from lecture-content.md.

    Supported formats:
    1. ## Slide N: Title ... ### Speaker Notes (current format)
    2. **SLIDE N: Title** ... ## Speaker Notes (legacy format)
    """
    with open(lecture_content_path, 'r', encoding='utf-8') as f:
        content = f.read()

    notes_list = []

    # Try current format first: ## Slide N: Title
    current_pattern = r'##\s+Slide\s+(\d+):\s*([^\n]+)'
    slides = re.split(r'(?=##\s+Slide\s+\d+:)', content)

    for slide_block in slides:
        if not slide_block.strip():
            continue

        # Extract slide number and title
        slide_match = re.search(current_pattern, slide_block)
        if not slide_match:
            continue

        slide_num = int(slide_match.group(1))
        slide_title = slide_match.group(2).strip()

        # Extract speaker notes section (### Speaker Notes for current format)
        notes_match = re.search(
            r'###\s*Speaker\s*Notes\s*\n(.*?)(?=\n---|\n##\s+Slide|\Z)',
            slide_block,
            re.DOTALL | re.IGNORECASE
        )

        if notes_match:
            notes_content = notes_match.group(1).strip()
            notes_content = clean_notes_content(notes_content)

            notes_list.append(SlideNotes(
                slide_number=slide_num,
                slide_title=slide_title,
                notes_content=notes_content
            ))

    # If no notes found with current format, try legacy format
    if not notes_list:
        legacy_pattern = r'\*\*SLIDE\s+(\d+):\s*([^*]+)\*\*'
        slides = re.split(r'(?=\*\*SLIDE\s+\d+:)', content)

        for slide_block in slides:
            if not slide_block.strip():
                continue

            slide_match = re.search(legacy_pattern, slide_block)
            if not slide_match:
                continue

            slide_num = int(slide_match.group(1))
            slide_title = slide_match.group(2).strip()

            notes_match = re.search(
                r'##\s*Speaker\s*Notes\s*\n(.*?)(?=\n---|\n\*\*SLIDE|\Z)',
                slide_block,
                re.DOTALL | re.IGNORECASE
            )

            if notes_match:
                notes_content = notes_match.group(1).strip()
                notes_content = clean_notes_content(notes_content)

                notes_list.append(SlideNotes(
                    slide_number=slide_num,
                    slide_title=slide_title,
                    notes_content=notes_content
                ))

    return notes_list


def clean_notes_content(notes: str) -> str:
    """Clean and format speaker notes for PowerPoint."""
    # Remove excessive whitespace
    notes = re.sub(r'\n{3,}', '\n\n', notes)

    # Convert markdown bold to plain text (PowerPoint notes don't support markdown)
    notes = re.sub(r'\*\*([^*]+)\*\*', r'\1', notes)

    # Convert markdown italic
    notes = re.sub(r'\*([^*]+)\*', r'\1', notes)

    # Clean up bullet points (keep them but normalize)
    notes = re.sub(r'^[-*]\s+', '- ', notes, flags=re.MULTILINE)

    return notes.strip()


# =============================================================================
# PPTX OPERATIONS
# =============================================================================

def find_pptx_file(week_folder: str) -> Optional[str]:
    """
    Find the PPTX file in the week folder or output/ subfolder.

    Searches: week root first, then output/ subfolder
    Excludes: slides.pptx (our output), any file with "batch" in name
    Prefers: Files with course code or "Lecture" in name
    """
    # Search both week root and output/ folder
    search_paths = [
        os.path.join(week_folder, "*.pptx"),
        os.path.join(week_folder, "output", "*.pptx")
    ]

    all_pptx = []
    for pattern in search_paths:
        all_pptx.extend(glob.glob(pattern))

    # Filter out unwanted files
    candidates = []
    for f in all_pptx:
        basename = os.path.basename(f).lower()
        # Exclude batch files and our output file slides.pptx
        if 'batch' in basename:
            continue
        if basename == 'slides.pptx':
            continue
        candidates.append(f)

    if not candidates:
        return None

    # If only one candidate, use it
    if len(candidates) == 1:
        return candidates[0]

    # Prefer files with "Lecture" in name
    for f in candidates:
        if 'lecture' in os.path.basename(f).lower():
            return f

    # Return first candidate
    return candidates[0]


def insert_speaker_notes(pptx_path: str, notes_list: list[SlideNotes], output_path: str = None) -> tuple[int, int]:
    """
    Insert speaker notes into presentation slides.

    Opens the PPTX, inserts notes, saves to output_path (or in-place if not specified).

    Returns: (notes_inserted, notes_skipped)
    """
    prs = Presentation(pptx_path)

    notes_inserted = 0
    notes_skipped = 0

    # Create a lookup by slide number
    notes_lookup = {n.slide_number: n for n in notes_list}

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        if slide_num in notes_lookup:
            notes = notes_lookup[slide_num]

            # Get or create notes slide
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_placeholder.text_frame

            # Clear existing notes and add new content
            notes_tf.text = notes.notes_content

            notes_inserted += 1
        else:
            notes_skipped += 1

    # Save to output path or in-place
    save_path = output_path if output_path else pptx_path
    prs.save(save_path)

    return notes_inserted, notes_skipped


# =============================================================================
# MAIN FUNCTION
# =============================================================================

def add_speaker_notes(course_code: str, week_number: int, base_path: str = None) -> InsertionReport:
    """
    Main function to add speaker notes to a week's PPTX.

    1. Find PPTX file in week folder (or output/ subfolder)
    2. Parse speaker notes from lecture-content.md
    3. Insert notes into each slide
    4. Save to output/slides.pptx
    """
    if base_path is None:
        base_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    # Construct paths
    course_folder = os.path.join(base_path, "courses", f"{course_code}*")
    course_matches = glob.glob(course_folder)

    if not course_matches:
        return InsertionReport(
            success=False,
            pptx_file="",
            slides_total=0,
            notes_inserted=0,
            notes_skipped=0,
            errors=[f"Course folder not found for: {course_code}"]
        )

    course_path = course_matches[0]
    week_folder = os.path.join(course_path, "weeks", f"week-{week_number:02d}")
    lecture_content_path = os.path.join(week_folder, "lecture-content.md")

    errors = []

    # 1. Find PPTX file
    print(f"\n1. Finding PPTX file in: {week_folder}")
    pptx_file = find_pptx_file(week_folder)

    if not pptx_file:
        return InsertionReport(
            success=False,
            pptx_file="",
            slides_total=0,
            notes_inserted=0,
            notes_skipped=0,
            errors=["No PPTX file found in week folder (excluding batch files and slides.pptx)"]
        )

    print(f"   Found: {os.path.basename(pptx_file)}")

    # 2. Parse speaker notes
    print(f"\n2. Parsing speaker notes from: {os.path.basename(lecture_content_path)}")

    if not os.path.exists(lecture_content_path):
        return InsertionReport(
            success=False,
            pptx_file=pptx_file,
            slides_total=0,
            notes_inserted=0,
            notes_skipped=0,
            errors=[f"Lecture content not found: {lecture_content_path}"]
        )

    notes_list = parse_speaker_notes(lecture_content_path)
    print(f"   Parsed {len(notes_list)} speaker notes sections")

    if not notes_list:
        return InsertionReport(
            success=False,
            pptx_file=pptx_file,
            slides_total=0,
            notes_inserted=0,
            notes_skipped=0,
            errors=["No speaker notes found in lecture-content.md"]
        )

    # 3. Count slides before insertion
    prs = Presentation(pptx_file)
    slides_total = len(prs.slides)
    print(f"   PPTX has {slides_total} slides")

    # 4. Insert speaker notes and save to output/ folder
    print(f"\n3. Inserting speaker notes...")
    output_folder = os.path.join(week_folder, "output")
    os.makedirs(output_folder, exist_ok=True)
    output_pptx = os.path.join(output_folder, "slides.pptx")
    notes_inserted, notes_skipped = insert_speaker_notes(pptx_file, notes_list, output_pptx)
    print(f"   Inserted: {notes_inserted}, Skipped: {notes_skipped}")

    # 5. Verify save
    print(f"\n4. Saved to: output/slides.pptx")

    # Determine success
    success = notes_inserted > 0 and not errors

    return InsertionReport(
        success=success,
        pptx_file=output_pptx,
        slides_total=slides_total,
        notes_inserted=notes_inserted,
        notes_skipped=notes_skipped,
        errors=errors
    )


def print_report(report: InsertionReport):
    """Print a formatted report of the insertion process."""
    print("\n" + "=" * 60)
    print("SPEAKER NOTES INSERTION REPORT")
    print("=" * 60)

    status = "SUCCESS" if report.success else "FAILED"
    print(f"\nStatus: {status}")
    print(f"File: {report.pptx_file}")

    print(f"\nSummary:")
    print(f"  - Total slides: {report.slides_total}")
    print(f"  - Notes inserted: {report.notes_inserted}")
    print(f"  - Slides without notes: {report.notes_skipped}")

    if report.notes_inserted > 0:
        coverage = (report.notes_inserted / report.slides_total * 100) if report.slides_total > 0 else 0
        print(f"  - Coverage: {coverage:.0f}%")

    if report.errors:
        print(f"\nErrors ({len(report.errors)}):")
        for e in report.errors:
            print(f"  - {e}")

    print("\n" + "=" * 60)


# =============================================================================
# CLI ENTRY POINT
# =============================================================================

def main():
    """Command-line entry point."""
    if len(sys.argv) < 3:
        print("Usage: python tools/add_speaker_notes.py [course-code] [week-number]")
        print("Example: python tools/add_speaker_notes.py BCI2AU 1")
        sys.exit(1)

    course_code = sys.argv[1]
    week_number = int(sys.argv[2])

    print(f"Adding speaker notes for {course_code} Week {week_number}")
    print("-" * 40)

    report = add_speaker_notes(course_code, week_number)
    print_report(report)

    sys.exit(0 if report.success else 1)


if __name__ == "__main__":
    main()
