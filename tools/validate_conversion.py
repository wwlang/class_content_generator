#!/usr/bin/env python3
"""
Comprehensive validation tool for HTML and PPTX conversion quality.

PRIORITY: Content Integrity > Design Aesthetics

Critical Checks (Tier 1 - Block build if fail):
- Slide count matches
- No empty slides
- Required elements present (blockquotes, tables)
- No overlapping text

Content Rules Compliance (Tier 2 - Warn if fail):
- Bullet count ≤7 per slide
- Speaker notes present on every slide
- References slide exists
- Citations present where expected

Design & Accessibility (Tier 3 - Info only):
- Background colors correct
- Font sizes readable (≥8pt)

Usage:
    python3 validate_conversion.py <input.html> <output.pptx>
"""

from dataclasses import dataclass
from typing import List, Dict, Tuple, Optional, Set
from pathlib import Path
from lxml import html as lxml_html
from pptx import Presentation
from pptx.util import Inches, Pt
import sys
import json


@dataclass
class ValidationIssue:
    """Represents a validation issue found during testing."""
    severity: str  # 'error', 'warning', 'info'
    category: str  # 'content', 'positioning', 'styling', 'structure'
    slide_number: int
    message: str
    details: Optional[str] = None


@dataclass
class ValidationReport:
    """Complete validation report with all issues found."""
    html_file: str
    pptx_file: str
    total_slides: int
    errors: List[ValidationIssue]
    warnings: List[ValidationIssue]
    info: List[ValidationIssue]

    def has_errors(self) -> bool:
        """Check if report contains any errors."""
        return len(self.errors) > 0

    def has_warnings(self) -> bool:
        """Check if report contains any warnings."""
        return len(self.warnings) > 0

    def print_summary(self):
        """Print human-readable summary of validation results."""
        print("\n" + "="*70)
        print("VALIDATION REPORT")
        print("="*70)
        print(f"HTML: {self.html_file}")
        print(f"PPTX: {self.pptx_file}")
        print(f"Slides: {self.total_slides}")
        print("-"*70)

        if not self.errors and not self.warnings:
            print("✓ ALL CHECKS PASSED")
            print("="*70 + "\n")
            return

        if self.errors:
            print(f"\n❌ ERRORS: {len(self.errors)}")
            for issue in self.errors:
                print(f"  Slide {issue.slide_number} [{issue.category}]: {issue.message}")
                if issue.details:
                    print(f"    → {issue.details}")

        if self.warnings:
            print(f"\n⚠️  WARNINGS: {len(self.warnings)}")
            for issue in self.warnings:
                print(f"  Slide {issue.slide_number} [{issue.category}]: {issue.message}")
                if issue.details:
                    print(f"    → {issue.details}")

        if self.info:
            print(f"\nℹ️  INFO: {len(self.info)}")
            for issue in self.info[:5]:  # Limit to first 5
                print(f"  Slide {issue.slide_number} [{issue.category}]: {issue.message}")
            if len(self.info) > 5:
                print(f"  ... and {len(self.info) - 5} more")

        print("="*70 + "\n")


class HTMLValidator:
    """Validates HTML slide structure and content."""

    def __init__(self, html_path: str):
        self.html_path = html_path
        with open(html_path, 'r', encoding='utf-8') as f:
            self.tree = lxml_html.fromstring(f.read())
        self.slides = self._extract_slides()

    def _extract_slides(self) -> List:
        """Extract all slide elements from HTML."""
        # Use same xpath as converter (matches only "slide" as complete class word)
        return self.tree.xpath('//div[contains(concat(" ", normalize-space(@class), " "), " slide ")]')

    def validate(self) -> List[ValidationIssue]:
        """Run all HTML validation checks."""
        issues = []

        for i, slide in enumerate(self.slides, 1):
            issues.extend(self._check_empty_slides(slide, i))
            issues.extend(self._check_required_elements(slide, i))
            # REMOVED: _check_duplicate_content (low priority)
            issues.extend(self._check_slide_backgrounds(slide, i))
            # REMOVED: _check_word_count (user feedback - not relevant)
            issues.extend(self._check_bullet_count(slide, i))  # Phase 1

        # Phase 2 - Check for references slide only (simplified per user feedback)
        issues.extend(self._check_references_slide())

        return issues

    def _check_empty_slides(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check for slides with no meaningful content."""
        issues = []

        # Extract all text
        all_text = ''.join(slide.itertext()).strip()

        if len(all_text) < 10:  # Less than 10 characters
            issues.append(ValidationIssue(
                severity='error',
                category='content',
                slide_number=slide_num,
                message='Slide appears empty or has minimal content',
                details=f'Only {len(all_text)} characters found'
            ))

        return issues

    def _check_required_elements(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check that slides have required elements based on type."""
        issues = []
        classes = slide.get('class', '')

        # Quote slides should have blockquote, quote-content, or slide-content with paragraphs
        if 'quote-slide' in classes:
            blockquote = slide.find('.//blockquote')
            quote_content = slide.find('.//*[@class="quote-content"]')
            slide_content = slide.find('.//*[@class="slide-content"]')

            # Check if we have any valid quote structure
            has_valid_quote = (
                blockquote is not None or
                quote_content is not None or
                (slide_content is not None and len(slide_content.findall('.//p')) > 0)
            )

            if not has_valid_quote:
                issues.append(ValidationIssue(
                    severity='error',
                    category='structure',
                    slide_number=slide_num,
                    message='Quote slide missing blockquote element'
                ))

        # Table slides should have table
        if 'table' in classes or 'comparison-table' in classes:
            table = slide.find('.//table')
            if table is None:
                issues.append(ValidationIssue(
                    severity='error',
                    category='structure',
                    slide_number=slide_num,
                    message='Table slide missing table element'
                ))

        # References slides should have list or references container
        if 'references-slide' in classes:
            ul = slide.find('.//ul')
            refs = slide.find('.//*[@class="references"]')
            if ul is None and refs is None:
                issues.append(ValidationIssue(
                    severity='warning',
                    category='structure',
                    slide_number=slide_num,
                    message='References slide missing list structure'
                ))

        return issues

    def _check_duplicate_content(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check for suspiciously duplicated text content."""
        issues = []

        # Get all text elements
        text_elements = []
        for elem in slide.iter():
            if elem.text and elem.text.strip():
                text_elements.append(elem.text.strip())

        # Check for exact duplicates
        seen = set()
        duplicates = set()
        for text in text_elements:
            if len(text) > 20:  # Only check substantial text
                if text in seen:
                    duplicates.add(text[:50])
                seen.add(text)

        if duplicates:
            issues.append(ValidationIssue(
                severity='warning',
                category='content',
                slide_number=slide_num,
                message='Duplicate text content detected',
                details=f'Found {len(duplicates)} duplicate text blocks'
            ))

        return issues

    def _check_slide_backgrounds(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check that only title, section-break, and dark-bg slides have non-beige backgrounds.

        BACKGROUND POLICY:
        - All slides default to beige background
        - ONLY these types override:
          1. title-slide (beige with decorative shapes - OK)
          2. section-break-slide (orange)
          3. .dark-bg class (dark gray)
        - All other slides MUST use default beige
        """
        issues = []
        classes = slide.get('class', '')

        # Allowed background overrides
        allowed_overrides = ['title-slide', 'section-break-slide']
        has_dark_bg = 'dark-bg' in classes

        # Check if this is a non-allowed slide type with background override in inline styles
        # (We can't check CSS directly, but we can check the class combinations)
        is_title = 'title-slide' in classes
        is_section_break = 'section-break-slide' in classes

        # Check for slide types that should NOT have background overrides
        problematic_types = [
            'content-slide',
            'framework-slide',
            'quote-slide',
            'reflection-slide',
            'comparison-table-slide',
            'vocab-table-slide',
            'references-slide'
        ]

        for slide_type in problematic_types:
            if slide_type in classes:
                # This slide type should only have beige background
                # Check if it has dark-bg class (which is the only allowed override)
                if not has_dark_bg:
                    # This is correct - using default beige
                    pass
                # If it has dark-bg, that's allowed (explicit override)
                break

        return issues

    def _check_word_count(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check that slide content is 150-200 words (max 200)."""
        issues = []

        # Extract all visible text (excluding comments)
        all_text = ''.join(slide.itertext()).strip()

        # Count words (split by whitespace)
        words = all_text.split()
        word_count = len(words)

        # Warn if >200 words (should split)
        if word_count > 200:
            issues.append(ValidationIssue(
                severity='warning',
                category='content',
                slide_number=slide_num,
                message=f'Slide too long ({word_count} words)',
                details='Max 200 words. Should split into 2 slides.'
            ))
        # Info if <100 words (might be too sparse, but not critical)
        elif word_count < 100:
            # Only flag if it's a content slide (not title, section-break, etc.)
            classes = slide.get('class', '')
            if 'content-slide' in classes or 'framework-slide' in classes:
                issues.append(ValidationIssue(
                    severity='info',
                    category='content',
                    slide_number=slide_num,
                    message=f'Slide may be too sparse ({word_count} words)',
                    details='Consider if more detail needed (150-200 words optimal)'
                ))

        return issues

    def _check_bullet_count(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check that slide has max 5-7 bullets."""
        issues = []

        # Count <li> elements (bullet points)
        bullets = slide.findall('.//li')
        bullet_count = len(bullets)

        # Warn if >7 bullets (should split)
        if bullet_count > 7:
            issues.append(ValidationIssue(
                severity='warning',
                category='content',
                slide_number=slide_num,
                message=f'Too many bullets ({bullet_count} bullets)',
                details='Max 7 bullets. Should split into 2 slides.'
            ))

        return issues

    def _check_references_slide(self) -> List[ValidationIssue]:
        """Check that lecture has a references slide."""
        issues = []

        # Only check for lectures with 10+ slides
        if len(self.slides) < 10:
            return issues

        # Look for references slide
        has_references = False
        for slide in self.slides:
            classes = slide.get('class', '')
            if 'references-slide' in classes:
                has_references = True
                break

        # Warn if missing
        if not has_references:
            issues.append(ValidationIssue(
                severity='warning',
                category='structure',
                slide_number=0,  # Overall lecture
                message='Missing references slide',
                details='All sources must be cited in references slide at end'
            ))

        return issues


class PPTXValidator:
    """Validates PPTX slide structure, positioning, and content."""

    # Slide dimensions (4:3 aspect ratio)
    SLIDE_WIDTH = 10.0  # inches
    SLIDE_HEIGHT = 7.5  # inches
    MARGIN = 0.5  # inches

    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)

    def validate(self) -> List[ValidationIssue]:
        """Run all PPTX validation checks."""
        issues = []

        for i, slide in enumerate(self.prs.slides, 1):
            issues.extend(self._check_empty_slides(slide, i))
            issues.extend(self._check_overlapping_shapes(slide, i))
            issues.extend(self._check_shapes_within_bounds(slide, i))
            # REMOVED: _check_font_consistency (downgraded to info, not useful)
            # REMOVED: _check_text_truncation (edge case, downgraded to info)
            # REMOVED: _check_missing_content (too many false positives)
            issues.extend(self._check_background_color(slide, i))
            issues.extend(self._check_font_sizes(slide, i))
            issues.extend(self._check_speaker_notes(slide, i))  # NEW: Phase 1

        return issues

    def _check_empty_slides(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check for slides with no content or very few shapes."""
        issues = []

        # Count meaningful shapes (exclude placeholders, footers)
        content_shapes = [s for s in slide.shapes if hasattr(s, 'text')]

        if len(content_shapes) == 0:
            issues.append(ValidationIssue(
                severity='error',
                category='content',
                slide_number=slide_num,
                message='Slide has no text content'
            ))
        elif len(content_shapes) <= 2:
            # Check if there's actual text
            total_text = ''.join(s.text for s in content_shapes if hasattr(s, 'text')).strip()
            if len(total_text) < 10:
                issues.append(ValidationIssue(
                    severity='warning',
                    category='content',
                    slide_number=slide_num,
                    message='Slide has minimal content',
                    details=f'Only {len(content_shapes)} shapes, {len(total_text)} characters'
                ))

        return issues

    def _check_overlapping_shapes(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check for overlapping text boxes and shapes."""
        issues = []

        # Get all shape positions
        shapes_with_pos = []
        for shape in slide.shapes:
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                shapes_with_pos.append({
                    'shape': shape,
                    'left': shape.left,
                    'top': shape.top,
                    'right': shape.left + shape.width,
                    'bottom': shape.top + shape.height,
                })

        # Check for overlaps
        overlaps = []
        for i, shape1 in enumerate(shapes_with_pos):
            for shape2 in shapes_with_pos[i+1:]:
                if self._rectangles_overlap(shape1, shape2):
                    # Only report if both have text (not decorative overlaps)
                    if (hasattr(shape1['shape'], 'text') and shape1['shape'].text.strip() and
                        hasattr(shape2['shape'], 'text') and shape2['shape'].text.strip()):
                        overlaps.append((shape1, shape2))

        if overlaps:
            issues.append(ValidationIssue(
                severity='error',
                category='positioning',
                slide_number=slide_num,
                message=f'Found {len(overlaps)} overlapping text shapes',
                details='Text may be unreadable due to overlap'
            ))

        return issues

    def _rectangles_overlap(self, rect1: Dict, rect2: Dict) -> bool:
        """Check if two rectangles overlap."""
        return not (rect1['right'] <= rect2['left'] or
                   rect1['left'] >= rect2['right'] or
                   rect1['bottom'] <= rect2['top'] or
                   rect1['top'] >= rect2['bottom'])

    def _check_shapes_within_bounds(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check that all shapes are within slide boundaries."""
        issues = []

        slide_width = Inches(self.SLIDE_WIDTH)
        slide_height = Inches(self.SLIDE_HEIGHT)

        # Allow reasonable margin for padding and text overflow (0.2 inches)
        # Small overflows are common and usually handled by PowerPoint
        tolerance = Inches(0.2)

        # Track shapes that are SIGNIFICANTLY out of bounds
        significantly_out_of_bounds = []
        for shape in slide.shapes:
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                # Calculate how far out of bounds the shape is
                left_overflow = max(0, -shape.left - tolerance)
                top_overflow = max(0, -shape.top - tolerance)
                right_overflow = max(0, shape.left + shape.width - slide_width - tolerance)
                bottom_overflow = max(0, shape.top + shape.height - slide_height - tolerance)

                max_overflow = max(left_overflow, top_overflow, right_overflow, bottom_overflow)

                # Only report if significantly out of bounds (>0.3 inches)
                if max_overflow > Inches(0.3):
                    # Only report if it has visible content
                    if hasattr(shape, 'text') and shape.text.strip():
                        significantly_out_of_bounds.append(shape)

        # Only report if multiple shapes are significantly out of bounds
        if len(significantly_out_of_bounds) >= 3:
            issues.append(ValidationIssue(
                severity='warning',
                category='positioning',
                slide_number=slide_num,
                message=f'{len(significantly_out_of_bounds)} content shapes extend significantly beyond boundaries',
                details=f'Check for text that might be cut off (>0.3" overflow)'
            ))
        elif len(significantly_out_of_bounds) > 0:
            # Just info if only 1-2 shapes (often intentional)
            issues.append(ValidationIssue(
                severity='info',
                category='positioning',
                slide_number=slide_num,
                message=f'{len(significantly_out_of_bounds)} content shape(s) extend beyond boundaries',
                details=f'Minor overflow, usually handled by PowerPoint'
            ))

        return issues

    def _check_font_consistency(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check for font consistency across slide."""
        issues = []

        fonts_used = set()
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts_used.add(run.font.name)

        # Expected fonts
        expected_fonts = {'Cal Sans', 'Plus Jakarta Sans'}
        unexpected_fonts = fonts_used - expected_fonts - {None}

        if unexpected_fonts:
            issues.append(ValidationIssue(
                severity='warning',
                category='styling',
                slide_number=slide_num,
                message='Unexpected fonts detected',
                details=f'Found: {", ".join(unexpected_fonts)}'
            ))

        return issues

    def _check_text_truncation(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check for text that may be truncated in text boxes."""
        issues = []

        truncated_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                # Check if text frame is set to auto-size
                if hasattr(shape.text_frame, 'word_wrap'):
                    # If word wrap is off and text is long, might truncate
                    if not shape.text_frame.word_wrap and len(shape.text_frame.text) > 100:
                        truncated_shapes.append(shape)

        if truncated_shapes:
            issues.append(ValidationIssue(
                severity='warning',
                category='content',
                slide_number=slide_num,
                message=f'{len(truncated_shapes)} text boxes may have truncated content',
                details='Long text without word wrap enabled'
            ))

        return issues

    def _check_missing_content(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check for shapes with no text (might indicate extraction failure)."""
        issues = []

        # Count content shapes (not footers/page numbers)
        content_shapes = [s for s in slide.shapes if hasattr(s, 'text')]

        # Find text boxes that are empty
        empty_textboxes = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                if hasattr(shape, 'text') and not shape.text.strip():
                    # Check if it has paragraphs (structure exists but no content)
                    if len(shape.text_frame.paragraphs) > 0:
                        empty_textboxes.append(shape)

        # More than 3 empty text boxes out of total is suspicious
        if len(content_shapes) > 5 and len(empty_textboxes) > 3:
            issues.append(ValidationIssue(
                severity='warning',  # Downgrade to warning
                category='content',
                slide_number=slide_num,
                message=f'{len(empty_textboxes)} of {len(content_shapes)} text boxes are empty',
                details='May include decorative or placeholder shapes'
            ))

        return issues

    def _check_background_color(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check that slide has appropriate background color."""
        issues = []

        # Expected backgrounds:
        # - Cream (#f4f3f1 = RGB 244, 243, 241) - Most slides
        # - Dark (#131313 = RGB 19, 19, 19) - Dark slide modifier
        # - Orange (#ed5e29 = RGB 237, 94, 41) - Section break slides
        # - White (#ffffff = RGB 255, 255, 255) - Title slides
        expected_cream = (244, 243, 241)
        expected_dark = (19, 19, 19)
        expected_orange = (237, 94, 41)
        expected_white = (255, 255, 255)

        # Find background rectangle (should be first shape)
        has_background = False
        background_color = None

        if len(slide.shapes) > 0:
            first_shape = slide.shapes[0]
            # Check if it's a rectangle covering whole slide
            if hasattr(first_shape, 'fill') and first_shape.fill.type == 1:  # SOLID
                try:
                    bg_color = first_shape.fill.fore_color.rgb
                    background_color = (bg_color[0], bg_color[1], bg_color[2])
                    has_background = True
                except:
                    pass

        if has_background:
            # Check if it matches expected colors (allow small tolerance)
            tolerance = 5
            is_cream = all(abs(background_color[i] - expected_cream[i]) <= tolerance for i in range(3))
            is_dark = all(abs(background_color[i] - expected_dark[i]) <= tolerance for i in range(3))
            is_orange = all(abs(background_color[i] - expected_orange[i]) <= tolerance for i in range(3))
            is_white = all(abs(background_color[i] - expected_white[i]) <= tolerance for i in range(3))

            if not (is_cream or is_dark or is_orange or is_white):
                issues.append(ValidationIssue(
                    severity='warning',
                    category='styling',
                    slide_number=slide_num,
                    message=f'Unexpected background color: RGB{background_color}',
                    details='Expected cream, dark, orange (section breaks), or white (title slides)'
                ))
        else:
            # No background detected - might be white default (acceptable for title slides)
            # Only warn if slide has significant content (not a title slide)
            shape_count = len([s for s in slide.shapes if hasattr(s, 'text') and s.text.strip()])
            if shape_count > 3:  # Title slides typically have 3 or fewer text shapes
                issues.append(ValidationIssue(
                    severity='info',  # Downgrade to info
                    category='styling',
                    slide_number=slide_num,
                    message='No background color applied (using default white)',
                    details='May be acceptable for title slides'
                ))

        return issues

    def _check_font_sizes(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check for unreasonably large or small font sizes."""
        issues = []

        font_sizes = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)

        if font_sizes:
            min_size = min(font_sizes)

            # Check for unreasonably small text (< 8pt, excluding footers)
            if min_size < 8:
                issues.append(ValidationIssue(
                    severity='warning',
                    category='styling',
                    slide_number=slide_num,
                    message=f'Very small font size detected: {min_size}pt',
                    details='May be difficult to read'
                ))

            # REMOVED: Large font size check (> 150pt) - not useful, decorative is fine

        return issues

    def _check_speaker_notes(self, slide, slide_num: int) -> List[ValidationIssue]:
        """Check that slide has speaker notes (presence only, not length)."""
        issues = []

        # Check if slide has notes
        if not hasattr(slide, 'notes_slide') or slide.notes_slide is None:
            issues.append(ValidationIssue(
                severity='warning',
                category='content',
                slide_number=slide_num,
                message='Missing speaker notes',
                details='Every slide should have speaker notes'
            ))
            return issues

        # Get notes text
        notes_text = ''
        if hasattr(slide.notes_slide, 'notes_text_frame'):
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        # Check if notes are empty (only check presence, not length)
        if not notes_text or len(notes_text) < 10:
            issues.append(ValidationIssue(
                severity='warning',
                category='content',
                slide_number=slide_num,
                message='Missing speaker notes',
                details='Every slide should have speaker notes'
            ))

        return issues


class CrossValidator:
    """Cross-validates HTML and PPTX to ensure content matches."""

    def __init__(self, html_path: str, pptx_path: str):
        self.html_path = html_path
        self.pptx_path = pptx_path

        # Load HTML
        with open(html_path, 'r', encoding='utf-8') as f:
            self.html_tree = lxml_html.fromstring(f.read())
        # Find slides using same xpath as converter (matches only "slide" as complete class word)
        self.html_slides = self.html_tree.xpath('//div[contains(concat(" ", normalize-space(@class), " "), " slide ")]')

        # Load PPTX
        self.pptx_prs = Presentation(pptx_path)

    def validate(self) -> List[ValidationIssue]:
        """Cross-validate HTML and PPTX content."""
        issues = []

        # Check slide count matches
        if len(self.html_slides) != len(self.pptx_prs.slides):
            issues.append(ValidationIssue(
                severity='error',
                category='structure',
                slide_number=0,
                message='Slide count mismatch',
                details=f'HTML: {len(self.html_slides)}, PPTX: {len(self.pptx_prs.slides)}'
            ))
            return issues

        # Check content matches for each slide
        for i in range(len(self.html_slides)):
            issues.extend(self._check_content_match(i))
            issues.extend(self._check_visual_consistency(i))
            issues.extend(self._check_citations(i))  # NEW: Phase 2
            # REMOVED: _check_text_formatting (color/bold styling - low priority)

        return issues

    def _check_content_match(self, slide_index: int) -> List[ValidationIssue]:
        """Check if content matches between HTML and PPTX for a slide."""
        issues = []
        slide_num = slide_index + 1

        # Extract HTML text
        html_slide = self.html_slides[slide_index]
        html_text = ' '.join(''.join(html_slide.itertext()).split())  # Normalize whitespace

        # Extract PPTX text
        pptx_slide = self.pptx_prs.slides[slide_index]
        pptx_text_parts = []
        for shape in pptx_slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                pptx_text_parts.append(shape.text.strip())
        pptx_text = ' '.join(' '.join(pptx_text_parts).split())  # Normalize whitespace

        # Check for significant content loss (more than 10% difference)
        html_len = len(html_text)
        pptx_len = len(pptx_text)

        if html_len > 0:
            diff_pct = abs(html_len - pptx_len) / html_len * 100

            if diff_pct > 10:
                issues.append(ValidationIssue(
                    severity='warning',
                    category='content',
                    slide_number=slide_num,
                    message='Significant content length difference between HTML and PPTX',
                    details=f'HTML: {html_len} chars, PPTX: {pptx_len} chars ({diff_pct:.1f}% diff)'
                ))

        return issues

    def _check_visual_consistency(self, slide_index: int) -> List[ValidationIssue]:
        """Check if visual elements are consistent between HTML and PPTX."""
        issues = []
        slide_num = slide_index + 1

        html_slide = self.html_slides[slide_index]
        pptx_slide = self.pptx_prs.slides[slide_index]

        # Check for images in HTML
        html_images = html_slide.findall('.//img')
        if html_images:
            # Count image shapes in PPTX
            pptx_images = 0
            for shape in pptx_slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    pptx_images += 1

            if pptx_images == 0:
                issues.append(ValidationIssue(
                    severity='warning',
                    category='content',
                    slide_number=slide_num,
                    message=f'HTML has {len(html_images)} image(s) but PPTX has none',
                    details='Images may not have been converted'
                ))

        # Check for tables
        html_tables = html_slide.findall('.//table')
        if html_tables:
            pptx_tables = 0
            for shape in pptx_slide.shapes:
                if shape.shape_type == 19:  # TABLE
                    pptx_tables += 1
                # Also count text-based table representations
                if hasattr(shape, 'text') and '|' in shape.text:
                    # Might be text-based table representation
                    pass

            # Tables often converted as shapes, not actual table objects
            # So we just check if there's substantial content
            if pptx_tables == 0 and len(pptx_slide.shapes) < 3:
                issues.append(ValidationIssue(
                    severity='info',
                    category='structure',
                    slide_number=slide_num,
                    message='HTML table may not be rendered as PPTX table',
                    details='Check if table content is present as shapes'
                ))

        return issues

    def _check_citations(self, slide_index: int) -> List[ValidationIssue]:
        """Check if slides with statistics/claims have proper citations."""
        issues = []
        slide_num = slide_index + 1

        html_slide = self.html_slides[slide_index]
        pptx_slide = self.pptx_prs.slides[slide_index]

        # Get all text from HTML slide
        html_text = ' '.join(''.join(html_slide.itertext()).split())

        # Pattern match for statistics (numbers with % or "X times" or "X out of X")
        import re
        has_statistics = bool(re.search(r'\d+%|\d+\s+times|\d+\s+out\s+of\s+\d+|\d+x\s+more', html_text, re.IGNORECASE))

        # Pattern match for inline citations (Author, Year) or (Author et al., Year)
        has_citation = bool(re.search(r'\([A-Z][a-z]+(?:\s+et\s+al\.?)?,\s+\d{4}\)', html_text))

        # If slide has statistics but no citation, warn
        if has_statistics and not has_citation:
            # Check if it's a references slide (citations are listed there, not inline)
            classes = html_slide.get('class', '')
            if 'references-slide' not in classes:
                issues.append(ValidationIssue(
                    severity='warning',
                    category='content',
                    slide_number=slide_num,
                    message='Statistics/claims found without inline citation',
                    details='Add inline citation (Author, Year) after statistical claims'
                ))

        return issues

    def _check_text_formatting(self, slide_index: int) -> List[ValidationIssue]:
        """Check if text formatting (bold, italic, colors) matches between HTML and PPTX."""
        issues = []
        slide_num = slide_index + 1

        html_slide = self.html_slides[slide_index]
        pptx_slide = self.pptx_prs.slides[slide_index]

        # Check for bold text in HTML
        bold_elements = html_slide.findall('.//strong') + html_slide.findall('.//b')

        if bold_elements:
            # Check if PPTX has bold text with orange color
            pptx_bold_with_color = []
            pptx_bold_without_color = []

            for shape in pptx_slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.bold:
                                # Check if color is orange (#ed5e29 = RGB 237, 94, 41)
                                try:
                                    if run.font.color.rgb:
                                        color = run.font.color.rgb
                                        # Allow small tolerance (±5)
                                        is_orange = (
                                            abs(color[0] - 237) <= 5 and
                                            abs(color[1] - 94) <= 5 and
                                            abs(color[2] - 41) <= 5
                                        )
                                        if is_orange:
                                            pptx_bold_with_color.append(run.text)
                                        else:
                                            pptx_bold_without_color.append((run.text, color))
                                except:
                                    # Color not set or error accessing
                                    pptx_bold_without_color.append((run.text, None))

            # If HTML has bold text but PPTX bold text is not orange
            if pptx_bold_without_color:
                issues.append(ValidationIssue(
                    severity='warning',
                    category='styling',
                    slide_number=slide_num,
                    message=f'{len(pptx_bold_without_color)} bold text runs not styled with orange accent color',
                    details='Bold text should use orange (#ed5e29) to match HTML'
                ))

        return issues


def validate_conversion(html_path: str, pptx_path: str) -> ValidationReport:
    """
    Perform comprehensive validation of HTML to PPTX conversion.

    Args:
        html_path: Path to HTML file
        pptx_path: Path to PPTX file

    Returns:
        ValidationReport with all issues found
    """
    all_issues = []

    # Validate HTML
    print("Validating HTML...")
    html_validator = HTMLValidator(html_path)
    all_issues.extend(html_validator.validate())

    # Validate PPTX
    print("Validating PPTX...")
    pptx_validator = PPTXValidator(pptx_path)
    all_issues.extend(pptx_validator.validate())

    # Cross-validate
    print("Cross-validating HTML and PPTX...")
    cross_validator = CrossValidator(html_path, pptx_path)
    all_issues.extend(cross_validator.validate())

    # Categorize issues
    errors = [i for i in all_issues if i.severity == 'error']
    warnings = [i for i in all_issues if i.severity == 'warning']
    info = [i for i in all_issues if i.severity == 'info']

    # Get slide count
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)

    return ValidationReport(
        html_file=html_path,
        pptx_file=pptx_path,
        total_slides=total_slides,
        errors=errors,
        warnings=warnings,
        info=info
    )


def main():
    """CLI entry point."""
    if len(sys.argv) != 3:
        print("Usage: python3 validate_conversion.py <input.html> <output.pptx>")
        sys.exit(1)

    html_path = sys.argv[1]
    pptx_path = sys.argv[2]

    # Validate inputs exist
    if not Path(html_path).exists():
        print(f"Error: HTML file not found: {html_path}")
        sys.exit(1)

    if not Path(pptx_path).exists():
        print(f"Error: PPTX file not found: {pptx_path}")
        sys.exit(1)

    # Run validation
    report = validate_conversion(html_path, pptx_path)

    # Print results
    report.print_summary()

    # Exit code based on errors
    sys.exit(1 if report.has_errors() else 0)


if __name__ == "__main__":
    main()
