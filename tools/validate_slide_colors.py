#!/usr/bin/env python3
"""
Validate slide colors in generated PPTX files.
"""

import sys
from pptx import Presentation

def check_slide_colors(pptx_path, slide_num=3):
    """Check colors in specified slide to verify bold text is orange."""
    prs = Presentation(pptx_path)

    if len(prs.slides) < slide_num:
        print(f"ERROR: PPTX only has {len(prs.slides)} slides, need at least {slide_num}")
        return False

    slide = prs.slides[slide_num - 1]  # Convert to 0-indexed

    print(f"\n=== Slide {slide_num} Analysis ===")
    print(f"Total shapes: {len(slide.shapes)}")

    found_lists = False
    orange_count = 0
    gray_count = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            # Check if this is a bulleted paragraph
            if len(paragraph.runs) > 0:
                for run in paragraph.runs:
                    text = run.text.strip()
                    if not text:
                        continue

                    # Check if run has color
                    try:
                        color_type = run.font.color.type
                        if color_type == 1:  # RGB color
                            rgb = run.font.color.rgb
                            hex_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                            is_bold = run.font.bold

                            # Check for orange (#ed5e29) or gray (#475569)
                            if hex_color.lower() in ['#ed5e29', '#ed5d29']:  # Orange (allow slight variation)
                                orange_count += 1
                                print(f"✓ ORANGE: '{text[:30]}...' bold={is_bold} color={hex_color}")
                            elif hex_color.lower() == '#475569':  # Gray
                                gray_count += 1
                                if is_bold and 'Footer' in text or 'Decorative' in text or 'Icon' in text or 'Image' in text:
                                    print(f"✗ PROBLEM: '{text[:30]}...' is BOLD but GRAY: {hex_color}")
                            else:
                                print(f"  Other: '{text[:30]}...' bold={is_bold} color={hex_color}")
                    except Exception as e:
                        pass

    print(f"\n=== Summary ===")
    print(f"Orange text runs: {orange_count}")
    print(f"Gray text runs: {gray_count}")

    if orange_count == 0:
        print("\n❌ FAIL: No orange text found! Bold headings should be orange.")
        return False
    else:
        print(f"\n✓ PASS: Found {orange_count} orange text runs")
        return True

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python validate_slide_colors.py <path-to-pptx> [slide-number]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_num = int(sys.argv[2]) if len(sys.argv) > 2 else 3
    success = check_slide_colors(pptx_path, slide_num)
    sys.exit(0 if success else 1)
