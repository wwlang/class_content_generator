#!/usr/bin/env python3
"""Diagnostic tool to analyze overlapping shapes on specific slides."""

from pptx import Presentation
from pptx.util import Inches
import sys


def analyze_slide_shapes(pptx_path, slide_num):
    """Analyze shapes on a specific slide to identify overlaps."""
    prs = Presentation(pptx_path)

    if slide_num < 1 or slide_num > len(prs.slides):
        print(f"Error: Slide {slide_num} out of range (1-{len(prs.slides)})")
        return

    slide = prs.slides[slide_num - 1]  # 0-indexed

    print(f"\n{'='*70}")
    print(f"SLIDE {slide_num} SHAPE ANALYSIS")
    print(f"{'='*70}\n")

    # Collect all shapes with positions
    shapes_data = []
    for idx, shape in enumerate(slide.shapes):
        if hasattr(shape, 'left') and hasattr(shape, 'top'):
            shape_info = {
                'idx': idx,
                'shape': shape,
                'left': shape.left / 914400,  # EMUs to inches
                'top': shape.top / 914400,
                'width': shape.width / 914400,
                'height': shape.height / 914400,
                'right': (shape.left + shape.width) / 914400,
                'bottom': (shape.top + shape.height) / 914400,
                'has_text': hasattr(shape, 'text') and bool(shape.text.strip()),
                'text_preview': ''
            }

            if shape_info['has_text']:
                text = shape.text.strip()
                shape_info['text_preview'] = text[:50] + '...' if len(text) > 50 else text

            shapes_data.append(shape_info)

    # Print all shapes
    print(f"Total shapes: {len(shapes_data)}\n")
    for s in shapes_data:
        print(f"Shape {s['idx']}:")
        print(f"  Position: ({s['left']:.2f}\", {s['top']:.2f}\")")
        print(f"  Size: {s['width']:.2f}\" x {s['height']:.2f}\"")
        print(f"  Bounds: L={s['left']:.2f}\" T={s['top']:.2f}\" R={s['right']:.2f}\" B={s['bottom']:.2f}\"")
        if s['has_text']:
            print(f"  Text: \"{s['text_preview']}\"")
        print()

    # Check for overlaps
    print(f"\n{'='*70}")
    print("OVERLAP ANALYSIS")
    print(f"{'='*70}\n")

    overlaps_found = []
    for i, s1 in enumerate(shapes_data):
        for s2 in shapes_data[i+1:]:
            if rectangles_overlap(s1, s2):
                # Only report if both have text
                if s1['has_text'] and s2['has_text']:
                    overlaps_found.append((s1, s2))

    if overlaps_found:
        print(f"Found {len(overlaps_found)} overlapping text shape pairs:\n")
        for s1, s2 in overlaps_found:
            print(f"Overlap between Shape {s1['idx']} and Shape {s2['idx']}:")
            print(f"  Shape {s1['idx']}: ({s1['left']:.2f}\", {s1['top']:.2f}\") - ({s1['right']:.2f}\", {s1['bottom']:.2f}\")")
            print(f"    Text: \"{s1['text_preview']}\"")
            print(f"  Shape {s2['idx']}: ({s2['left']:.2f}\", {s2['top']:.2f}\") - ({s2['right']:.2f}\", {s2['bottom']:.2f}\")")
            print(f"    Text: \"{s2['text_preview']}\"")

            # Calculate overlap region
            overlap_left = max(s1['left'], s2['left'])
            overlap_right = min(s1['right'], s2['right'])
            overlap_top = max(s1['top'], s2['top'])
            overlap_bottom = min(s1['bottom'], s2['bottom'])
            overlap_width = overlap_right - overlap_left
            overlap_height = overlap_bottom - overlap_top

            print(f"  Overlap region: {overlap_width:.3f}\" wide x {overlap_height:.3f}\" tall")
            print()
    else:
        print("No overlapping text shapes found!")

    print(f"{'='*70}\n")


def rectangles_overlap(rect1, rect2):
    """Check if two rectangles overlap."""
    return not (rect1['right'] <= rect2['left'] or
               rect1['left'] >= rect2['right'] or
               rect1['bottom'] <= rect2['top'] or
               rect1['top'] >= rect2['bottom'])


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: python diagnose_overlaps.py <pptx_file> <slide_number>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_num = int(sys.argv[2])

    analyze_slide_shapes(pptx_path, slide_num)
