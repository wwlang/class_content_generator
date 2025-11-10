#!/usr/bin/env python3
"""
Analyze PowerPoint design elements from a reference PPTX file.
Extracts colors, fonts, layouts, and styling patterns.
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from collections import Counter
import sys

def analyze_pptx(file_path):
    """Analyze design elements from a PowerPoint file."""
    prs = Presentation(file_path)

    print("=" * 80)
    print(f"ANALYZING: {file_path}")
    print("=" * 80)

    # Slide dimensions
    print(f"\n📐 SLIDE DIMENSIONS:")
    print(f"   Width:  {prs.slide_width / 914400:.2f} inches ({prs.slide_width} EMU)")
    print(f"   Height: {prs.slide_height / 914400:.2f} inches ({prs.slide_height} EMU)")
    print(f"   Aspect: {prs.slide_width / prs.slide_height:.2f}:1")

    # Analyze first few slides in detail
    print(f"\n📊 TOTAL SLIDES: {len(prs.slides)}")

    # Collect design elements
    fonts_used = []
    colors_used = []
    font_sizes = []

    for slide_idx, slide in enumerate(list(prs.slides)[:5], 1):  # Analyze first 5 slides
        print(f"\n{'=' * 80}")
        print(f"SLIDE {slide_idx}")
        print(f"{'=' * 80}")

        # Check for background
        try:
            if slide.background.fill.type:
                print(f"Background fill type: {slide.background.fill.type}")
                try:
                    rgb = slide.background.fill.fore_color.rgb
                    print(f"Background color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]}) = #{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                    colors_used.append(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                except:
                    pass
        except:
            pass

        print(f"\nShapes: {len(slide.shapes)}")

        for shape_idx, shape in enumerate(slide.shapes):
            print(f"\n  Shape {shape_idx + 1}: {shape.shape_type} ({shape.name})")

            # Position and size
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                print(f"    Position: ({shape.left / 914400:.2f}\", {shape.top / 914400:.2f}\")")
                print(f"    Size: {shape.width / 914400:.2f}\" × {shape.height / 914400:.2f}\"")

            # Fill color
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type:
                        print(f"    Fill type: {shape.fill.type}")
                        if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                            rgb = shape.fill.fore_color.rgb
                            hex_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                            print(f"    Fill color: {hex_color}")
                            colors_used.append(hex_color)
                except:
                    pass

            # Text content
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text.strip()
                if text:
                    print(f"    Text: {text[:60]}{'...' if len(text) > 60 else ''}")

                    # Analyze text formatting
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph.text.strip():
                            print(f"      Paragraph {para_idx + 1}:")
                            print(f"        Alignment: {paragraph.alignment}")
                            print(f"        Level: {paragraph.level}")

                            for run_idx, run in enumerate(paragraph.runs):
                                if run.text.strip():
                                    print(f"        Run {run_idx + 1}: '{run.text[:40]}{'...' if len(run.text) > 40 else ''}'")

                                    # Font
                                    if run.font.name:
                                        print(f"          Font: {run.font.name}")
                                        fonts_used.append(run.font.name)

                                    # Size
                                    if run.font.size:
                                        size_pt = run.font.size.pt
                                        print(f"          Size: {size_pt}pt")
                                        font_sizes.append(size_pt)

                                    # Color
                                    try:
                                        if hasattr(run.font.color, 'rgb'):
                                            rgb = run.font.color.rgb
                                            hex_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                            print(f"          Color: {hex_color}")
                                            colors_used.append(hex_color)
                                    except:
                                        pass

                                    # Bold/Italic
                                    if run.font.bold:
                                        print(f"          Bold: Yes")
                                    if run.font.italic:
                                        print(f"          Italic: Yes")

    # Summary
    print(f"\n{'=' * 80}")
    print("DESIGN SUMMARY")
    print(f"{'=' * 80}")

    if fonts_used:
        font_counts = Counter(fonts_used)
        print(f"\n🔤 FONTS USED:")
        for font, count in font_counts.most_common():
            print(f"   {font}: {count} uses")

    if colors_used:
        color_counts = Counter(colors_used)
        print(f"\n🎨 COLOR PALETTE:")
        for color, count in color_counts.most_common(15):
            print(f"   {color}: {count} uses")

    if font_sizes:
        size_counts = Counter(font_sizes)
        print(f"\n📏 FONT SIZES:")
        for size, count in sorted(size_counts.items(), reverse=True):
            print(f"   {size}pt: {count} uses")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python3 analyze_pptx_design.py <path_to_pptx>")
        sys.exit(1)

    analyze_pptx(sys.argv[1])
