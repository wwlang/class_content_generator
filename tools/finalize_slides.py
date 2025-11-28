#!/usr/bin/env python3
"""
Finalize Slides: Merge Gemini batches + Insert Speaker Notes

This script implements a hybrid workflow:
1. Gemini creates visual slides in batches (better visuals, images, infographics)
2. This script merges batches and inserts speaker notes from lecture-content.md

Usage:
    python tools/finalize_slides.py [course-code] [week-number]

Example:
    python tools/finalize_slides.py BCI2AU 1
"""

import os
import re
import sys
import glob
from pathlib import Path
from dataclasses import dataclass
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# =============================================================================
# DATA STRUCTURES
# =============================================================================

@dataclass
class SlideNotes:
    """Speaker notes for a single slide."""
    slide_number: int
    slide_title: str
    notes_content: str


@dataclass
class ValidationResult:
    """Result of a quality validation check."""
    check_name: str
    passed: bool
    message: str
    severity: str  # "error", "warning", "info"


@dataclass
class ProcessingReport:
    """Final report of the finalization process."""
    success: bool
    batches_merged: int
    total_slides: int
    notes_inserted: int
    notes_missing: int
    validations: list
    output_file: str
    warnings: list
    errors: list


# =============================================================================
# SPEAKER NOTES PARSER
# =============================================================================

def parse_speaker_notes(lecture_content_path: str) -> list[SlideNotes]:
    """
    Parse speaker notes from lecture-content.md.

    Supported formats:
    1. ## Slide N: Title ... ### Speaker Notes (current format)
    2. **SLIDE N: Title** ... ## Speaker Notes (legacy format)
    """
    with open(lecture_content_path, 'r', encoding='utf-8') as f:
        content = f.read()

    notes_list = []

    # Try current format first: ## Slide N: Title
    current_pattern = r'##\s+Slide\s+(\d+):\s*([^\n]+)'
    slides = re.split(r'(?=##\s+Slide\s+\d+:)', content)

    for slide_block in slides:
        if not slide_block.strip():
            continue

        # Extract slide number and title
        slide_match = re.search(current_pattern, slide_block)
        if not slide_match:
            continue

        slide_num = int(slide_match.group(1))
        slide_title = slide_match.group(2).strip()

        # Extract speaker notes section (### Speaker Notes for current format)
        notes_match = re.search(
            r'###\s*Speaker\s*Notes\s*\n(.*?)(?=\n---|\n##\s+Slide|\Z)',
            slide_block,
            re.DOTALL | re.IGNORECASE
        )

        if notes_match:
            notes_content = notes_match.group(1).strip()
            notes_content = clean_notes_content(notes_content)

            notes_list.append(SlideNotes(
                slide_number=slide_num,
                slide_title=slide_title,
                notes_content=notes_content
            ))

    # If no notes found with current format, try legacy format
    if not notes_list:
        legacy_pattern = r'\*\*SLIDE\s+(\d+):\s*([^*]+)\*\*'
        slides = re.split(r'(?=\*\*SLIDE\s+\d+:)', content)

        for slide_block in slides:
            if not slide_block.strip():
                continue

            slide_match = re.search(legacy_pattern, slide_block)
            if not slide_match:
                continue

            slide_num = int(slide_match.group(1))
            slide_title = slide_match.group(2).strip()

            notes_match = re.search(
                r'##\s*Speaker\s*Notes\s*\n(.*?)(?=\n---|\n\*\*SLIDE|\Z)',
                slide_block,
                re.DOTALL | re.IGNORECASE
            )

            if notes_match:
                notes_content = notes_match.group(1).strip()
                notes_content = clean_notes_content(notes_content)

                notes_list.append(SlideNotes(
                    slide_number=slide_num,
                    slide_title=slide_title,
                    notes_content=notes_content
                ))

    return notes_list


def clean_notes_content(notes: str) -> str:
    """Clean and format speaker notes for PowerPoint."""
    # Remove excessive whitespace
    notes = re.sub(r'\n{3,}', '\n\n', notes)

    # Convert markdown bold to plain text (PowerPoint notes don't support markdown)
    notes = re.sub(r'\*\*([^*]+)\*\*', r'\1', notes)

    # Convert markdown italic
    notes = re.sub(r'\*([^*]+)\*', r'\1', notes)

    # Clean up bullet points (keep them but normalize)
    notes = re.sub(r'^[-*]\s+', '- ', notes, flags=re.MULTILINE)

    return notes.strip()


# =============================================================================
# PPTX OPERATIONS
# =============================================================================

def find_batch_files(week_folder: str, week_number: int) -> list[str]:
    """
    Find all batch PPTX files for a week, supporting multiple naming patterns.

    Supported patterns:
    1. Preferred: week-01-batch-1.pptx, week-01-batch-2.pptx
    2. Gemini default: "Topic Name.pptx" (batch 1), "Topic Name - Batch 2.pptx"
    3. Gemini variation: "Topic Name (Week 1).pptx", "Topic Name (Week 1) - Batch 2.pptx"
    """
    # Try preferred naming pattern first
    preferred_pattern = os.path.join(week_folder, f"week-{week_number:02d}-batch-*.pptx")
    preferred_files = glob.glob(preferred_pattern)

    if preferred_files:
        # Sort by batch number
        def get_batch_num(filepath):
            match = re.search(r'batch-(\d+)', filepath)
            return int(match.group(1)) if match else 0
        return sorted(preferred_files, key=get_batch_num)

    # Try Gemini naming patterns
    all_pptx = glob.glob(os.path.join(week_folder, "*.pptx"))

    # Exclude final output file if it exists
    all_pptx = [f for f in all_pptx if not f.endswith("slides.pptx")]

    if not all_pptx:
        return []

    # Categorize files
    batch_files = []
    for f in all_pptx:
        basename = os.path.basename(f)

        # Check for "Batch N" pattern (Gemini uses this for batches 2+)
        batch_match = re.search(r'[Bb]atch\s*(\d+)', basename)
        if batch_match:
            batch_num = int(batch_match.group(1))
        else:
            # No batch number = first batch
            batch_num = 1

        batch_files.append((batch_num, f))

    # Sort by batch number and return just the file paths
    batch_files.sort(key=lambda x: x[0])
    return [f for _, f in batch_files]


def merge_presentations(batch_files: list[str], output_path: str) -> Presentation:
    """
    Merge multiple PPTX files into one presentation.

    Uses a ZIP-level merge approach to preserve all media (images, etc.)
    """
    if not batch_files:
        raise ValueError("No batch files to merge")

    import zipfile
    import tempfile
    import shutil
    from lxml import etree

    # Create a working copy of the first presentation
    with tempfile.TemporaryDirectory() as temp_dir:
        base_dir = os.path.join(temp_dir, "base")

        # Extract first presentation as base
        with zipfile.ZipFile(batch_files[0], 'r') as zf:
            zf.extractall(base_dir)

        # Track the highest slide number and media file
        slide_count = count_slides_in_extracted(base_dir)
        media_count = count_media_in_extracted(base_dir)

        # Process each additional batch
        for batch_file in batch_files[1:]:
            batch_temp = os.path.join(temp_dir, f"batch_{os.path.basename(batch_file)}")

            with zipfile.ZipFile(batch_file, 'r') as zf:
                zf.extractall(batch_temp)

            # Merge this batch into base
            slide_count, media_count = merge_extracted_presentations(
                base_dir, batch_temp, slide_count, media_count
            )

        # Repackage as PPTX
        temp_output = os.path.join(temp_dir, "merged.pptx")
        create_pptx_from_extracted(base_dir, temp_output)

        # Load with python-pptx for further processing
        return Presentation(temp_output)


def count_slides_in_extracted(pptx_dir: str) -> int:
    """Count slides in an extracted PPTX directory."""
    slides_dir = os.path.join(pptx_dir, "ppt", "slides")
    if not os.path.exists(slides_dir):
        return 0
    return len([f for f in os.listdir(slides_dir) if f.startswith("slide") and f.endswith(".xml")])


def count_media_in_extracted(pptx_dir: str) -> int:
    """Count media files in an extracted PPTX directory."""
    media_dir = os.path.join(pptx_dir, "ppt", "media")
    if not os.path.exists(media_dir):
        return 0
    return len(os.listdir(media_dir))


def merge_extracted_presentations(base_dir: str, batch_dir: str, slide_offset: int, media_offset: int) -> tuple[int, int]:
    """
    Merge slides and media from batch into base.

    Returns updated (slide_count, media_count).
    """
    from lxml import etree
    import shutil

    # Namespace definitions
    NSMAP = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    REL_NS = 'http://schemas.openxmlformats.org/package/2006/relationships'

    batch_slides_dir = os.path.join(batch_dir, "ppt", "slides")
    base_slides_dir = os.path.join(base_dir, "ppt", "slides")
    batch_media_dir = os.path.join(batch_dir, "ppt", "media")
    base_media_dir = os.path.join(base_dir, "ppt", "media")

    os.makedirs(base_media_dir, exist_ok=True)
    os.makedirs(os.path.join(base_slides_dir, "_rels"), exist_ok=True)

    # Build media mapping: old name -> new name
    media_mapping = {}
    if os.path.exists(batch_media_dir):
        for media_file in sorted(os.listdir(batch_media_dir)):
            media_offset += 1
            ext = os.path.splitext(media_file)[1]
            new_name = f"image{media_offset}{ext}"
            media_mapping[media_file] = new_name

            # Copy media file
            shutil.copy2(
                os.path.join(batch_media_dir, media_file),
                os.path.join(base_media_dir, new_name)
            )

    # Get batch slide files
    batch_slides = sorted([f for f in os.listdir(batch_slides_dir)
                          if f.startswith("slide") and f.endswith(".xml")])

    new_slide_ids = []

    for batch_slide in batch_slides:
        slide_offset += 1
        new_slide_name = f"slide{slide_offset}.xml"
        new_slide_ids.append(slide_offset)

        # Copy and update slide XML
        batch_slide_path = os.path.join(batch_slides_dir, batch_slide)
        new_slide_path = os.path.join(base_slides_dir, new_slide_name)
        shutil.copy2(batch_slide_path, new_slide_path)

        # Copy and update slide relationships
        batch_rels_path = os.path.join(batch_slides_dir, "_rels", f"{batch_slide}.rels")
        new_rels_path = os.path.join(base_slides_dir, "_rels", f"{new_slide_name}.rels")

        if os.path.exists(batch_rels_path):
            # Update media references in relationships
            tree = etree.parse(batch_rels_path)
            root = tree.getroot()

            for rel in root.findall(f'{{{REL_NS}}}Relationship'):
                target = rel.get('Target')
                if target and '../media/' in target:
                    old_media = os.path.basename(target)
                    if old_media in media_mapping:
                        rel.set('Target', f'../media/{media_mapping[old_media]}')

            tree.write(new_rels_path, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Update presentation.xml to include new slides
    pres_path = os.path.join(base_dir, "ppt", "presentation.xml")
    pres_rels_path = os.path.join(base_dir, "ppt", "_rels", "presentation.xml.rels")
    content_types_path = os.path.join(base_dir, "[Content_Types].xml")

    # Update presentation.xml.rels
    if os.path.exists(pres_rels_path):
        tree = etree.parse(pres_rels_path)
        root = tree.getroot()

        # Find highest rId
        max_rid = 0
        for rel in root.findall(f'{{{REL_NS}}}Relationship'):
            rid = rel.get('Id', '')
            if rid.startswith('rId'):
                try:
                    max_rid = max(max_rid, int(rid[3:]))
                except ValueError:
                    pass

        # Add new slide relationships
        slide_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
        new_rids = []
        for slide_id in new_slide_ids:
            max_rid += 1
            rid = f"rId{max_rid}"
            new_rids.append((rid, slide_id))

            new_rel = etree.SubElement(root, f'{{{REL_NS}}}Relationship')
            new_rel.set('Id', rid)
            new_rel.set('Type', slide_type)
            new_rel.set('Target', f'slides/slide{slide_id}.xml')

        tree.write(pres_rels_path, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Update presentation.xml sldIdLst
    if os.path.exists(pres_path):
        tree = etree.parse(pres_path)
        root = tree.getroot()

        # Find sldIdLst
        sld_id_lst = root.find(f'.//{{{NSMAP["p"]}}}sldIdLst')
        if sld_id_lst is not None:
            # Find max slide id
            max_sld_id = 256  # PowerPoint starts slide IDs at 256
            for sld_id in sld_id_lst.findall(f'{{{NSMAP["p"]}}}sldId'):
                try:
                    max_sld_id = max(max_sld_id, int(sld_id.get('id', 256)))
                except ValueError:
                    pass

            # Add new slide references
            for rid, slide_id in new_rids:
                max_sld_id += 1
                new_sld = etree.SubElement(sld_id_lst, f'{{{NSMAP["p"]}}}sldId')
                new_sld.set('id', str(max_sld_id))
                new_sld.set(f'{{{NSMAP["r"]}}}id', rid)

        tree.write(pres_path, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Update [Content_Types].xml
    if os.path.exists(content_types_path):
        CT_NS = 'http://schemas.openxmlformats.org/package/2006/content-types'
        tree = etree.parse(content_types_path)
        root = tree.getroot()

        slide_type = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
        for slide_id in new_slide_ids:
            override = etree.SubElement(root, f'{{{CT_NS}}}Override')
            override.set('PartName', f'/ppt/slides/slide{slide_id}.xml')
            override.set('ContentType', slide_type)

        tree.write(content_types_path, xml_declaration=True, encoding='UTF-8', standalone=True)

    return slide_offset, media_offset


def create_pptx_from_extracted(extracted_dir: str, output_path: str):
    """Create a PPTX file from an extracted directory."""
    import zipfile

    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(extracted_dir):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, extracted_dir)
                zf.write(file_path, arcname)


def copy_slide_contents(src_slide, dest_slide):
    """Copy all shapes from source slide to destination slide (legacy - not used in ZIP merge)."""
    from copy import deepcopy

    # Remove default shapes from destination
    for shape in list(dest_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy shapes from source
    for shape in src_slide.shapes:
        el = shape._element
        new_el = deepcopy(el)
        dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')


def get_slide_text(slide) -> str:
    """Extract all text from a slide for pattern matching."""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text_parts.append(shape.text)
        elif hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                text_parts.append(para.text)
    return " ".join(text_parts).lower()


def remove_extra_slides(prs: Presentation, first_slide_title: str = None) -> dict:
    """
    Remove Gemini-generated extra slides:
    - "Image Sources" slides (Gemini adds these at end of each batch)
    - Duplicate title slides from batches 2+ (same title as first slide)

    Returns dict with removal statistics.
    """
    slides_to_remove = []
    stats = {
        "image_sources_removed": 0,
        "duplicate_titles_removed": 0,
        "total_removed": 0
    }

    # Get first slide title if not provided
    if first_slide_title is None and prs.slides:
        first_slide = prs.slides[0]
        if first_slide.shapes.title:
            first_slide_title = first_slide.shapes.title.text.lower().strip()

    # Scan all slides
    for i, slide in enumerate(prs.slides):
        slide_text = get_slide_text(slide)
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.lower().strip()

        # Check for "Image Sources" slides
        if "image source" in slide_text or "image sources" in slide_text:
            slides_to_remove.append(i)
            stats["image_sources_removed"] += 1
            continue

        # Check for duplicate title slides (but not the first one)
        if i > 0 and first_slide_title:
            if slide_title == first_slide_title:
                slides_to_remove.append(i)
                stats["duplicate_titles_removed"] += 1
                continue

    # Remove slides in reverse order (to preserve indices)
    for i in reversed(slides_to_remove):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    stats["total_removed"] = len(slides_to_remove)
    return stats


def parse_references_from_lecture(lecture_content_path: str) -> list[str]:
    """Extract references section from lecture-content.md."""
    with open(lecture_content_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Find References section
    refs_match = re.search(
        r'##\s*References\s*\n(.*?)(?=\n##\s|\Z)',
        content,
        re.DOTALL | re.IGNORECASE
    )

    if not refs_match:
        return []

    refs_text = refs_match.group(1).strip()

    # Extract individual references (each starts with bullet or numbered)
    references = []
    for line in refs_text.split('\n'):
        line = line.strip()
        # Remove markdown bullets
        line = re.sub(r'^[-*•]\s*', '', line)
        line = re.sub(r'^\d+\.\s*', '', line)
        if line and len(line) > 10:  # Ignore empty or very short lines
            references.append(line)

    return references


def has_references_slide(prs: Presentation) -> bool:
    """Check if the presentation already has a References slide."""
    if not prs.slides:
        return False

    # Check last few slides
    for slide in list(prs.slides)[-3:]:
        slide_text = get_slide_text(slide)
        if "references" in slide_text.lower() or "bibliography" in slide_text.lower():
            return True

    return False


def add_references_slide(prs: Presentation, references: list[str]) -> bool:
    """
    Add a References slide to the end of the presentation.

    Returns True if slide was added, False if references list is empty.
    """
    if not references:
        return False

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Typically blank

    # Add new slide
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    from pptx.util import Inches, Pt

    # Add title textbox
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_tf = title_box.text_frame
    title_tf.text = "References"
    title_para = title_tf.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Add references content
    refs_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9), Inches(6)
    )
    refs_tf = refs_box.text_frame
    refs_tf.word_wrap = True

    # Add each reference
    for i, ref in enumerate(references[:10]):  # Limit to 10 to fit slide
        if i == 0:
            refs_tf.text = ref
        else:
            p = refs_tf.add_paragraph()
            p.text = ref
            p.space_before = Pt(6)

        # Style the paragraph
        para = refs_tf.paragraphs[i]
        para.font.size = Pt(14)

    return True


def insert_speaker_notes(prs: Presentation, notes_list: list[SlideNotes]) -> tuple[int, int]:
    """
    Insert speaker notes into presentation slides.

    Returns: (notes_inserted, notes_missing)
    """
    notes_inserted = 0
    notes_missing = 0

    # Create a lookup by slide number
    notes_lookup = {n.slide_number: n for n in notes_list}

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        if slide_num in notes_lookup:
            notes = notes_lookup[slide_num]

            # Get or create notes slide
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_placeholder.text_frame

            # Clear existing notes and add new content
            notes_tf.text = notes.notes_content

            notes_inserted += 1
        else:
            notes_missing += 1

    return notes_inserted, notes_missing


# =============================================================================
# QUALITY VALIDATIONS
# =============================================================================

def run_validations(
    prs: Presentation,
    notes_list: list[SlideNotes],
    batch_files: list[str],
    removal_stats: dict = None,
    refs_added: bool = False
) -> list[ValidationResult]:
    """Run all quality validation checks."""
    results = []

    # 1. Slide count validation
    results.append(validate_slide_count(prs, notes_list))

    # 2. Speaker notes coverage
    results.append(validate_notes_coverage(prs, notes_list))

    # 3. Batch continuity
    results.append(validate_batch_continuity(batch_files))

    # 4. Extra slides removed (info only)
    if removal_stats and removal_stats.get("total_removed", 0) > 0:
        results.append(ValidationResult(
            check_name="Extra Slides Removed",
            passed=True,
            message=f"Removed {removal_stats['total_removed']} extra slides (Image Sources: {removal_stats.get('image_sources_removed', 0)}, Duplicate titles: {removal_stats.get('duplicate_titles_removed', 0)})",
            severity="info"
        ))

    # 5. References slide check
    results.append(validate_references_slide(prs, refs_added))

    # 6. Minimum font size check
    results.extend(validate_font_sizes(prs, min_size=14))

    # 7. Title presence check
    results.append(validate_slide_titles(prs))

    return results


def validate_references_slide(prs: Presentation, refs_added: bool) -> ValidationResult:
    """Check if References slide exists."""
    has_refs = has_references_slide(prs)

    if has_refs:
        msg = "References slide present"
        if refs_added:
            msg += " (auto-generated)"
        return ValidationResult(
            check_name="References Present",
            passed=True,
            message=msg,
            severity="info"
        )
    else:
        return ValidationResult(
            check_name="References Present",
            passed=False,
            message="No References slide found",
            severity="warning"
        )


def validate_slide_count(prs: Presentation, notes_list: list[SlideNotes]) -> ValidationResult:
    """Check if slide count matches expected notes count."""
    prs_count = len(prs.slides)
    notes_count = len(notes_list)

    if prs_count == notes_count:
        return ValidationResult(
            check_name="Slide Count Match",
            passed=True,
            message=f"Slide count ({prs_count}) matches lecture content",
            severity="info"
        )
    else:
        return ValidationResult(
            check_name="Slide Count Match",
            passed=False,
            message=f"Slide count mismatch: PPTX has {prs_count} slides, lecture content has {notes_count}",
            severity="warning"
        )


def validate_notes_coverage(prs: Presentation, notes_list: list[SlideNotes]) -> ValidationResult:
    """Check if all slides will have speaker notes."""
    prs_count = len(prs.slides)
    notes_count = len(notes_list)
    coverage = (notes_count / prs_count * 100) if prs_count > 0 else 0

    if coverage >= 100:
        return ValidationResult(
            check_name="Speaker Notes Coverage",
            passed=True,
            message=f"100% speaker notes coverage ({notes_count}/{prs_count} slides)",
            severity="info"
        )
    elif coverage >= 80:
        return ValidationResult(
            check_name="Speaker Notes Coverage",
            passed=True,
            message=f"{coverage:.0f}% speaker notes coverage ({notes_count}/{prs_count} slides)",
            severity="warning"
        )
    else:
        return ValidationResult(
            check_name="Speaker Notes Coverage",
            passed=False,
            message=f"Low speaker notes coverage: {coverage:.0f}% ({notes_count}/{prs_count} slides)",
            severity="error"
        )


def validate_batch_continuity(batch_files: list[str]) -> ValidationResult:
    """Check that batch files are sequential (1, 2, 3, not 1, 3, 5)."""
    if not batch_files:
        return ValidationResult(
            check_name="Batch Continuity",
            passed=False,
            message="No batch files found",
            severity="error"
        )

    # Extract batch numbers using flexible pattern
    batch_nums = []
    for f in batch_files:
        basename = os.path.basename(f)
        # Check for "batch-N" or "Batch N" pattern
        match = re.search(r'[Bb]atch[- ]?(\d+)', basename)
        if match:
            batch_nums.append(int(match.group(1)))
        else:
            # No batch number = first batch (Gemini naming)
            batch_nums.append(1)

    # For files found, we expect sequential numbering
    # Note: With Gemini naming, we might have [1, 2, 3] from files like:
    # - "Topic.pptx" (no batch = 1)
    # - "Topic - Batch 2.pptx"
    # - "Topic - Batch 3.pptx"
    expected = list(range(1, len(batch_nums) + 1))

    if sorted(batch_nums) == expected:
        return ValidationResult(
            check_name="Batch Continuity",
            passed=True,
            message=f"All {len(batch_nums)} batches present and sequential",
            severity="info"
        )
    else:
        return ValidationResult(
            check_name="Batch Continuity",
            passed=False,
            message=f"Batch sequence issue: found {sorted(batch_nums)}, expected {expected}",
            severity="error"
        )


def validate_font_sizes(prs: Presentation, min_size: int = 14) -> list[ValidationResult]:
    """Check for fonts below minimum readable size."""
    results = []
    small_font_slides = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        has_small_font = False

        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size.pt < min_size:
                            has_small_font = True
                            break

        if has_small_font:
            small_font_slides.append(slide_num)

    if small_font_slides:
        results.append(ValidationResult(
            check_name="Font Size Check",
            passed=False,
            message=f"Slides with fonts < {min_size}pt: {small_font_slides[:5]}{'...' if len(small_font_slides) > 5 else ''}",
            severity="warning"
        ))
    else:
        results.append(ValidationResult(
            check_name="Font Size Check",
            passed=True,
            message=f"All fonts >= {min_size}pt (readable at distance)",
            severity="info"
        ))

    return results


def validate_slide_titles(prs: Presentation) -> ValidationResult:
    """Check that most slides have identifiable titles."""
    slides_without_title = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        has_title = False

        # Check for title placeholder
        if slide.shapes.title:
            has_title = True
        else:
            # Check for any large text that might be a title
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size.pt >= 24:
                                has_title = True
                                break

        if not has_title:
            slides_without_title.append(slide_num)

    if len(slides_without_title) > len(prs.slides) * 0.2:
        return ValidationResult(
            check_name="Slide Titles",
            passed=False,
            message=f"Many slides missing titles: {slides_without_title[:5]}{'...' if len(slides_without_title) > 5 else ''}",
            severity="warning"
        )
    else:
        return ValidationResult(
            check_name="Slide Titles",
            passed=True,
            message="Most slides have identifiable titles",
            severity="info"
        )


# =============================================================================
# MAIN PROCESSING
# =============================================================================

def finalize_slides(course_code: str, week_number: int, base_path: str = None) -> ProcessingReport:
    """
    Main function to finalize slides for a week.

    1. Find batch files
    2. Merge into single presentation
    3. Parse speaker notes from lecture content
    4. Insert speaker notes
    5. Run validations
    6. Save final presentation
    """
    if base_path is None:
        base_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    # Construct paths
    course_folder = os.path.join(base_path, "courses", f"{course_code}*")
    course_matches = glob.glob(course_folder)

    if not course_matches:
        return ProcessingReport(
            success=False,
            batches_merged=0,
            total_slides=0,
            notes_inserted=0,
            notes_missing=0,
            validations=[],
            output_file="",
            warnings=[],
            errors=[f"Course folder not found for: {course_code}"]
        )

    course_path = course_matches[0]
    week_folder = os.path.join(course_path, "weeks", f"week-{week_number:02d}")
    lecture_content_path = os.path.join(week_folder, "lecture-content.md")
    output_path = os.path.join(week_folder, "slides.pptx")

    errors = []
    warnings = []

    # 1. Find batch files
    print(f"\n1. Finding batch files in: {week_folder}")
    batch_files = find_batch_files(week_folder, week_number)

    if not batch_files:
        return ProcessingReport(
            success=False,
            batches_merged=0,
            total_slides=0,
            notes_inserted=0,
            notes_missing=0,
            validations=[],
            output_file="",
            warnings=[],
            errors=[f"No batch files found matching: week-{week_number:02d}-batch-*.pptx"]
        )

    print(f"   Found {len(batch_files)} batch files:")
    for f in batch_files:
        print(f"   - {os.path.basename(f)}")

    # 2. Parse speaker notes
    print(f"\n2. Parsing speaker notes from: {os.path.basename(lecture_content_path)}")

    if not os.path.exists(lecture_content_path):
        return ProcessingReport(
            success=False,
            batches_merged=len(batch_files),
            total_slides=0,
            notes_inserted=0,
            notes_missing=0,
            validations=[],
            output_file="",
            warnings=[],
            errors=[f"Lecture content not found: {lecture_content_path}"]
        )

    notes_list = parse_speaker_notes(lecture_content_path)
    print(f"   Parsed {len(notes_list)} speaker notes sections")

    # 3. Merge presentations
    print(f"\n3. Merging {len(batch_files)} batch files...")

    try:
        merged_prs = merge_presentations(batch_files, output_path)
        total_slides_before = len(merged_prs.slides)
        print(f"   Merged presentation has {total_slides_before} slides")
    except Exception as e:
        return ProcessingReport(
            success=False,
            batches_merged=0,
            total_slides=0,
            notes_inserted=0,
            notes_missing=0,
            validations=[],
            output_file="",
            warnings=[],
            errors=[f"Failed to merge presentations: {str(e)}"]
        )

    # 3.5. Remove extra slides (Image Sources, duplicate titles)
    print(f"\n3.5. Removing extra slides (Image Sources, duplicate titles)...")
    removal_stats = remove_extra_slides(merged_prs)
    if removal_stats["total_removed"] > 0:
        print(f"   Removed {removal_stats['total_removed']} slides:")
        if removal_stats["image_sources_removed"]:
            print(f"      - Image Sources: {removal_stats['image_sources_removed']}")
        if removal_stats["duplicate_titles_removed"]:
            print(f"      - Duplicate titles: {removal_stats['duplicate_titles_removed']}")
    else:
        print(f"   No extra slides to remove")

    total_slides = len(merged_prs.slides)

    # 4. Insert speaker notes
    print(f"\n4. Inserting speaker notes...")
    notes_inserted, notes_missing = insert_speaker_notes(merged_prs, notes_list)
    print(f"   Inserted: {notes_inserted}, Missing: {notes_missing}")

    # 4.5. Add References slide if missing
    print(f"\n4.5. Checking for References slide...")
    refs_added = False
    if not has_references_slide(merged_prs):
        references = parse_references_from_lecture(lecture_content_path)
        if references:
            refs_added = add_references_slide(merged_prs, references)
            if refs_added:
                print(f"   Added References slide with {len(references)} citations")
                total_slides = len(merged_prs.slides)
        else:
            print(f"   No references found in lecture content")
            warnings.append("No References slide and no references found in lecture-content.md")
    else:
        print(f"   References slide already present")

    # 5. Run validations
    print(f"\n5. Running quality validations...")
    validations = run_validations(merged_prs, notes_list, batch_files, removal_stats, refs_added)

    passed = sum(1 for v in validations if v.passed)
    failed = len(validations) - passed
    print(f"   {passed} passed, {failed} issues found")

    # Collect warnings and errors from validations
    for v in validations:
        if not v.passed:
            if v.severity == "error":
                errors.append(v.message)
            elif v.severity == "warning":
                warnings.append(v.message)

    # 6. Save final presentation
    print(f"\n6. Saving final presentation...")
    merged_prs.save(output_path)
    print(f"   Saved to: {output_path}")

    # Determine overall success
    has_errors = any(v.severity == "error" and not v.passed for v in validations)

    return ProcessingReport(
        success=not has_errors,
        batches_merged=len(batch_files),
        total_slides=total_slides,
        notes_inserted=notes_inserted,
        notes_missing=notes_missing,
        validations=validations,
        output_file=output_path,
        warnings=warnings,
        errors=errors
    )


def print_report(report: ProcessingReport):
    """Print a formatted report of the finalization process."""
    print("\n" + "=" * 60)
    print("FINALIZATION REPORT")
    print("=" * 60)

    status = "SUCCESS" if report.success else "FAILED"
    print(f"\nStatus: {status}")
    print(f"Output: {report.output_file}")

    print(f"\nProcessing Summary:")
    print(f"  - Batches merged: {report.batches_merged}")
    print(f"  - Total slides: {report.total_slides}")
    print(f"  - Speaker notes inserted: {report.notes_inserted}")
    print(f"  - Speaker notes missing: {report.notes_missing}")

    print(f"\nValidation Results:")
    for v in report.validations:
        icon = "PASS" if v.passed else "FAIL"
        severity_icon = {"error": "!!!", "warning": "!", "info": " "}[v.severity]
        print(f"  [{icon}] {severity_icon} {v.check_name}: {v.message}")

    if report.warnings:
        print(f"\nWarnings ({len(report.warnings)}):")
        for w in report.warnings:
            print(f"  - {w}")

    if report.errors:
        print(f"\nErrors ({len(report.errors)}):")
        for e in report.errors:
            print(f"  - {e}")

    print("\n" + "=" * 60)


# =============================================================================
# CLI ENTRY POINT
# =============================================================================

def main():
    """Command-line entry point."""
    if len(sys.argv) < 3:
        print("Usage: python tools/finalize_slides.py [course-code] [week-number]")
        print("Example: python tools/finalize_slides.py BCI2AU 1")
        sys.exit(1)

    course_code = sys.argv[1]
    week_number = int(sys.argv[2])

    print(f"Finalizing slides for {course_code} Week {week_number}")
    print("-" * 40)

    report = finalize_slides(course_code, week_number)
    print_report(report)

    sys.exit(0 if report.success else 1)


if __name__ == "__main__":
    main()
